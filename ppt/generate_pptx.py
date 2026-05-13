# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Theme Colors ──
GREEN_900 = RGBColor(0x0D, 0x33, 0x20)
GREEN_800 = RGBColor(0x14, 0x52, 0x2F)
GREEN_700 = RGBColor(0x1A, 0x6B, 0x3C)
GREEN_600 = RGBColor(0x2E, 0x8B, 0x57)
GREEN_500 = RGBColor(0x3C, 0xB3, 0x71)
GREEN_400 = RGBColor(0x66, 0xCD, 0xAA)
GREEN_100 = RGBColor(0xD4, 0xF0, 0xE4)
GREEN_50  = RGBColor(0xEA, 0xFA, 0xF3)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x00, 0x00, 0x00)
GRAY_800  = RGBColor(0x2D, 0x2D, 0x2D)
GRAY_600  = RGBColor(0x55, 0x55, 0x55)
GRAY_400  = RGBColor(0x9C, 0xA3, 0xAF)
GRAY_200  = RGBColor(0xE2, 0xE6, 0xEA)
GRAY_100  = RGBColor(0xF0, 0xF2, 0xF4)
GRAY_50   = RGBColor(0xF8, 0xF9, 0xFA)
GOLD      = RGBColor(0xC8, 0x95, 0x3A)
RED       = RGBColor(0xEF, 0x53, 0x50)
BLUE      = RGBColor(0x1E, 0x40, 0xAF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height
FONT_TITLE = 'Microsoft YaHei'
FONT_BODY  = 'Microsoft YaHei'
FONT_MONO  = 'Consolas'

# ── Helper Functions ──
def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.fill.solid()
        if line_width:
            shape.line.width = line_width
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape

def add_circle(slide, left, top, size, fill_color=GREEN_600):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape

def add_line(slide, left, top, width, height, color=GREEN_500, width_pt=4):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=GRAY_800, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name=FONT_BODY, line_spacing=1.2, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(2)
    if line_spacing:
        p.line_spacing = Pt(int(font_size * line_spacing))
    return txBox

def add_multiline_text(slide, left, top, width, height, lines, font_name=FONT_BODY):
    """lines: list of (text, font_size, color, bold)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, fs, clr, bld) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(fs)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = font_name
        p.space_after = Pt(4)
    return txBox

def add_page_indicator(slide, num, total):
    add_text_box(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
                 f'{num}', font_size=10, color=GRAY_400, alignment=PP_ALIGN.RIGHT,
                 font_name=FONT_MONO)

def add_header(slide, title, subtitle=''):
    # dot
    add_circle(slide, Inches(0.6), Inches(0.5), Inches(0.15), GREEN_600)
    # title
    add_text_box(slide, Inches(0.9), Inches(0.4), Inches(8), Inches(0.6),
                 title, font_size=26, color=GREEN_800, bold=True)
    # subtitle
    if subtitle:
        add_text_box(slide, Inches(11.5), Inches(0.4), Inches(1.5), Inches(0.4),
                     subtitle, font_size=11, color=GRAY_400, alignment=PP_ALIGN.RIGHT)
    # separator line
    add_line(slide, Inches(0.6), Inches(1.05), Inches(12.1), Pt(1), color=GRAY_200, width_pt=1)

def add_card(slide, left, top, width, height, icon, title, desc,
             tag_text='', tag_color=GREEN_100, tag_text_color=GREEN_700):
    bg = add_rounded_rect(slide, left, top, width, height, WHITE)
    bg.line.color.rgb = GRAY_100
    bg.line.fill.solid()
    bg.line.width = Pt(1)
    # icon bg
    add_rounded_rect(slide, left + Inches(0.15), top + Inches(0.15), Inches(0.5), Inches(0.5), GREEN_100)
    add_text_box(slide, left + Inches(0.18), top + Inches(0.18), Inches(0.44), Inches(0.44),
                 icon, font_size=16, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # title
    add_text_box(slide, left + Inches(0.15), top + Inches(0.75), width - Inches(0.3), Inches(0.35),
                 title, font_size=12, color=GRAY_800, bold=True)
    # desc
    if isinstance(desc, list):
        for j, dline in enumerate(desc):
            add_text_box(slide, left + Inches(0.15), top + Inches(1.05) + Inches(j * 0.2),
                         width - Inches(0.3), Inches(0.2),
                         dline, font_size=9, color=GRAY_600)
    else:
        add_text_box(slide, left + Inches(0.15), top + Inches(1.05), width - Inches(0.3),
                     height - Inches(1.2), desc, font_size=9, color=GRAY_600, line_spacing=1.3)

# ── Slide Contents ──

# ===== SLIDE 1: COVER =====
s = add_blank_slide()
add_bg(s, GREEN_600)
# gradient overlay - darken top area
add_rect(s, 0, 0, W, Inches(3.5), GREEN_800)
# badge
badge = add_text_box(s, 0, Inches(1.5), W, Inches(0.6),
                     'P R O J E C T   P A N O R A M A',
                     font_size=13, color=RGBColor(0xFF, 0xFF, 0xFF), bold=False,
                     alignment=PP_ALIGN.CENTER, font_name=FONT_MONO)
# title
add_text_box(s, 0, Inches(2.6), W, Inches(1.0),
             '桂林银发旅游小程序',
             font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(s, 0, Inches(3.6), W, Inches(0.6),
             '专为银发群体打造的桂林智慧导览与文旅电商平台',
             font_size=16, color=RGBColor(0xFF,0xFF,0xFF), bold=False, alignment=PP_ALIGN.CENTER)
# tags
tags_data = [
    ('微信小程序', Inches(3.5), Inches(4.6)),
    ('Vue 3 管理后台', Inches(5.2), Inches(4.6)),
    ('Express + TypeScript', Inches(7.1), Inches(4.6)),
    ('SQL.js 数据库', Inches(9.1), Inches(4.6)),
    ('Railway 部署', Inches(10.9), Inches(4.6)),
]
for tag_text, l, t in tags_data:
    shape = add_rounded_rect(s, l, t, Inches(1.7), Inches(0.45), None)
    shape.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    shape.line.fill.solid()
    shape.line.width = Pt(0.5)
    add_text_box(s, l, t, Inches(1.7), Inches(0.45), tag_text,
                 font_size=11, color=WHITE, bold=False, alignment=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)

# ===== Helper to create standard slides =====
def make_slide(title, sub=''):
    s = add_blank_slide()
    add_header(s, title, sub)
    return s

def stat_block(slide, left, top, num_text, label):
    add_text_box(slide, left, top, Inches(1.4), Inches(0.6),
                 num_text, font_size=26, color=GREEN_700, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + Inches(0.55), Inches(1.4), Inches(0.3),
                 label, font_size=10, color=GRAY_600, alignment=PP_ALIGN.CENTER)

def highlight_box(slide, left, top, width, height, h4, ptext):
    bg = add_rounded_rect(slide, left, top, width, height, GREEN_50)
    bg.line.color.rgb = GREEN_400
    bg.line.fill.solid()
    bg.line.width = Pt(1)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), Inches(0.35),
                 h4, font_size=12, color=GREEN_700, bold=True)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.42), width - Inches(0.3),
                 height - Inches(0.52), ptext, font_size=9, color=GRAY_600, line_spacing=1.3)


# ===== SLIDE 2: TOC =====
s = make_slide('目录', 'Contents')
toc = [
    '1  项目概述与定位',      '2  系统架构全景图',
    '3  技术栈详解',          '4  页面路由与导航体系',
    '5  首页设计与功能模块',    '6  景点导览模块详解',
    '7  安全防护模块详解',     '8  个人中心模块详解',
    '9  辅助功能页面全景',     '10 电商功能模块',
    '11 银发友好设计体系',     '12 安全防护体系架构',
    '13 紧急求助SOS流程',     '14 后端架构与安全措施',
    '15 数据库设计（9张表）',  '16 API接口全景',
    '17 天气API多级Fallback', '18 组件化与工具模块',
    '19 Web管理后台（Vue 3）', '20 部署架构与运维策略',
    '21 全局状态与CSS变量',    '22 技术亮点与创新总结',
    '23 项目文件结构总览',     '24 致谢',
]
for i, txt in enumerate(toc):
    col = i % 2
    row = i // 2
    left = Inches(0.7) + col * Inches(6.1)
    top = Inches(1.4) + row * Inches(0.42)
    add_text_box(s, left, top, Inches(5.8), Inches(0.38),
                 txt, font_size=11, color=GRAY_800)

# ===== SLIDE 3: Overview =====
s = make_slide('项目概述与定位', 'Overview')
# left column
highlight_box(s, Inches(0.6), Inches(1.3), Inches(5.8), Inches(0.9),
              '项目名称',
              '桂林银发旅游小程序（Guilin Silver Hair Tourism）')
highlight_box(s, Inches(0.6), Inches(2.35), Inches(5.8), Inches(0.9),
              '目标用户',
              '银发群体（老年人）—— 到桂林旅游的长者及其家属')
highlight_box(s, Inches(0.6), Inches(3.4), Inches(5.8), Inches(0.9),
              '核心定位',
              '智慧导览 + 文旅电商 + 安全防护 + 健康关怀 四位一体')
highlight_box(s, Inches(0.6), Inches(4.45), Inches(5.8), Inches(0.7),
              '版本信息',
              'v1.0.0 | 开发完成 | 2026年5月')
# right - stats
stat_block(s, Inches(7.0), Inches(1.3), '18', '小程序页面')
stat_block(s, Inches(8.6), Inches(1.3), '9', '数据库表')
stat_block(s, Inches(7.0), Inches(2.1), '20+', 'API接口')
stat_block(s, Inches(8.6), Inches(2.1), '4', 'TabBar入口')
stat_block(s, Inches(7.0), Inches(2.9), '7', '大景点详情')
stat_block(s, Inches(8.6), Inches(2.9), '3', '服务端技术栈')
# pain points
add_text_box(s, Inches(7.0), Inches(3.9), Inches(6), Inches(0.4),
             '解决的痛点', font_size=13, color=GREEN_700, bold=True)
pains = [
    '🎯 老人看手机字太小 → 三档字体缩放',
    '🎯 老人不熟悉智能操作 → 全程语音播报导航',
    '🎯 出游安全隐患多 → 天气预警 + 路滑评估 + SOS求助',
    '🎯 景区信息不透明 → 老人友好评分 + 适老化提示',
]
for i, p in enumerate(pains):
    add_text_box(s, Inches(7.0), Inches(4.35) + i * Inches(0.35), Inches(6), Inches(0.3),
                 p, font_size=11, color=GRAY_800)


# ===== SLIDE 4: Architecture =====
s = make_slide('系统架构全景图', 'Architecture')
layers = [
    ('📱', '微信小程序', '原生框架 + 自定义TabBar\n18个页面 + 4个组件\nWechatSI语音插件\n高德地图SDK', '表现层', GREEN_600, GREEN_100),
    ('⚙️', 'Express 后端', 'TypeScript + Node.js\nSQL.js 嵌入式数据库\nHelmet + Rate Limit\nZod 参数验证', '服务层', GOLD, RGBColor(0xFE,0xF3,0xC7)),
    ('💻', 'Vue 3 Web后台', 'Vue 3 + TypeScript\nVite 构建 + Pinia 状态\nVue Router 路由\n10个管理页面', '管理层', BLUE, RGBColor(0xDB,0xEA,0xFE)),
]
for i, (icon, h5, items, tag, tc, tbg) in enumerate(layers):
    left = Inches(0.7) + i * Inches(4.1)
    bg = add_rounded_rect(s, left, Inches(1.4), Inches(3.7), Inches(3.6), WHITE)
    bg.line.color.rgb = GRAY_200
    bg.line.fill.solid()
    bg.line.width = Pt(1.5)
    # icon
    add_text_box(s, left, Inches(1.55), Inches(3.7), Inches(0.7),
                 icon, font_size=30, alignment=PP_ALIGN.CENTER)
    # h5
    add_text_box(s, left, Inches(2.3), Inches(3.7), Inches(0.4),
                 h5, font_size=15, color=GREEN_700, bold=True, alignment=PP_ALIGN.CENTER)
    # items
    add_text_box(s, left + Inches(0.3), Inches(2.75), Inches(3.1), Inches(1.6),
                 items, font_size=10, color=GRAY_600, alignment=PP_ALIGN.CENTER, line_spacing=1.5)
    # tag
    tag_bg = add_rounded_rect(s, left + Inches(1.0), Inches(4.4), Inches(1.7), Inches(0.35), tbg)
    add_text_box(s, left + Inches(1.0), Inches(4.4), Inches(1.7), Inches(0.35),
                 tag, font_size=9, color=tc, bold=True, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # arrow
    if i < 2:
        add_text_box(s, left + Inches(3.7), Inches(3.0), Inches(0.4), Inches(0.4),
                     '→', font_size=22, color=GREEN_400, alignment=PP_ALIGN.CENTER)
# external services
highlight_box(s, Inches(0.6), Inches(5.2), Inches(12.1), Inches(1.2),
              '🌐 外部服务依赖',
              '微信开放平台 — 登录授权（code换openid）  |  Open-Meteo — 免费天气API（主）  |  高德地图 — 天气Fallback + 地图SDK  |  有道词典 — TTS语音Fallback  |  Railway — 云部署平台')

# ===== SLIDE 5: Tech Stack =====
s = make_slide('技术栈详解', 'Tech Stack')
stacks = [
    ('📱', '微信小程序', '原生框架 WXML + WXSS + JS\n自定义TabBar，WechatSI语音插件\n高德地图SDK', '前端', GREEN_600, GREEN_100),
    ('🖥️', 'Vue 3 Web后台', 'Vue 3 + TypeScript + Vite\nPinia状态管理 + Vue Router路由\n10个管理页面', '管理端', BLUE, RGBColor(0xDB,0xEA,0xFE)),
    ('⚡', 'Express 后端', 'Node.js + TypeScript + Express\nSQL.js嵌入式数据库，Pino日志\nRailway部署', '服务端', GOLD, RGBColor(0xFE,0xF3,0xC7)),
    ('🗄️', 'SQL.js 数据库', 'SQLite编译为WebAssembly\n零配置嵌入式运行\n每小时自动备份保留24份', '存储', RGBColor(0x5B,0x21,0xB6), RGBColor(0xED,0xE9,0xFE)),
    ('🛡️', '安全中间件', 'Helmet安全头 + CORS\nRate Limit(100次/15分钟)\nZod参数验证 + 审计日志', '安全', RED, RGBColor(0xFE,0xE2,0xE2)),
]
for i, (icon, h4, desc, tag, tc, tbg) in enumerate(stacks):
    row = i // 3
    col = i % 3
    left = Inches(0.7) + col * Inches(4.1)
    top = Inches(1.4) + row * Inches(2.7)
    bg_c = add_rounded_rect(s, left, top, Inches(3.7), Inches(2.3), WHITE)
    bg_c.line.color.rgb = GRAY_100
    bg_c.line.fill.solid()
    bg_c.line.width = Pt(1)
    add_rounded_rect(s, left + Inches(0.15), top + Inches(0.15), Inches(0.5), Inches(0.5), GREEN_100)
    add_text_box(s, left + Inches(0.18), top + Inches(0.17), Inches(0.44), Inches(0.44),
                 icon, font_size=18, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(s, left + Inches(0.8), top + Inches(0.15), Inches(2.6), Inches(0.4),
                 h4, font_size=13, color=GRAY_800, bold=True)
    add_text_box(s, left + Inches(0.15), top + Inches(0.8), Inches(3.4), Inches(1.0),
                 desc, font_size=9, color=GRAY_600, line_spacing=1.4)
    tag_r = add_rounded_rect(s, left + Inches(0.15), top + Inches(1.72), Inches(0.9), Inches(0.3), tbg)
    add_text_box(s, left + Inches(0.15), top + Inches(1.72), Inches(0.9), Inches(0.3),
                 tag, font_size=8, color=tc, bold=True, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Railway card - 6th item
row = 1; col = 2
left = Inches(0.7) + col * Inches(4.1)
top = Inches(1.4) + row * Inches(2.7)
bg_c = add_rounded_rect(s, left, top, Inches(3.7), Inches(2.3), WHITE)
bg_c.line.color.rgb = GRAY_100
bg_c.line.fill.solid()
bg_c.line.width = Pt(1)
add_rounded_rect(s, left + Inches(0.15), top + Inches(0.15), Inches(0.5), Inches(0.5), GREEN_100)
add_text_box(s, left + Inches(0.18), top + Inches(0.17), Inches(0.44), Inches(0.44),
             '☁️', font_size=18, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text_box(s, left + Inches(0.8), top + Inches(0.15), Inches(2.6), Inches(0.4),
             'Railway 部署', font_size=13, color=GRAY_800, bold=True)
add_text_box(s, left + Inches(0.15), top + Inches(0.8), Inches(3.4), Inches(1.0),
             '云端一键部署，持久化卷挂载数据\n环境变量配置\n优雅关停SIGTERM/SIGINT', font_size=9, color=GRAY_600, line_spacing=1.4)
tag_r = add_rounded_rect(s, left + Inches(0.15), top + Inches(1.72), Inches(0.9), Inches(0.3), GREEN_100)
add_text_box(s, left + Inches(0.15), top + Inches(1.72), Inches(0.9), Inches(0.3),
             '运维', font_size=8, color=GREEN_700, bold=True, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ===== SLIDE 6: Pages & Nav =====
s = make_slide('页面路由与导航体系', 'Pages & Navigation')
add_text_box(s, Inches(0.7), Inches(1.3), Inches(5), Inches(0.35),
             'TABBAR 主页（4个）', font_size=12, color=GREEN_600, bold=True)
tab_data = [
    ['1', '🏠', '首页', 'pages/index/index', '轮播图、功能模块、热门推荐、天气'],
    ['2', '⛰️', '景点导览', 'pages/scenic/scenic', '7大景点、分类筛选、老人友好评分'],
    ['3', '🛡️', '安全防护', 'pages/safety/safety', '天气预警、路滑评估、出行建议'],
    ['4', '👤', '个人中心', 'pages/user/user', '登录、订单、收藏、SOS、语音设置'],
]
# simple table header
for j, (hdr_txt, w) in enumerate([('序号', 0.7), ('图标', 0.6), ('名称', 1.5), ('路径', 2.5), ('核心功能', 3.5)]):
    l = Inches(0.7) + sum([0.7, 0.6, 1.5, 2.5, 3.5][:j])
    add_text_box(s, l, Inches(1.7), Inches(w), Inches(0.4),
                 hdr_txt, font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, l, Inches(1.7), Inches(w), Inches(0.4), GREEN_700)
for i, row in enumerate(tab_data):
    widths = [0.7, 0.6, 1.5, 2.5, 3.5]
    for j, (cell, w) in enumerate(zip(row, widths)):
        l = Inches(0.7) + sum(widths[:j])
        t = Inches(2.15) + i * Inches(0.38)
        add_rect(s, l, t, Inches(w), Inches(0.35), GRAY_50 if i % 2 == 1 else WHITE)
        add_rect(s, l, t, Inches(w), Inches(0.35), None, GRAY_100, Pt(0.5))
        add_text_box(s, l + Inches(0.05), t + Inches(0.03), Inches(w) - Inches(0.1), Inches(0.3),
                     cell, font_size=9, color=GRAY_800)

# auxiliary pages
add_text_box(s, Inches(0.7), Inches(3.8), Inches(8), Inches(0.35),
             '辅助功能页 + 电商页（14个）', font_size=12, color=GREEN_600, bold=True)
aux = [
    '🗺️ 地图导览', '🍜 桂林美食', '🚌 交通出行', '🚦 实时路况', '💚 健康记录',
    '🆘 紧急求助', '🤖 智能助手', '📂 分类浏览', '🔍 商品搜索', '📋 商品列表',
    '🛒 购物车',   '📦 订单列表', '📍 地址管理', '📝 商品详情',
]
for i, txt in enumerate(aux):
    col = i % 5
    row = i // 5
    l = Inches(0.7) + col * Inches(2.44)
    t = Inches(4.25) + row * Inches(0.75)
    c = add_rounded_rect(s, l, t, Inches(2.2), Inches(0.55), WHITE)
    c.line.color.rgb = GRAY_100
    c.line.fill.solid()
    c.line.width = Pt(1)
    add_text_box(s, l, t, Inches(2.2), Inches(0.55),
                 txt, font_size=10, color=GRAY_800, bold=True, alignment=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)

# ===== SLIDE 7: Home Page =====
s = make_slide('首页设计与功能模块', 'Home Page')
zones = [
    ('区域1：安全栏 .safety-bar', '天气emoji + 实时温度 + "桂林·银发宜居"标签，来自Open-Meteo/高德API'),
    ('区域2：轮播图 .swiper-box', '360rpx高度，3张精美Banner（b1.jpg/b2.jpg/b3.jpg），图片加载失败自动降级'),
    ('区域3：功能模块网格 .module-grid', '2行×4列，8个入口：景区导览、地图导览、桂林美食、交通出行、安全保障、健康记录、紧急求助、全部商品'),
    ('区域4：热门推荐 .hot-section', '2×3商品卡片网格，支持骨架屏、错误重试、空状态三种UI'),
]
for i, (z_title, z_desc) in enumerate(zones):
    t = Inches(1.35) + i * Inches(0.85)
    highlight_box(s, Inches(0.6), t, Inches(5.8), Inches(0.75), z_title, z_desc)
# right - loading strategy
add_text_box(s, Inches(6.8), Inches(1.35), Inches(5.5), Inches(0.35),
             '数据加载策略', font_size=13, color=GREEN_700, bold=True)
strats = [
    '📡 天气：onLoad + onShow 双重加载，本地fallback兜底（22°C多云）',
    '📡 热门推荐：API优先 → 失败后6条本地硬编码兜底',
    '📡 分类数据：onLaunch预加载至全局缓存',
    '📡 轮播图：静态资源，binderror自动移除损坏图片',
]
for i, st in enumerate(strats):
    add_text_box(s, Inches(6.8), Inches(1.75) + i * Inches(0.35), Inches(5.5), Inches(0.3),
                 st, font_size=10, color=GRAY_800)
# state machine
add_text_box(s, Inches(6.8), Inches(3.3), Inches(5.5), Inches(0.35),
             '热门推荐状态机', font_size=13, color=GREEN_700, bold=True)
states = [
    'loadingHot=true          → 显示骨架屏动画',
    'loadingHot=false + error → 显示错误状态 + 重新加载按钮',
    'loadingHot=false + 空数据 → 显示空状态提示',
    'loadingHot=false + 有数据 → 显示2×3商品卡片网格',
]
for i, st in enumerate(states):
    add_text_box(s, Inches(6.8), Inches(3.7) + i * Inches(0.32), Inches(5.5), Inches(0.28),
                 st, font_size=9, color=GRAY_600, font_name=FONT_MONO)

# ===== SLIDE 8: Scenic =====
s = make_slide('景点导览模块详解', 'Scenic Spots')
spots = [
    ['漓江精华段', '5A', '¥215', '4.2', '0级'],
    ['象鼻山公园', '5A', '¥55',  '3.8', '~80级'],
    ['两江四湖',   '5A', '¥220', '4.3', '0级'],
    ['阳朔西街',   '4A', '免费', '4.5', '0级'],
    ['龙脊梯田',   '4A', '¥80',  '2.5', '~200级'],
    ['芦笛岩',     '4A', '¥90',  '3.5', '~60级'],
    ['杨堤码头',   '精华', '¥160', '3.9', '0级'],
]
t_headers = ['景点', '等级', '价格', '老人评分', '台阶']
t_widths = [1.8, 0.8, 0.8, 0.9, 0.8]
t_top = Inches(1.35)
for j, (hdr, w) in enumerate(zip(t_headers, t_widths)):
    l = Inches(0.6) + sum(t_widths[:j])
    add_rect(s, l, t_top, Inches(w), Inches(0.38), GREEN_700)
    add_text_box(s, l, t_top, Inches(w), Inches(0.38),
                 hdr, font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
for i, row in enumerate(spots):
    for j, (cell, w) in enumerate(zip(row, t_widths)):
        l = Inches(0.6) + sum(t_widths[:j])
        t = t_top + Inches(0.38) + i * Inches(0.36)
        add_rect(s, l, t, Inches(w), Inches(0.35), GRAY_50 if i % 2 == 0 else WHITE)
        add_rect(s, l, t, Inches(w), Inches(0.35), None, GRAY_200, Pt(0.5))
        add_text_box(s, l, t, Inches(w), Inches(0.35),
                     cell, font_size=9, color=GRAY_800, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# right info
add_text_box(s, Inches(6.3), Inches(1.35), Inches(6.2), Inches(0.35),
             '老人友好信息体系', font_size=13, color=GREEN_700, bold=True)
items_elderly = [
    '⭐ 老人友好评分：综合步行距离、台阶数、休息区打分',
    '🚶 步行距离：精确到百米级别',
    '🪜 台阶数量：详细描述台阶数和难度',
    '🪑 休息区信息：标明是否有休息亭/座椅',
    '⚠️ 安全提示：雨天湿滑、台阶陡峭等针对性提醒',
    '🚌 交通指南：公交路线和站点信息',
    '🕐 开放时间：详细的营业时间段',
]
for i, item in enumerate(items_elderly):
    add_text_box(s, Inches(6.3), Inches(1.75) + i * Inches(0.32), Inches(6.2), Inches(0.28),
                 item, font_size=9, color=GRAY_800)

add_text_box(s, Inches(6.3), Inches(4.05), Inches(6.2), Inches(0.35),
             '交互功能', font_size=13, color=GREEN_700, bold=True)
add_text_box(s, Inches(6.3), Inches(4.42), Inches(6.2), Inches(0.25),
             '👆 展开/收起：点击卡片展开详情，再次点击收起', font_size=9)
add_text_box(s, Inches(6.3), Inches(4.72), Inches(6.2), Inches(0.25),
             '❤️ 收藏系统：基于localStorage，跨页实时同步', font_size=9)
add_text_box(s, Inches(6.3), Inches(5.02), Inches(6.2), Inches(0.25),
             '🔊 语音播报：展开时自动朗读景点介绍', font_size=9)

# ===== SLIDE 9: Safety =====
s = make_slide('安全防护模块详解', 'Safety Center')
# left cards
safety_cards = [
    ('🌤️ 天气与路况卡片', '实时显示桂林天气（温度/湿度/风向），集成路滑风险评级（3级：安全/注意/预警），根据天气自动匹配出行建议'),
    ('📋 出行建议', '智能推荐/避开列表：晴天推荐象鼻山+两江四湖+阳朔西街，雨天推荐芦笛岩+靖江王城室内景点，避开龙脊梯田+杨堤竹筏'),
    ('🚨 紧急操作三按钮', '一键拨打120急救电话 / 一键获取GPS位置并复制分享链接 / 跳转紧急联系人管理页面'),
]
for i, (h4, desc) in enumerate(safety_cards):
    t = Inches(1.35) + i * Inches(1.25)
    highlight_box(s, Inches(0.6), t, Inches(5.8), Inches(1.05), h4, desc)

# right - safety tips
add_text_box(s, Inches(6.8), Inches(1.35), Inches(5.5), Inches(0.35),
             '银发安全提示（6条分级）', font_size=12, color=GREEN_700, bold=True)
tips = [
    ['⚠️ 高', '景区湿滑台阶', '使用扶手，一步一阶，勿停留拍照'],
    ['⚠️ 高', '雨后石板路', '绕行或等地面干燥，穿防滑鞋'],
    ['⚡ 中', '上下旅游巴士', '抓牢车门扶手，车停稳后再上下'],
    ['⚡ 中', '漓江游船甲板', '不走到甲板边缘，穿平底鞋'],
    ['✅ 低', '酒店卫生间', '使用防滑拖鞋，铺防滑垫'],
    ['✅ 低', '购物街区地砖', '注意积水和油渍，走路不看手机'],
]
for i, (level, scene, advice) in enumerate(tips):
    t = Inches(1.8) + i * Inches(0.35)
    add_text_box(s, Inches(6.8), t, Inches(1.0), Inches(0.3),
                 level, font_size=8, color=GRAY_800, bold=True)
    add_text_box(s, Inches(7.8), t, Inches(2.0), Inches(0.3),
                 scene, font_size=8, color=GRAY_800)
    add_text_box(s, Inches(9.8), t, Inches(2.8), Inches(0.3),
                 advice, font_size=8, color=GRAY_600)

add_text_box(s, Inches(6.8), Inches(4.1), Inches(5.5), Inches(0.35),
             '健康提醒（5条）', font_size=12, color=GREEN_700, bold=True)
health_tips = ['💧 携带饮用水', '🌂 随身带折叠伞/遮阳帽', '💊 常用药品随身带',
               '🕐 每游览1小时休息10-15分钟', '📱 保持手机电量，开启定位']
for i, ht in enumerate(health_tips):
    add_text_box(s, Inches(6.8), Inches(4.5) + i * Inches(0.3), Inches(5.5), Inches(0.25),
                 ht, font_size=9, color=GRAY_800)

# ===== SLIDE 10: User Center =====
s = make_slide('个人中心模块详解', 'User Center')
# menu
add_text_box(s, Inches(0.7), Inches(1.35), Inches(5.5), Inches(0.35),
             '菜单功能（6项）', font_size=13, color=GREEN_700, bold=True)
menus = [
    ('📋 我的订单', '订单列表 + 创建订单 + 模拟支付'),
    ('📍 地址管理', '收货地址增删改 + 默认地址'),
    ('❤️ 我的收藏', '收藏景点列表 + 一键清空'),
    ('🏥 健康记录', '血压/心率/备注记录'),
    ('👨‍👩‍👧 紧急联系人', '最多3个，支持设为主联系人'),
    ('🛡️ 安全设置', '跳转安全防护页面'),
]
for i, (h4, desc) in enumerate(menus):
    col = i % 2
    row = i // 2
    l = Inches(0.7) + col * Inches(3.1)
    t = Inches(1.8) + row * Inches(1.5)
    bg = add_rounded_rect(s, l, t, Inches(2.85), Inches(1.2), WHITE)
    bg.line.color.rgb = GRAY_100
    bg.line.fill.solid()
    bg.line.width = Pt(1)
    add_text_box(s, l + Inches(0.1), t + Inches(0.1), Inches(2.65), Inches(0.35),
                 h4, font_size=11, color=GRAY_800, bold=True)
    add_text_box(s, l + Inches(0.1), t + Inches(0.5), Inches(2.65), Inches(0.6),
                 desc, font_size=9, color=GRAY_600)

# right - settings
add_text_box(s, Inches(7.0), Inches(1.35), Inches(5.5), Inches(0.35),
             '设置功能', font_size=13, color=GREEN_700, bold=True)
settings = [
    ('🔊 语音提示开关', '控制全局TTS语音播报，开启后页面切换/按钮点击均有语音反馈'),
    ('🌐 语言辅助（4语切换）', '中文 / English / 日本語 / 한국어，切换后广播到所有活跃页面'),
    ('🔤 字体缩放（3档）', '标准(默认) | 较大(1.2倍) | 超大(1.5倍)，CSS变量全局响应'),
    ('🆘 SOS紧急求助', '一键获取GPS位置 → 调用API通知紧急联系人 → 自动拨打'),
    ('🔐 微信登录/登出', 'wx.login获取code → /users/wxlogin换取openid → 持久化'),
]
for i, (h4, desc) in enumerate(settings):
    t = Inches(1.8) + i * Inches(0.85)
    highlight_box(s, Inches(7.0), t, Inches(5.5), Inches(0.72), h4, desc)

# ===== SLIDE 11: Aux Pages =====
s = make_slide('辅助功能页面全景', 'Auxiliary Pages')
aux_pages = [
    ('🗺️', '地图导览', '高德地图SDK，桂林景点标记，GPS定位+路径导航，fallback-first渲染'),
    ('🍜', '桂林美食', '桂林特色美食推荐（米粉、啤酒鱼等），图文展示+价格信息'),
    ('🚌', '交通出行', '公交路线、旅游大巴、自驾路线等多种交通方案推荐'),
    ('🚦', '实时路况', '景区周边路况信息展示，帮助老人合理规划出行时间和路线'),
    ('💚', '健康记录', '血压记录/心率记录，本地+服务端双存储'),
    ('🆘', '紧急联系人', '最多3个，姓名+电话+关系，支持is_primary'),
    ('🤖', '智能助手', 'AI智能问答助手，帮助老人解答旅行相关疑问'),
    ('📂', '分类浏览', '按分类浏览：景点/住宿/美食/线路/特产'),
    ('🔍', '商品搜索', '关键词搜索+防抖(300ms)+历史记录+LIK模糊匹配'),
]
for i, (icon, h4, desc) in enumerate(aux_pages):
    col = i % 3
    row = i // 3
    l = Inches(0.6) + col * Inches(4.15)
    t = Inches(1.4) + row * Inches(2.0)
    bg = add_rounded_rect(s, l, t, Inches(3.8), Inches(1.7), WHITE)
    bg.line.color.rgb = GRAY_100
    bg.line.fill.solid()
    bg.line.width = Pt(1)
    add_rounded_rect(s, l + Inches(0.15), t + Inches(0.15), Inches(0.45), Inches(0.45), GREEN_100)
    add_text_box(s, l + Inches(0.17), t + Inches(0.17), Inches(0.41), Inches(0.41),
                 icon, font_size=18, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(s, l + Inches(0.15), t + Inches(0.72), Inches(3.5), Inches(0.35),
                 h4, font_size=12, color=GRAY_800, bold=True)
    add_text_box(s, l + Inches(0.15), t + Inches(1.1), Inches(3.5), Inches(0.5),
                 desc, font_size=8, color=GRAY_600, line_spacing=1.3)

# ===== SLIDE 12: E-Commerce =====
s = make_slide('电商功能模块', 'E-Commerce')
add_text_box(s, Inches(0.7), Inches(1.35), Inches(5.5), Inches(0.35),
             '完整电商闭环', font_size=13, color=GREEN_700, bold=True)
flow_steps = ['分类浏览', '商品搜索', '商品列表', '商品详情', '加入购物车', '创建订单', '模拟支付']
for i, step in enumerate(flow_steps):
    l = Inches(0.7) + i * Inches(1.55)
    bg = add_rounded_rect(s, l, Inches(1.8), Inches(1.25), Inches(0.42), GREEN_100)
    bg.line.color.rgb = GREEN_400
    bg.line.fill.solid()
    bg.line.width = Pt(1)
    add_text_box(s, l, Inches(1.8), Inches(1.25), Inches(0.42),
                 step, font_size=8, color=GREEN_700, bold=True, alignment=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)

add_text_box(s, Inches(0.7), Inches(2.45), Inches(5.5), Inches(0.35),
             '5大商品分类', font_size=13, color=GREEN_700, bold=True)
cats = ['🏔️ 漓江景点', '🏨 酒店住宿', '🍜 桂林美食', '🚩 旅游线路', '🎁 桂林特产']
for i, cat in enumerate(cats):
    l = Inches(0.7) + i * Inches(2.0)
    add_text_box(s, l, Inches(2.85), Inches(1.8), Inches(0.35),
                 cat, font_size=10, color=GREEN_700, bold=True)

add_text_box(s, Inches(0.7), Inches(3.3), Inches(5.5), Inches(0.35),
             '购物车管理', font_size=13, color=GREEN_700, bold=True)
cart_items = [
    '➕ 加入购物车：POST /cart/add，已存在则累加数量',
    '✏️ 修改数量：PUT /cart/:itemId，quantity≤0时自动删除',
    '🗑️ 删除商品：DELETE /cart/:itemId', '🧹 清空购物车：POST /cart/clear',
]
for i, ci in enumerate(cart_items):
    add_text_box(s, Inches(0.7), Inches(3.7) + i * Inches(0.3), Inches(5.5), Inches(0.25),
                 ci, font_size=9, color=GRAY_800)

# right - order flow
add_text_box(s, Inches(7.0), Inches(1.35), Inches(5.5), Inches(0.35),
             '订单流程', font_size=13, color=GREEN_700, bold=True)
order_steps = [
    '创建订单 → 生成订单号(ORD+时间戳)',
    '→ 计算总价(quantity × unit_price)',
    '→ 状态: pending(待支付)',
    '→ POST /orders/:id/paid 模拟支付',
    '→ 状态: paid(已支付)',
]
for i, os in enumerate(order_steps):
    add_text_box(s, Inches(7.0), Inches(1.75) + i * Inches(0.28), Inches(5.5), Inches(0.25),
                 os, font_size=9, color=GRAY_600, font_name=FONT_MONO)

add_text_box(s, Inches(7.0), Inches(3.3), Inches(5.5), Inches(0.35),
             '地址管理', font_size=13, color=GREEN_700, bold=True)
add_text_box(s, Inches(7.0), Inches(3.72), Inches(5.5), Inches(0.25),
             '📝 创建地址：姓名 + 手机号(Zod正则验证) + 详细地址', font_size=9)
add_text_box(s, Inches(7.0), Inches(4.05), Inches(5.5), Inches(0.25),
             '⭐ 默认地址：设置is_default时自动取消其他默认', font_size=9)


# ===== SLIDE 13: Elderly Design =====
s = make_slide('银发友好设计体系', 'Elderly-Friendly Design')
design_cards = [
    ('🔤', '三档字体缩放', '标准 (normal)：默认字体大小\n较大 (large)：整体放大1.2倍\n超大 (huge)：整体放大1.5倍\n\n实现：CSS类选择器 + getCurrentPages()广播\n.font-mode-small / .font-mode-large / .font-mode-huge'),
    ('🔊', '全程语音播报', '主方案：WechatSI插件 (textToSpeech)\n备选方案：有道词典TTS API\n\n触发场景：页面切换、按钮点击、景点展开\n安全提示朗读、SOS求助确认\n支持4语种：zh_CN / en_US / ja_JP / ko_KR'),
    ('🌐', '语言辅助切换', '中文 | English | 日本語 | 한국어\n\n存储到globalData + wx.setStorageSync持久化\n通过getCurrentPages()广播到所有活跃页面\n影响：界面文案、语音播报语言'),
]
for i, (icon, h4, desc) in enumerate(design_cards):
    l = Inches(0.6) + i * Inches(4.2)
    bg = add_rounded_rect(s, l, Inches(1.4), Inches(3.85), Inches(3.6), WHITE)
    bg.line.color.rgb = GRAY_100
    bg.line.fill.solid()
    bg.line.width = Pt(1)
    add_rounded_rect(s, l + Inches(0.15), t := Inches(1.55), Inches(0.45), Inches(0.45), GREEN_100)
    add_text_box(s, l + Inches(0.17), t, Inches(0.41), Inches(0.41),
                 icon, font_size=18, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(s, l + Inches(0.15), Inches(2.1), Inches(3.55), Inches(0.35),
                 h4, font_size=12, color=GRAY_800, bold=True)
    add_text_box(s, l + Inches(0.15), Inches(2.5), Inches(3.55), Inches(2.3),
                 desc, font_size=9, color=GRAY_600, line_spacing=1.4)

# font flow
highlight_box(s, Inches(0.6), Inches(5.2), Inches(12.1), Inches(1.2),
              '字体切换流程',
              '用户点击 → user.js setFontSize() → app.updateFontSize(mode) → globalData更新 → Storage持久化 → getCurrentPages()广播 → 各页面setData({fontSizeMode}) → CSS类名变化 → 字体生效')

# ===== SLIDE 14: Safety Architecture =====
s = make_slide('安全防护体系架构', 'Safety Architecture')
layers_safety = [
    ('🌤️', '第一层：天气实时监测', 'Open-Meteo免费API → 高德天气API(Fallback) → 本地静态数据兜底，三重保障。30分钟缓存。'),
    ('⚠️', '第二层：路滑风险评估', '根据天气状况自动评估3级风险：level 1(安全/晴朗) → level 2(注意/阴湿) → level 3(预警/雨雪)'),
    ('📋', '第三层：智能出行建议', '根据天气自动推荐适合景点(suitable列表)和避开景点(avoid列表)'),
    ('🛡️', '第四层：6条分级安全提示', '高危(⚠️)：景区湿滑台阶+雨后石板路 | 中危(⚡)：上下巴士+游船甲板 | 低危(✅)：酒店+街区'),
    ('💚', '第五层：5条健康提醒', '补水、防晒防雨、随身药品、定时休息、保持通讯——覆盖银发出游全场景'),
    ('🆘', '第六层：SOS紧急求助', 'GPS定位 + 通知所有紧急联系人 + 多重fallback(联系人失败→110，位置失败→直接拨110)'),
]
for i, (icon, h5, desc) in enumerate(layers_safety):
    t = Inches(1.4) + i * Inches(0.82)
    add_line(s, Inches(0.6), t, Pt(4), Inches(0.65), GREEN_500, 4)
    add_text_box(s, Inches(0.8), t + Inches(0.02), Inches(0.5), Inches(0.5),
                 icon, font_size=18, alignment=PP_ALIGN.CENTER)
    add_text_box(s, Inches(1.4), t + Inches(0.04), Inches(4.5), Inches(0.3),
                 h5, font_size=11, color=GRAY_800, bold=True)
    add_text_box(s, Inches(1.4), t + Inches(0.35), Inches(11), Inches(0.3),
                 desc, font_size=8, color=GRAY_600)

# ===== SLIDE 15: SOS =====
s = make_slide('紧急求助SOS流程', 'SOS Flow')
add_text_box(s, Inches(0.7), Inches(1.35), Inches(12), Inches(0.4),
             '正常路径：点击SOS → 二次确认 → wx.getLocation获取GPS → POST /sos/alert → 通知紧急联系人 → 语音播报 → 弹窗拨打',
             font_size=11, color=GRAY_800, alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(0.7), Inches(1.85), Inches(5.5), Inches(0.35),
             '✅ 定位成功路径', font_size=13, color=GREEN_700, bold=True)
success_steps = [
    'Step 1：调用 POST /sos/alert 发送 {openid, latitude, longitude}',
    'Step 2：后端查询is_primary=1的紧急联系人',
    'Step 3：语音播报"已获取位置并通知紧急联系人"',
    'Step 4a：有主联系人 → 弹窗显示姓名+电话 → 确认拨打',
    'Step 4b：无联系人 → 直接拨打110',
]
for i, st in enumerate(success_steps):
    highlight_box(s, Inches(0.7), Inches(2.3) + i * Inches(0.7), Inches(5.5), Inches(0.6), st, '')

add_text_box(s, Inches(7.0), Inches(1.85), Inches(5.5), Inches(0.35),
             '❌ 异常处理路径', font_size=13, color=RED, bold=True)
error_steps = [
    ('API调用失败', '语音播报 → 直接拨打110'),
    ('GPS定位失败', '语音播报 → 直接拨打110'),
    ('无紧急联系人', '跳过联系人通知 → 直接拨打110'),
]
for i, (title, desc) in enumerate(error_steps):
    l = Inches(7.0)
    t = Inches(2.3) + i * Inches(0.8)
    add_line(s, l, t, Pt(4), Inches(0.55), RED, 4)
    add_text_box(s, l + Inches(0.2), t + Inches(0.02), Inches(5.2), Inches(0.28),
                 title, font_size=11, color=RED, bold=True)
    add_text_box(s, l + Inches(0.2), t + Inches(0.3), Inches(5.2), Inches(0.25),
                 desc, font_size=9, color=GRAY_600)

# ===== SLIDE 16: Backend =====
s = make_slide('后端架构与安全措施', 'Backend & Security')
add_text_box(s, Inches(0.7), Inches(1.35), Inches(5.5), Inches(0.35),
             'Express服务器配置', font_size=13, color=GREEN_700, bold=True)
configs = [
    '🛡️ Helmet：HTTP安全头（CSP/X-Frame/XSS防护等）',
    '📦 Compression：gzip压缩(level 6)，减少传输量',
    '🌐 CORS：可配置允许域名，支持credentials',
    '🚦 Rate Limit：每IP每15分钟100次请求',
    '📏 Body Limit：请求体限制1MB',
    '📝 Pino日志：结构化JSON日志输出',
    '🖼️ 静态文件：图片7天缓存，支持CDN配置',
    '💚 健康检查：/health端点返回状态+时间戳',
]
for i, cfg in enumerate(configs):
    add_text_box(s, Inches(0.7), Inches(1.75) + i * Inches(0.32), Inches(5.5), Inches(0.28),
                 cfg, font_size=9, color=GRAY_800)

add_text_box(s, Inches(7.0), Inches(1.35), Inches(5.5), Inches(0.35),
             '中间件体系', font_size=13, color=GREEN_700, bold=True)
mw_items = [
    ('🔍 requestLogger 请求日志', '记录每个请求的method、url、状态码、响应时间'),
    ('✅ Zod 参数验证', '7个验证Schema + 手机号正则 /^1[3-9]\\\\d{9}$/'),
    ('📋 auditLog 审计日志', '记录敏感操作：login、create_order、create_address'),
    ('💾 SQL参数化查询', 'db.prepare() + stmt.bind() 防止SQL注入'),
]
for i, (h4, desc) in enumerate(mw_items):
    t = Inches(1.8) + i * Inches(1.0)
    highlight_box(s, Inches(7.0), t, Inches(5.5), Inches(0.8), h4, desc)

# ===== SLIDE 17: Database =====
s = make_slide('数据库设计（9张表）', 'Database Schema')
tables = [
    ['categories', '商品分类', 'id, name, slug(唯一), order(排序)', '-'],
    ['products', '商品/景点', 'id, name, price, stock, image_url, is_active', 'category_id'],
    ['user_profiles', '用户信息', 'id, openid(唯一), nickname, avatar_url, phone', '-'],
    ['addresses', '收货地址', 'id, full_name, phone, address_line, city, is_default', 'user_id'],
    ['orders', '订单', 'id, order_no(唯一), total_price, status(pending→paid)', 'user_id'],
    ['order_items', '订单明细', 'id, quantity, unit_price', 'order_id, product_id'],
    ['cart_items', '购物车', 'id, quantity, created_at', 'user_id, product_id'],
    ['emergency_contacts', '紧急联系人', 'id, name, phone, relationship, is_primary', 'user_id'],
    ['health_records', '健康记录', 'id, blood_pressure, heart_rate, notes, record_date', 'user_id'],
]
cols = ['表名', '用途', '核心字段', '外键']
c_widths = [1.8, 1.5, 4.5, 2.0]
t_top = Inches(1.3)
for j, (hdr, w) in enumerate(zip(cols, c_widths)):
    l = Inches(0.6) + sum(c_widths[:j])
    add_rect(s, l, t_top, Inches(w), Inches(0.4), GREEN_700)
    add_text_box(s, l, t_top, Inches(w), Inches(0.4),
                 hdr, font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
for i, row in enumerate(tables):
    for j, (cell, w) in enumerate(zip(row, c_widths)):
        l = Inches(0.6) + sum(c_widths[:j])
        t = t_top + Inches(0.4) + i * Inches(0.38)
        add_rect(s, l, t, Inches(w), Inches(0.36), GRAY_50 if i % 2 == 0 else WHITE)
        add_rect(s, l, t, Inches(w), Inches(0.36), None, GRAY_200, Pt(0.5))
        fs = 8 if j == 2 else 9
        add_text_box(s, l + Inches(0.04), t + Inches(0.02), Inches(w) - Inches(0.08), Inches(0.32),
                     cell, font_size=fs, color=GRAY_800)

highlight_box(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.4),
              '数据库特性',
              'SQL.js — SQLite编译为WebAssembly，零外部依赖，嵌入式运行  |  '
              '自动建表 — initDatabase()自动创建9张表 + 13条示例数据  |  '
              '自动备份 — 每小时备份到 backups/，保留最近24份  |  '
              '灾难恢复 — 主库损坏时自动从最新备份恢复')

# ===== SLIDE 18: API =====
s = make_slide('API接口全景（20+端点）', 'API Endpoints')
api_groups = {
    '商品模块': [
        ('GET', '/api/categories', '分类列表(5分钟缓存)'),
        ('GET', '/api/products', '商品列表(多条件+分页)'),
        ('GET', '/api/products/:id', '商品详情'),
    ],
    '用户模块': [
        ('POST', '/api/users/wxlogin', '微信登录(code→openid)'),
        ('POST', '/api/users/login', '登录/注册(审计日志)'),
        ('GET', '/api/users/:openid', '用户详情'),
    ],
    '购物车 & 订单': [
        ('GET/POST', '/api/cart', '查看/添加购物车'),
        ('PUT/DELETE', '/api/cart/:itemId', '修改数量/删除'),
        ('POST', '/api/cart/clear', '清空购物车'),
        ('GET', '/api/orders', '订单列表'),
        ('POST', '/api/order/create', '创建订单(审计日志)'),
        ('POST', '/api/orders/:id/paid', '模拟支付'),
    ],
    '地址 & 天气': [
        ('GET', '/api/addresses', '地址列表'),
        ('POST', '/api/address/create', '创建地址'),
        ('GET', '/api/weather', '实时天气(30分钟缓存)'),
        ('GET', '/api/weather/forecast', '天气预报(1-7天)'),
    ],
    '紧急联系人': [
        ('GET', '/api/emergency-contacts', '联系人列表'),
        ('POST', '/api/emergency-contacts/create', '添加(最多3个)'),
        ('PUT', '/api/emergency-contacts/:id', '设为主联系人'),
        ('DELETE', '/api/emergency-contacts/:id', '删除联系人'),
    ],
    '健康 & SOS': [
        ('GET', '/api/health-records', '健康记录列表'),
        ('POST', '/api/health-records', '添加健康记录'),
        ('POST', '/api/sos/alert', '触发SOS求助'),
    ],
}
y = Inches(1.25)
for group_name, endpoints in api_groups.items():
    add_text_box(s, Inches(0.7), y, Inches(3), Inches(0.3),
                 group_name, font_size=10, color=GREEN_600, bold=True)
    y += Inches(0.3)
    for method, path, desc in endpoints:
        add_text_box(s, Inches(0.7), y, Inches(0.8), Inches(0.22),
                     method, font_size=8, color=GREEN_700, bold=True, font_name=FONT_MONO)
        add_text_box(s, Inches(1.5), y, Inches(3.5), Inches(0.22),
                     path, font_size=8, color=GRAY_800, font_name=FONT_MONO)
        add_text_box(s, Inches(5.2), y, Inches(4.0), Inches(0.22),
                     desc, font_size=8, color=GRAY_600)
        y += Inches(0.24)
    y += Inches(0.15)
    if y > Inches(6.5) and group_name != list(api_groups.keys())[-1]:
        # second column
        pass

highlight_box(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.6),
              '统一响应格式',
              '成功：{"code": 0, "data": ...}  |  失败：{"code": 1, "error": "..."}')

# ===== SLIDE 19: Weather Fallback =====
s = make_slide('天气API多级Fallback机制', 'Weather Fallback')
weather_layers = [
    ('🌍', 'Open-Meteo', '免费全球天气API\n无需API Key\n经纬度查询', '主数据源', GREEN_600, GREEN_100),
    ('🗺️', '高德天气API', '城市编码查询\n需AMAP_KEY\nrestapi.amap.com', '备选1', GOLD, RGBColor(0xFE,0xF3,0xC7)),
    ('📦', '本地Fallback', '静态数据兜底\n多云 22°C 65%\n南风2级 桂林', '兜底', RED, RGBColor(0xFE,0xE2,0xE2)),
]
for i, (icon, h5, items, tag, tc, tbg) in enumerate(weather_layers):
    l = Inches(0.7) + i * Inches(4.2)
    if i > 0:
        add_text_box(s, l - Inches(0.4), Inches(2.5), Inches(0.35), Inches(0.4),
                     '↓', font_size=16, color=GRAY_400, alignment=PP_ALIGN.CENTER)
    bg = add_rounded_rect(s, l, Inches(1.4), Inches(3.7), Inches(2.5), WHITE)
    bg.line.color.rgb = GREEN_400 if i == 0 else (GOLD if i == 1 else GRAY_400)
    bg.line.fill.solid()
    bg.line.width = Pt(1.5)
    add_text_box(s, l, Inches(1.55), Inches(3.7), Inches(0.6),
                 icon, font_size=28, alignment=PP_ALIGN.CENTER)
    add_text_box(s, l, Inches(2.2), Inches(3.7), Inches(0.4),
                 h5, font_size=13, color=GREEN_700, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(s, l + Inches(0.3), Inches(2.65), Inches(3.1), Inches(0.8),
                 items, font_size=9, color=GRAY_600, alignment=PP_ALIGN.CENTER, line_spacing=1.5)
    tag_r = add_rounded_rect(s, l + Inches(0.9), Inches(3.5), Inches(1.9), Inches(0.32), tbg)
    add_text_box(s, l + Inches(0.9), Inches(3.5), Inches(1.9), Inches(0.32),
                 tag, font_size=9, color=tc, bold=True, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text_box(s, Inches(0.7), Inches(4.2), Inches(5), Inches(0.35),
             '天气缓存策略', font_size=13, color=GREEN_700, bold=True)
add_text_box(s, Inches(0.7), Inches(4.6), Inches(5), Inches(0.25),
             '⏱️ 实时天气：30分钟缓存TTL', font_size=9)
add_text_box(s, Inches(0.7), Inches(4.9), Inches(5), Inches(0.25),
             '⏱️ 天气预报：同样30分钟缓存，支持1-7天', font_size=9)
add_text_box(s, Inches(0.7), Inches(5.2), Inches(5), Inches(0.25),
             '🌐 天气编码映射：Open-Meteo WMO Code → 中文天气描述', font_size=9)
add_text_box(s, Inches(0.7), Inches(5.5), Inches(5), Inches(0.25),
             '🧭 风向转换：角度度数 → 中文风向(北/东北/东...)', font_size=9)

# ===== SLIDE 20: Components & Utils =====
s = make_slide('组件化与工具模块', 'Components & Utilities')
add_text_box(s, Inches(0.7), Inches(1.35), Inches(5.5), Inches(0.35),
             '4个可复用组件', font_size=13, color=GREEN_700, bold=True)
comps = [
    ('🖼️ banner-view', '轮播图组件，支持自动播放、指示器、错误兜底'),
    ('⏳ loading-view', '加载中组件，带旋转动画'),
    ('❌ error-view', '错误状态组件，带重新加载按钮'),
    ('📭 empty-view', '空状态组件，带友好提示信息'),
    ('📊 custom-tab-bar', '4个Tab：首页/景点导览/安全防护/个人中心，支持暗黑模式'),
]
for i, (h4, desc) in enumerate(comps):
    t = Inches(1.8) + i * Inches(0.7)
    highlight_box(s, Inches(0.7), t, Inches(5.5), Inches(0.6), h4, desc)

add_text_box(s, Inches(7.0), Inches(1.35), Inches(5.5), Inches(0.35),
             '6个工具模块', font_size=13, color=GREEN_700, bold=True)
utils_ = [
    ('📡 api.js', '统一请求封装：loading计数、超时(8s)、401弹窗、17个API函数'),
    ('🔊 voice.js', 'TTS引擎：WechatSI(主)→有道(备)→静默降级，4语种'),
    ('🌤️ weather.js', '天气工具：图标映射、路滑评估(3级)、suitable/avoid列表'),
    ('🧭 navigate.js', '导航工具：go/switchTab/goBack，600ms语音延迟后跳转'),
    ('📐 page-base.js', '页面基类：状态机、字体同步、防重复加载'),
    ('🖼️ image-paths.js', '图片路径集中管理：景点/美食/模块/轮播/TabBar/兜底'),
]
for i, (h4, desc) in enumerate(utils_):
    t = Inches(1.8) + i * Inches(0.7)
    highlight_box(s, Inches(7.0), t, Inches(5.5), Inches(0.6), h4, desc)

# ===== SLIDE 21: Vue Admin =====
s = make_slide('Web管理后台（Vue 3）', 'Vue Admin Panel')
add_text_box(s, Inches(0.7), Inches(1.35), Inches(5.5), Inches(0.35),
             '技术架构', font_size=13, color=GREEN_700, bold=True)
highlight_box(s, Inches(0.7), Inches(1.8), Inches(5.5), Inches(0.8),
              'Vue 3 + TypeScript + Vite',
              'Composition API + <script setup>语法，Vite快速HMR开发体验')
highlight_box(s, Inches(0.7), Inches(2.8), Inches(5.5), Inches(0.8),
              'Pinia 状态管理',
              'auth store（openid/nickname/isLoggedIn）+ cart store + order store，localStorage持久化')
highlight_box(s, Inches(0.7), Inches(3.8), Inches(5.5), Inches(0.8),
              'Vue Router 路由守卫',
              'beforeEach检查requiresAuth，未登录重定向到/user页')

add_text_box(s, Inches(7.0), Inches(1.35), Inches(5.5), Inches(0.35),
             '10个管理页面', font_size=13, color=GREEN_700, bold=True)
pages_data = [['首页', '/', '否'], ['分类浏览', '/category', '否'], ['商品列表', '/list', '否'],
              ['商品详情', '/details', '否'], ['搜索', '/search', '否'], ['地图', '/map', '否'],
              ['用户中心', '/user', '否'], ['购物车', '/cart', '是 🔒'], ['订单列表', '/orders', '是 🔒'],
              ['地址管理', '/address', '是 🔒']]
for i, (page, path, auth) in enumerate(pages_data):
    col = i // 5
    row = i % 5
    l = Inches(7.0) + col * Inches(2.8)
    t = Inches(1.8) + row * Inches(0.5)
    add_text_box(s, l, t, Inches(1.2), Inches(0.35),
                 page, font_size=9, color=GRAY_800)
    add_text_box(s, l + Inches(1.2), t, Inches(1.2), Inches(0.35),
                 path, font_size=8, color=GRAY_600, font_name=FONT_MONO)
    add_text_box(s, l + Inches(2.4), t, Inches(0.4), Inches(0.35),
                 auth, font_size=9, color=RED if '🔒' in auth else GRAY_600)

# ===== SLIDE 22: Deployment =====
s = make_slide('部署架构与运维策略', 'Deployment & DevOps')
add_text_box(s, Inches(0.7), Inches(1.35), Inches(5.5), Inches(0.35),
             'Railway 一键部署', font_size=13, color=GREEN_700, bold=True)
dep_items = [
    ('🚀 部署流程', 'GitHub Push → Railway自动检测 → 读取railway.json → 安装依赖 → 启动服务'),
    ('📦 railway.json', 'build命令(npm install) + start命令(ts-node src/index.ts)'),
    ('💾 持久化卷', 'Mount Path: /app/data → 数据库+备份持久存储，容器重启不丢失'),
    ('🔧 环境变量', 'PORT、NODE_ENV、WX_APPID、WX_SECRET、AMAP_KEY、ALLOWED_ORIGINS'),
    ('🔄 优雅关停', 'SIGTERM/SIGINT → stopBackupTimer() → process.exit(0)'),
]
for i, (h4, desc) in enumerate(dep_items):
    t = Inches(1.8) + i * Inches(0.82)
    highlight_box(s, Inches(0.7), t, Inches(5.5), Inches(0.7), h4, desc)

add_text_box(s, Inches(7.0), Inches(1.35), Inches(5.5), Inches(0.35),
             '数据库备份策略', font_size=13, color=GREEN_700, bold=True)
highlight_box(s, Inches(7.0), Inches(1.8), Inches(5.5), Inches(0.9),
              '⏰ 自动备份机制',
              '每小时自动备份 → backups/shop_backup_*.db，保留最近24份')
highlight_box(s, Inches(7.0), Inches(2.9), Inches(5.5), Inches(0.9),
              '🔄 灾难恢复',
              '主DB加载失败 → 扫描backups → 按时间倒序尝试恢复 → 全部失败则新建DB并插入示例数据')

add_text_box(s, Inches(7.0), Inches(4.1), Inches(5.5), Inches(0.35),
             '目录结构', font_size=13, color=GREEN_700, bold=True)
tree = ('yinfa/\n'
        '├── yinfa-ts/          # 微信小程序\n'
        '│   ├── pages/         # 18页面\n'
        '│   ├── components/    # 4组件\n'
        '│   ├── utils/         # 6工具\n'
        '│   ├── frontend/      # Vue 3\n'
        '│   ├── backend/server/# Express\n'
        '│   └── image/         # 图片\n'
        '├── data/              # 数据(运行时)\n'
        '├── railway.json\n'
        '└── Procfile')
add_text_box(s, Inches(7.0), Inches(4.5), Inches(5.5), Inches(2.5),
             tree, font_size=8, color=GRAY_600, font_name=FONT_MONO, line_spacing=1.2)

# ===== SLIDE 23: State & CSS =====
s = make_slide('全局状态管理 & CSS变量系统', 'State & Variables')
add_text_box(s, Inches(0.7), Inches(1.35), Inches(5.5), Inches(0.35),
             'globalData 全局状态', font_size=13, color=GREEN_700, bold=True)
gdata = [
    ('userInfo', 'Object', '用户信息(openid/nickname/avatar/phone)'),
    ('voiceEnabled', 'Boolean', '语音开关'),
    ('currentLang', 'String', '当前语言(zh/en/ja/ko)'),
    ('fontSizeMode', 'String', '字体模式(normal/large/huge)'),
    ('cachedCategories', 'Array', '分类数据缓存'),
    ('loginChecked', 'Boolean', '登录状态已检查'),
]
for i, (field, typ, desc) in enumerate(gdata):
    t = Inches(1.8) + i * Inches(0.4)
    add_text_box(s, Inches(0.7), t, Inches(1.5), Inches(0.3),
                 field, font_size=9, color=GRAY_800, font_name=FONT_MONO)
    add_text_box(s, Inches(2.3), t, Inches(1.0), Inches(0.3),
                 typ, font_size=9, color=GRAY_600, font_name=FONT_MONO)
    add_text_box(s, Inches(3.3), t, Inches(3.0), Inches(0.3),
                 desc, font_size=9, color=GRAY_800)

add_text_box(s, Inches(0.7), Inches(4.4), Inches(5.5), Inches(0.5),
             '状态广播：字体/语言/语音变更 → getCurrentPages() → setData同步 → Storage持久化',
             font_size=9, color=GRAY_600)

add_text_box(s, Inches(7.0), Inches(1.35), Inches(5.5), Inches(0.35),
             'CSS变量设计系统', font_size=13, color=GREEN_700, bold=True)
css_vars = [
    ('--font-xs..xxl', '24~52rpx', '字号体系'),
    ('--primary', '#2E8B57 海绿', '主色调'),
    ('--accent', '#FF6B35 珊瑚橙', '强调色'),
    ('--sos', '#EF5350 警示红', 'SOS色'),
    ('--gold', '#FFD700 金色', '金牌色'),
    ('--radius-sm/md/lg', '8/12/16rpx', '圆角'),
]
for i, (var, val, desc) in enumerate(css_vars):
    t = Inches(1.8) + i * Inches(0.45)
    add_text_box(s, Inches(7.0), t, Inches(2.2), Inches(0.3),
                 var, font_size=9, color=GRAY_800, font_name=FONT_MONO)
    add_text_box(s, Inches(9.2), t, Inches(1.8), Inches(0.3),
                 val, font_size=9, color=GREEN_700)
    add_text_box(s, Inches(11.0), t, Inches(1.6), Inches(0.3),
                 desc, font_size=9, color=GRAY_600)

highlight_box(s, Inches(7.0), Inches(4.6), Inches(5.5), Inches(1.0),
              '关键设计决策',
              '所有CSS变量均配备硬值fallback：\nfont-size: var(--font-md, 32rpx)\n兼容不支持CSS变量的老旧设备')

# ===== SLIDE 24: Highlights =====
s = make_slide('技术亮点与创新总结', 'Highlights')
hl_items = [
    ('🎨', 'CSS变量系统', '统一设计Token，硬值fallback兼容老旧设备，完整字号体系'),
    ('🧩', '组件模块化', '4个通用组件，统一加载/错误/空状态处理'),
    ('⚡', 'Fallback-First渲染', '地图先展示本地marker，API失败兜底，杜绝白屏'),
    ('📡', '状态广播机制', '字体/语言/语音 → 遍历所有页面 → setData同步'),
    ('🛡️', '多层安全防护', 'Helmet+CORS+RateLimit+Zod+审计日志+SQL参数化'),
    ('🔄', '数据持久化', 'SQL.js+每小时备份24份+灾难恢复+持久化卷'),
    ('🌐', '多级API Fallback', '天气:Open-Meteo→高德→本地 | TTS:WechatSI→有道→静默'),
    ('👴', '银发友好设计', '三档字体+全程语音(4语种)+老人友好评分+6级安全提示'),
    ('📊', '状态机驱动UI', 'loading→error→empty→data，page-base.js标准化'),
]
for i, (icon, h4, desc) in enumerate(hl_items):
    col = i % 3
    row = i // 3
    l = Inches(0.6) + col * Inches(4.2)
    t = Inches(1.4) + row * Inches(2.0)
    bg = add_rounded_rect(s, l, t, Inches(3.85), Inches(1.7), WHITE)
    bg.line.color.rgb = GRAY_100
    bg.line.fill.solid()
    bg.line.width = Pt(1)
    ti = t / 914400
    add_rounded_rect(s, l + Inches(0.15), Inches(ti + 0.15), Inches(0.45), Inches(0.45), GREEN_100)
    add_text_box(s, l + Inches(0.17), Inches(ti + 0.17), Inches(0.41), Inches(0.41),
                 icon, font_size=18, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(s, l + Inches(0.15), Inches(ti + 0.72), Inches(3.55), Inches(0.35),
                 h4, font_size=12, color=GRAY_800, bold=True)
    add_text_box(s, l + Inches(0.15), Inches(ti + 1.1), Inches(3.55), Inches(0.5),
                 desc, font_size=8, color=GRAY_600, line_spacing=1.3)

# ===== SLIDE 25: Thank You =====
s = add_blank_slide()
add_bg(s, GREEN_600)
add_rect(s, 0, 0, W, Inches(4.0), GREEN_800)
add_text_box(s, 0, Inches(1.2), W, Inches(1.0),
             '感谢聆听', font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
thank_text = ('桂林银发旅游小程序 v1.0.0\n\n'
              '18个页面 · 9张数据库表 · 20+ API接口\n'
              '4层安全防护 · 三重数据备份 · 多级Fallback机制\n\n'
              '专为银发群体打造的桂林智慧导览与文旅电商平台\n'
              '让每一位长者都能安心、便捷地畅游桂林山水')
add_text_box(s, 0, Inches(3.2), W, Inches(2.5),
             thank_text, font_size=14, color=RGBColor(0xFF,0xFF,0xFF), alignment=PP_ALIGN.CENTER, line_spacing=1.6)
# bottom tags
end_tags = ['微信小程序', 'Vue 3', 'Express', 'TypeScript', 'SQL.js', 'Railway']
for i, tag in enumerate(end_tags):
    l = Inches(2.5) + i * Inches(1.55)
    bg = add_rounded_rect(s, l, Inches(5.8), Inches(1.3), Inches(0.4), None)
    bg.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    bg.line.fill.solid()
    bg.line.width = Pt(0.5)
    add_text_box(s, l, Inches(5.8), Inches(1.3), Inches(0.4),
                 tag, font_size=10, color=WHITE, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text_box(s, 0, Inches(6.5), W, Inches(0.5),
             '2026年5月 · 桂林', font_size=12, color=RGBColor(0xFF,0xFF,0xFF), alignment=PP_ALIGN.CENTER)

# ── Save ──
output_path = r'g:\github.cangku\yinfa\ppt\桂林银发旅游小程序_项目全景介绍.pptx'
prs.save(output_path)
print(f'PPT saved to: {output_path}')
print(f'Total slides: {len(prs.slides)}')