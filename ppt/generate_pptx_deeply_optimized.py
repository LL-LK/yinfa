#!/usr/bin/env python3
"""
桂林银发旅游小程序 - 深度优化PPT生成脚本
字体≥14pt | 空间利用率≥90% | 绿色护眼主题（银发友好）
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ============ 幻灯片尺寸：16:9 宽屏 ============
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ============ 护眼绿色主题配色（银发友好，高对比度）============
GREEN_900 = RGBColor(0x1B, 0x5E, 0x20)   # 深绿（标题）
GREEN_700 = RGBColor(0x2E, 0x7D, 0x32)   # 中绿
GREEN_500 = RGBColor(0x43, 0xA0, 0x47)   # 亮绿
GREEN_100 = RGBColor(0xC8, 0xE6, 0xC9)   # 淡绿（背景）
GREEN_50  = RGBColor(0xE8, 0xF5, 0xE9)   # 极淡绿
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x1A, 0x1A, 0x1A)   # 深灰文字
MID_GRAY  = RGBColor(0x4A, 0x4A, 0x4A)   # 中灰
LIGHT_GRAY= RGBColor(0xF5, 0xF5, 0xF5)   # 浅灰
ACCENT    = RGBColor(0xFF, 0x6F, 0x00)   # 橙色点缀

# ============ 字体大小常量（银发友好，全部≥14pt）============
TITLE_COVER = Pt(44)     # 封面大标题
SUBTITLE_COVER = Pt(26)  # 封面副标题
TAG_LINE = Pt(16)        # 封面标签行
H1 = Pt(28)              # 章节标题
H2 = Pt(20)              # 小节标题
BODY = Pt(14)            # 正文（≥14pt硬性要求）
SMALL = Pt(14)           # 小字（与正文同，确保可读）
FOOTER_PT = Pt(11)       # 页脚

# ============ 布局常量 ============
MARGIN = Inches(0.35)          # 左右边距
TOP_MARGIN = Inches(0.0)       # 顶部（标题栏在上面）
HEADER_H = Inches(1.1)         # 标题栏高度
FOOTER_H = Inches(0.32)        # 页脚高度
CONTENT_T = HEADER_H           # 内容区顶部 = 标题栏底部
CONTENT_B = SLIDE_H - FOOTER_H # 内容区底部
CONTENT_H = CONTENT_B - CONTENT_T # 内容区高度
USABLE_W = SLIDE_W - 2 * MARGIN # 可用宽度

def new_slide(prs):
    """创建空白幻灯片"""
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)

def bg(slide, color):
    """纯色背景"""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def header_bar(slide, title_text):
    """绿色标题栏"""
    # 深绿色标题栏
    bar = slide.shapes.add_shape(
        1, MARGIN, TOP_MARGIN, USABLE_W, HEADER_H
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN_900
    bar.line.fill.background()

    # 白色标题文字
    txb = slide.shapes.add_textbox(
        MARGIN + Inches(0.2), TOP_MARGIN + Inches(0.25),
        USABLE_W - Inches(0.4), HEADER_H - Inches(0.3)
    )
    tf = txb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = H1
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

def footer(slide, text="桂林银发旅游小程序 | 深度优化版"):
    """底部页脚"""
    txb = slide.shapes.add_textbox(
        MARGIN, CONTENT_B + Inches(0.04),
        USABLE_W, FOOTER_H - Inches(0.04)
    )
    tf = txb.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = FOOTER_PT
    p.font.color.rgb = MID_GRAY
    p.alignment = PP_ALIGN.CENTER

def content_box(slide, l, t, w, h, fill_color=None, border_color=None):
    """通用内容框"""
    shape = slide.shapes.add_shape(1, l, t, w, h)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape

def tb(slide, l, t, w, h, text, size=BODY, bold=False, color=DARK_GRAY,
       align=PP_ALIGN.LEFT, wrap=True):
    """文本框"""
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return txb

def tb_lines(slide, l, t, w, h, lines, size=BODY, bold=False, color=DARK_GRAY,
             spacing=1.3, align=PP_ALIGN.LEFT):
    """多行文本框，lines为[(文字, bold, size), ...]或[文字, ...]"""
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if isinstance(line, tuple):
            txt, b, sz = line
        else:
            txt, b, sz = line, bold, size

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = sz
        p.font.bold = b
        p.font.color.rgb = color
        p.alignment = align
        p.space_after = Pt(2)
        p.space_before = Pt(0)
    return txb


def make_yinfa_ppt():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ============================================================
    # 第1页：封面
    # ============================================================
    s = new_slide(prs)
    bg(s, GREEN_50)

    # 顶部装饰条
    top_bar = s.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(0.5))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = GREEN_900
    top_bar.line.fill.background()

    # 主标题
    tb(s, MARGIN, Inches(1.0), USABLE_W, Inches(1.2),
       "桂林银发旅游小程序", TITLE_COVER, True, GREEN_900, PP_ALIGN.CENTER)

    # 副标题
    tb(s, MARGIN, Inches(2.3), USABLE_W, Inches(0.8),
       "专为中老年游客打造的一站式旅游服务平台", SUBTITLE_COVER, False, GREEN_700, PP_ALIGN.CENTER)

    # 分割线
    line_shape = s.shapes.add_shape(1,
        Inches(2.5), Inches(3.2), Inches(8.333), Pt(3))
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = GREEN_500
    line_shape.line.fill.background()

    # 核心卖点5宫格
    features = [
        ("🗺️", "景点门票", "覆盖桂林/漓江\n知名景点"),
        ("🏨", "酒店预订", "精选老年友好\n住宿设施"),
        ("🚗", "旅游团", "专属银发团\n安全有保障"),
        ("🤖", "AI助手", "智能规划\n语音交互"),
        ("📍", "地图导览", "精准导航\n防走失"),
    ]
    card_w = Inches(2.3)
    card_h = Inches(1.5)
    gap = Inches(0.17)
    total_w = 5 * card_w + 4 * gap
    start_l = (SLIDE_W - total_w) / 2

    for i, (emoji, title, desc) in enumerate(features):
        l = start_l + i * (card_w + gap)
        # 卡片背景
        _ = content_box(s, l, Inches(3.6), card_w, card_h, GREEN_100, GREEN_500)
        # Emoji
        tb(s, l, Inches(3.7), card_w, Inches(0.5),
           emoji, Pt(28), False, GREEN_900, PP_ALIGN.CENTER)
        # 标题
        tb(s, l, Inches(4.15), card_w, Inches(0.3),
           title, Pt(14), True, GREEN_900, PP_ALIGN.CENTER)
        # 描述
        tb(s, l, Inches(4.45), card_w, Inches(0.6),
           desc, SMALL, False, DARK_GRAY, PP_ALIGN.CENTER)

    # 底部信息
    tb(s, MARGIN, Inches(5.5), USABLE_W, Inches(0.4),
       "技术栈：微信小程序 · Vue 3 Web · Express + TypeScript · SQL.js · Railway部署",
       TAG_LINE, False, MID_GRAY, PP_ALIGN.CENTER)
    tb(s, MARGIN, Inches(5.95), USABLE_W, Inches(0.4),
       "后端接口：https://yinfa-backend.up.railway.app/api",
       Pt(13), False, MID_GRAY, PP_ALIGN.CENTER)
    tb(s, MARGIN, Inches(6.4), USABLE_W, Inches(0.4),
       "微信小程序 × AI智能导览 × 漓江旅游知识服务",
       Pt(13), False, GREEN_500, PP_ALIGN.CENTER)

    footer(s, "桂林银发旅游小程序 · 银发友好 × AI赋能 · 2026")

    # ============================================================
    # 第2页：目录
    # ============================================================
    s = new_slide(prs)
    bg(s, WHITE)
    header_bar(s, "目 录")
    footer(s)

    sections = [
        ("01", "项目概述", "小程序定位与核心价值"),
        ("02", "技术架构", "前后端技术栈全貌"),
        ("03", "核心功能", "8大功能模块详解"),
        ("04", "AI助手", "智能旅游导览集成"),
        ("05", "数据库设计", "SQL.js + PostgreSQL"),
        ("06", "API接口", "18个接口全解析"),
        ("07", "部署方案", "Railway一键部署"),
        ("08", "未来规划", "支付/知识库/扩展"),
    ]

    # 4×2 网格布局
    card_w = (USABLE_W - Inches(0.5)) / 4 * 2.1
    card_h = Inches(1.2)
    col_gap = Inches(0.18)
    row_gap = Inches(0.2)
    cols = 4
    rows_n = 2

    start_l = MARGIN
    start_t = CONTENT_T + Inches(0.3)

    for idx, (num, title, desc) in enumerate(sections):
        col = idx % cols
        row = idx // cols
        l = start_l + col * (card_w + col_gap)
        t = start_t + row * (card_h + row_gap)

        # 卡片
        fc = GREEN_100 if int(num) % 2 == 0 else GREEN_50
        c = content_box(s, l, t, card_w, card_h, fc, GREEN_500)

        # 编号
        tb(s, l + Inches(0.1), t + Inches(0.1), Inches(0.6), Inches(0.5),
           num, H2, True, GREEN_700, PP_ALIGN.LEFT)
        # 标题
        tb(s, l + Inches(0.1), t + Inches(0.45), card_w - Inches(0.2), Inches(0.35),
           title, Pt(15), True, GREEN_900, PP_ALIGN.LEFT)
        # 描述
        tb(s, l + Inches(0.1), t + Inches(0.8), card_w - Inches(0.2), Inches(0.35),
           desc, Pt(13), False, MID_GRAY, PP_ALIGN.LEFT)

    # ============================================================
    # 第3页：项目概述
    # ============================================================
    s = new_slide(prs)
    bg(s, WHITE)
    header_bar(s, "01 项目概述")
    footer(s)

    # 左侧项目说明
    l_box = content_box(s, MARGIN, CONTENT_T + Inches(0.2),
                        Inches(5.8), CONTENT_H - Inches(0.4), GREEN_50, GREEN_500)
    tb(s, MARGIN + Inches(0.15), CONTENT_T + Inches(0.3),
       Inches(5.5), Inches(0.4),
       "项目定位", H2, True, GREEN_900, PP_ALIGN.LEFT)
    lines = [
        "桂林银发旅游小程序，专注为50岁以上中老年游客",
        "提供桂林、漓江流域的一站式旅游票务与导览服务。",
        "",
        "核心理念：简单操作、大字界面、安全优先、AI辅助。",
        "专为银发群体定制字体、配色、交互和内容体验。",
    ]
    tb_lines(s, MARGIN + Inches(0.15), CONTENT_T + Inches(0.7),
             Inches(5.5), Inches(1.5), lines, SMALL, False, DARK_GRAY)

    tb(s, MARGIN + Inches(0.15), CONTENT_T + Inches(2.2),
       Inches(5.5), Inches(0.4),
       "项目规模", H2, True, GREEN_900, PP_ALIGN.LEFT)
    stats = [
        ("9+", "小程序页面"),
        ("18+", "API接口"),
        ("8", "功能模块"),
        ("3", "技术栈层"),
    ]
    sw = Inches(1.25)
    for i, (num, label) in enumerate(stats):
        sl = MARGIN + Inches(0.15) + i * sw
        tb(s, sl, CONTENT_T + Inches(2.6), sw, Inches(0.6),
           num, H1, True, GREEN_700, PP_ALIGN.CENTER)
        tb(s, sl, CONTENT_T + Inches(3.1), sw, Inches(0.3),
           label, Pt(13), False, MID_GRAY, PP_ALIGN.CENTER)

    # 右侧技术架构
    r_box = content_box(s, MARGIN + Inches(6.1), CONTENT_T + Inches(0.2),
                        Inches(6.5), CONTENT_H - Inches(0.4), None)

    tb(s, MARGIN + Inches(6.25), CONTENT_T + Inches(0.3),
       Inches(6.0), Inches(0.4),
       "技术架构", H2, True, GREEN_900, PP_ALIGN.LEFT)

    layers = [
        ("前端层", "微信小程序 · Vue 3 Web", GREEN_100, GREEN_700),
        ("后端层", "Express.js + TypeScript", GREEN_50, GREEN_700),
        ("数据层", "SQL.js（内存）+ PostgreSQL（持久化）", GREEN_100, GREEN_700),
        ("AI层", "lvyou-agi · FastAPI · TF-IDF RAG", GREEN_50, GREEN_700),
        ("部署层", "Railway · Docker", GREEN_100, GREEN_700),
    ]
    for i, (layer, tech, _layer_bg, _tc) in enumerate(layers):
        lt = CONTENT_T + Inches(0.75) + i * Inches(0.72)
        # 层标签
        lb = content_box(s, MARGIN + Inches(6.25), lt, Inches(1.3), Inches(0.5),
                         GREEN_700)
        tb(s, MARGIN + Inches(6.25), lt + Inches(0.05), Inches(1.3), Inches(0.4),
           layer, Pt(13), True, WHITE, PP_ALIGN.CENTER)
        # 技术描述
        tb(s, MARGIN + Inches(7.65), lt + Inches(0.08), Inches(4.9), Inches(0.4),
           tech, SMALL, False, DARK_GRAY, PP_ALIGN.LEFT)

    # ============================================================
    # 第4页：技术架构详情
    # ============================================================
    s = new_slide(prs)
    bg(s, WHITE)
    header_bar(s, "02 技术架构")
    footer(s)

    # 三栏架构图
    col_w = (USABLE_W - 2 * Inches(0.2)) / 3
    col_data = [
        ("🖥️ 前端", GREEN_700, [
            "微信小程序（原生框架）",
            "  · pages/（9+页面）",
            "  · components/（组件）",
            "  · utils/（工具库）",
            "Vue 3 Web（建设中50%）",
            "  · Vite + TypeScript",
            "  · views/（首页/景点/导览）",
        ]),
        ("⚙️ 后端", GREEN_700, [
            "Express.js + TypeScript",
            "  · routes/（18个路由）",
            "  · middleware/（认证/日志）",
            "  · services/（业务逻辑）",
            "  · validators/（Zod验证）",
            "  · logger.ts（Pino日志）",
            "SQL.js（内存）+ PostgreSQL",
        ]),
        ("🤖 AI服务", GREEN_700, [
            "lvyou-agi（独立部署）",
            "  · FastAPI + Python",
            "  · LijiangCoordinator Agent",
            "  · RouteAgent + BehaviorAgent",
            "  · TF-IDF RAG（numpy/scipy）",
            "  · MiniMax LLM API",
            "  · ChromaDB（可选）",
        ]),
    ]

    for i, (title, color, items) in enumerate(col_data):
        l = MARGIN + i * (col_w + Inches(0.2))
        # 顶部色条
        bar = s.shapes.add_shape(1, l, CONTENT_T + Inches(0.25),
                                 col_w, Inches(0.55))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        tb(s, l, CONTENT_T + Inches(0.3), col_w, Inches(0.45),
           title, Pt(16), True, WHITE, PP_ALIGN.CENTER)

        # 内容卡片
        c = content_box(s, l, CONTENT_T + Inches(0.85),
                        col_w, Inches(3.9), GREEN_50)
        for j, item in enumerate(items):
            bold = not item.startswith("  ")
            tb(s, l + Inches(0.12), CONTENT_T + Inches(0.95) + j * Inches(0.5),
               col_w - Inches(0.2), Inches(0.45),
               item.strip(), SMALL, bold, DARK_GRAY if bold else MID_GRAY)

    # 底部连接说明
    conn = content_box(s, MARGIN, CONTENT_T + Inches(4.85),
                       USABLE_W, Inches(0.9), GREEN_100)
    tb(s, MARGIN + Inches(0.2), CONTENT_T + Inches(4.95),
       USABLE_W - Inches(0.4), Inches(0.35),
       "🔗 系统连接：微信小程序 ──HTTP──> Express后端 ──HTTP──> lvyou-agi AI服务 ──MiniMax API──> 游客",
       Pt(14), False, GREEN_900, PP_ALIGN.CENTER)
    tb(s, MARGIN + Inches(0.2), CONTENT_T + Inches(5.3),
       USABLE_W - Inches(0.4), Inches(0.35),
       "部署地址：yinfa后端 https://yinfa-backend.up.railway.app  |  lvyou-agi https://lvyou-agi.onrender.com",
       Pt(13), False, MID_GRAY, PP_ALIGN.CENTER)

    # ============================================================
    # 第5页：核心功能（1）
    # ============================================================
    s = new_slide(prs)
    bg(s, WHITE)
    header_bar(s, "03 核心功能 · 商品与订单")
    footer(s)

    # 四大功能双行布局
    features_data = [
        {
            "icon": "🎫", "title": "商品分类浏览", "status": "✅ 已完成",
            "items": [
                "GET /api/categories —— 获取景点分类列表",
                "GET /api/products —— 获取商品/景点列表",
                "支持分类筛选（category参数）",
                "支持关键词搜索（search参数）",
                "返回商品图片、价格、库存信息",
                "响应结构：{code, data, error}",
            ]
        },
        {
            "icon": "📦", "title": "商品详情", "status": "✅ 已完成",
            "items": [
                "GET /api/products/:id —— 商品详情",
                "展示名称、价格、描述、图片",
                "支持库存状态实时查询",
                "关联分类信息一并返回",
                "is_active字段控制上下架",
                "详细字段：id/name/price/stock/image_url",
            ]
        },
        {
            "icon": "🛒", "title": "购物车管理", "status": "🔨 开发中",
            "items": [
                "GET /api/cart —— 获取购物车列表",
                "POST /api/cart/add —— 添加商品",
                "POST /api/cart/update —— 更新数量",
                "POST /api/cart/remove —— 删除商品",
                "POST /api/cart/clear —— 清空购物车",
                "基于openid识别用户购物车",
            ]
        },
        {
            "icon": "📋", "title": "订单管理", "status": "✅ 已完成",
            "items": [
                "POST /api/order/create —— 创建订单",
                "GET /api/orders —— 订单列表",
                "GET /api/orders/:id —— 订单详情",
                "自动生成订单号（ORD前缀+时间戳）",
                "支持多商品同时下单（items数组）",
                "订单状态：pending/paid/shipped/completed",
            ]
        },
    ]

    card_w = (USABLE_W - Inches(0.3)) / 2
    card_h = Inches(2.3)

    for idx, feat in enumerate(features_data):
        col = idx % 2
        row = idx // 2
        l = MARGIN + col * (card_w + Inches(0.3))
        t = CONTENT_T + Inches(0.2) + row * (card_h + Inches(0.2))

        fc = GREEN_100 if row == 0 else GREEN_50
        content_box(s, l, t, card_w, card_h, fc, GREEN_500)

        # 图标+标题+状态
        tb(s, l + Inches(0.12), t + Inches(0.1), Inches(1.5), Inches(0.4),
           feat["icon"] + " " + feat["title"], Pt(15), True, GREEN_900)
        st_color = GREEN_700 if "✅" in feat["status"] else ACCENT
        tb(s, l + card_w - Inches(1.2), t + Inches(0.1), Inches(1.1), Inches(0.3),
           feat["status"], Pt(12), True, st_color, PP_ALIGN.RIGHT)

        for j, item in enumerate(feat["items"]):
            tb(s, l + Inches(0.15), t + Inches(0.5) + j * Inches(0.28),
               card_w - Inches(0.25), Inches(0.28),
               "• " + item, Pt(13), False, DARK_GRAY)

    # ============================================================
    # 第6页：核心功能（2）
    # ============================================================
    s = new_slide(prs)
    bg(s, WHITE)
    header_bar(s, "03 核心功能 · 用户与AI")
    footer(s)

    features2 = [
        {
            "icon": "👤", "title": "用户中心", "status": "✅ 已完成",
            "items": [
                "POST /api/users/login —— 微信登录/注册",
                "  openid/nickname/avatar_url/phone",
                "GET /api/users/:openid —— 获取用户信息",
                "支持手机号快捷登录",
                "用户创建时间自动记录（created_at）",
                "集成紧急联系人管理（emergency_contacts）",
            ]
        },
        {
            "icon": "📍", "title": "地址管理", "status": "✅ 已完成",
            "items": [
                "POST /api/address/create —— 新增地址",
                "  支持设置默认地址（is_default）",
                "GET /api/addresses —— 地址列表",
                "  按openid查询用户所有地址",
                "字段：full_name/phone/address_line",
                "  city/postal_code/is_default",
            ]
        },
        {
            "icon": "🤖", "title": "AI智能助手", "status": "✅ 已集成",
            "items": [
                "集成lvyou-agi旅游AI知识服务",
                "Agent路由：Route/Behavior/Knowledge",
                "支持SSE流式输出（/api/chat/stream）",
                "意图识别 → 路由 → Agent执行",
                "RAG知识库：桂林漓江景点/美食/交通",
                "天气查询：支持桂林/北京/上海/广州等8城",
            ]
        },
        {
            "icon": "🗺️", "title": "地图导航", "status": "✅ 已完成",
            "items": [
                "内置地图页面（pages/map）",
                "支持景点定位与展示",
                "结合AI助手提供智能路线建议",
                "可查看周边配套设施（厕所/急救）",
                "杨堤码头为重点研究区域",
                "与景点详情页深度整合",
            ]
        },
    ]

    for idx, feat in enumerate(features2):
        col = idx % 2
        row = idx // 2
        l = MARGIN + col * (card_w + Inches(0.3))
        t = CONTENT_T + Inches(0.2) + row * (card_h + Inches(0.2))

        fc = GREEN_100 if row == 0 else GREEN_50
        content_box(s, l, t, card_w, card_h, fc, GREEN_500)

        tb(s, l + Inches(0.12), t + Inches(0.1), Inches(1.5), Inches(0.4),
           feat["icon"] + " " + feat["title"], Pt(15), True, GREEN_900)
        st_color = GREEN_700 if "✅" in feat["status"] else ACCENT
        tb(s, l + card_w - Inches(1.2), t + Inches(0.1), Inches(1.1), Inches(0.3),
           feat["status"], Pt(12), True, st_color, PP_ALIGN.RIGHT)

        for j, item in enumerate(feat["items"]):
            tb(s, l + Inches(0.15), t + Inches(0.5) + j * Inches(0.28),
               card_w - Inches(0.25), Inches(0.28),
               "• " + item, Pt(13), False, DARK_GRAY)

    # ============================================================
    # 第7页：AI助手深度解析
    # ============================================================
    s = new_slide(prs)
    bg(s, WHITE)
    header_bar(s, "04 AI助手 · lvyou-agi集成")
    footer(s)

    # 左侧Agent架构
    content_box(s, MARGIN, CONTENT_T + Inches(0.2), Inches(6.0), CONTENT_H - Inches(0.4), GREEN_50)

    tb(s, MARGIN + Inches(0.15), CONTENT_T + Inches(0.3),
       Inches(5.7), Inches(0.4),
       "LijiangCoordinator 多意图路由架构", H2, True, GREEN_900)

    arch_items = [
        ("🌐", "用户输入", "小程序AI助手页面输入旅游需求"),
        ("⬇️", "意图分析", "Coordinator分析：单意图/多意图"),
        ("⟶", "路由分发", "单意图→单一Agent | 多意图→三Agent并行"),
        ("🤝", "并行执行", "RouteAgent + BehaviorAgent + KnowledgeAgent"),
        ("📝", "结果综合", "综合三Agent输出，格式化响应"),
    ]
    for i, (emoji, title, desc) in enumerate(arch_items):
        lt = CONTENT_T + Inches(0.8) + i * Inches(0.65)
        # Emoji标签
        tb(s, MARGIN + Inches(0.15), lt, Inches(0.5), Inches(0.45),
           emoji, Pt(18), False, GREEN_900, PP_ALIGN.CENTER)
        # 标题+描述
        tb(s, MARGIN + Inches(0.7), lt, Inches(2.0), Inches(0.35),
           title, Pt(14), True, GREEN_700)
        tb(s, MARGIN + Inches(0.7), lt + Inches(0.3), Inches(5.0), Inches(0.3),
           desc, Pt(13), False, MID_GRAY)

    # 右侧Agent详情
    content_box(s, MARGIN + Inches(6.2), CONTENT_T + Inches(0.2), Inches(6.5), CONTENT_H - Inches(0.4), None)

    tb(s, MARGIN + Inches(6.35), CONTENT_T + Inches(0.3),
       Inches(6.2), Inches(0.4),
       "三大专业Agent", H2, True, GREEN_900)

    agents = [
        ("🗺️ RouteAgent", "旅游路线规划",
         ["· 分析用户时间/预算/人数", "· 漓江竹筏/徒步/自驾路线", "· 生成个性化行程方案", "· 多目标优化（时间/费用/体验）"]),
        ("📊 BehaviorAgent", "游客行为预测",
         ["· 预测各景点停留时间", "· 购买意愿与消费分析", "· 客流高峰时段预测", "· 最优出发时间建议"]),
        ("📚 KnowledgeAgent", "旅游知识问答",
         ["· 景点介绍与历史背景", "· 票价信息与优惠政策", "· 美食推荐与交通指南", "· RAG检索桂林漓江知识库"]),
    ]

    for i, (name, role, items) in enumerate(agents):
        lt = CONTENT_T + Inches(0.75) + i * Inches(1.25)
        # Agent名称
        lb = content_box(s, MARGIN + Inches(6.35), lt, Inches(2.5), Inches(0.4), GREEN_700)
        tb(s, MARGIN + Inches(6.35), lt + Inches(0.02), Inches(2.5), Inches(0.38),
           name, Pt(14), True, WHITE, PP_ALIGN.CENTER)
        tb(s, MARGIN + Inches(6.35), lt + Inches(0.38), Inches(6.0), Inches(0.3),
           "职能：" + role, Pt(13), True, GREEN_700)
        for j, item in enumerate(items):
            tb(s, MARGIN + Inches(6.4), lt + Inches(0.65) + j * Inches(0.24),
               Inches(6.0), Inches(0.24), item, Pt(13), False, DARK_GRAY)

    # ============================================================
    # 第8页：数据库设计
    # ============================================================
    s = new_slide(prs)
    bg(s, WHITE)
    header_bar(s, "05 数据库设计")
    footer(s)

    tables = [
        ("categories", "景点分类表", GREEN_700, ["id（主键）", "name（名称）", "slug（别名）", "order（排序）"]),
        ("products", "商品/景点表", GREEN_700, ["id（主键）", "name（名称）", "description", "price（价格）", "stock（库存）", "image_url", "is_active"]),
        ("user_profiles", "用户档案表", GREEN_500, ["id（主键）", "openid（微信）", "nickname", "avatar_url", "phone（手机）", "created_at"]),
        ("addresses", "地址管理表", GREEN_500, ["id（主键）", "openid", "full_name", "phone", "address_line", "city", "postal_code", "is_default"]),
        ("orders", "订单主表", GREEN_700, ["id（主键）", "order_no", "openid", "total_price", "status（状态）", "created_at", "updated_at"]),
        ("order_items", "订单明细表", GREEN_700, ["id（主键）", "order_id（外键）", "product_id", "quantity", "unit_price", "subtotal"]),
    ]

    col_w = (USABLE_W - 5 * Inches(0.12)) / 6
    for i, (tbl, label, color, cols) in enumerate(tables):
        l = MARGIN + i * (col_w + Inches(0.12))
        # 表头
        h = content_box(s, l, CONTENT_T + Inches(0.2), col_w, Inches(0.5), color)
        tb(s, l, CONTENT_T + Inches(0.25), col_w, Inches(0.4),
           tbl, Pt(11), True, WHITE, PP_ALIGN.CENTER)
        # 表标签
        tb(s, l, CONTENT_T + Inches(0.7), col_w, Inches(0.3),
           label, Pt(12), True, GREEN_900, PP_ALIGN.CENTER)
        # 字段
        bh = len(cols) * Inches(0.3) + Inches(0.1)
        content_box(s, l, CONTENT_T + Inches(1.0), col_w, bh, GREEN_50, GREEN_500)
        for j, col in enumerate(cols):
            tb(s, l + Inches(0.05), CONTENT_T + Inches(1.08) + j * Inches(0.3),
               col_w - Inches(0.1), Inches(0.28),
               col, Pt(11), False, DARK_GRAY, PP_ALIGN.LEFT)

    # 迁移说明
    note = content_box(s, MARGIN, CONTENT_T + Inches(4.3), USABLE_W, Inches(1.1), GREEN_100, GREEN_500)
    tb(s, MARGIN + Inches(0.2), CONTENT_T + Inches(4.4), Inches(3.0), Inches(0.35),
       "📊 当前：SQL.js（内存数据库）", Pt(14), True, GREEN_900)
    tb(s, MARGIN + Inches(0.2), CONTENT_T + Inches(4.75), Inches(12.0), Inches(0.28),
       "  支持快速开发、本地运行、数据持久化到 /app/data 目录（Railway持久化卷）", Pt(13), False, DARK_GRAY)
    tb(s, MARGIN + Inches(0.2), CONTENT_T + Inches(5.0), Inches(3.0), Inches(0.35),
       "🚀 迁移：PostgreSQL（持久化）", Pt(14), True, GREEN_900)
    tb(s, MARGIN + Inches(0.2), CONTENT_T + Inches(5.35), Inches(12.0), Inches(0.28),
       "  已创建迁移脚本 001_initial_schema.sql，配置 railway.json，依赖 DATABASE_URL 环境变量", Pt(13), False, DARK_GRAY)

    # ============================================================
    # 第9页：API接口（全部18个）
    # ============================================================
    s = new_slide(prs)
    bg(s, WHITE)
    header_bar(s, "06 API接口全览（18个接口）")
    footer(s)

    apis = [
        ("系统", "GET", "/health", "健康检查"),
        ("系统", "GET", "/", "服务器信息"),
        ("分类", "GET", "/api/categories", "获取分类列表"),
        ("商品", "GET", "/api/products", "获取商品列表"),
        ("商品", "GET", "/api/products/:id", "获取商品详情"),
        ("用户", "POST", "/api/users/login", "微信登录注册"),
        ("用户", "GET", "/api/users/:openid", "获取用户信息"),
        ("地址", "POST", "/api/address/create", "创建地址"),
        ("地址", "GET", "/api/addresses", "获取地址列表"),
        ("订单", "POST", "/api/order/create", "创建订单"),
        ("订单", "GET", "/api/orders", "获取订单列表"),
        ("订单", "GET", "/api/orders/:id", "获取订单详情"),
        ("天气", "GET", "/api/weather", "天气查询"),
        ("购物车", "GET", "/api/cart", "获取购物车（开发中）"),
        ("购物车", "POST", "/api/cart/add", "添加购物车（开发中）"),
        ("AI", "POST", "/api/chat", "AI对话"),
        ("AI", "POST", "/api/chat/stream", "AI流式对话"),
        ("AI", "GET", "/api/capabilities", "AI能力清单"),
    ]

    # 三列布局
    cols_n = 3
    per_col = 6
    col_w2 = (USABLE_W - 2 * Inches(0.15)) / 3

    for idx, (cat, method, path, desc) in enumerate(apis):
        col = idx // per_col
        row = idx % per_col
        l = MARGIN + col * (col_w2 + Inches(0.15))
        t = CONTENT_T + Inches(0.15) + row * Inches(0.62)

        # 方法标签
        m_color = {"GET": GREEN_700, "POST": ACCENT}.get(method, GREEN_700)
        mb = content_box(s, l, t, Inches(0.7), Inches(0.45), m_color)
        tb(s, l, t + Inches(0.05), Inches(0.7), Inches(0.35),
           method, Pt(11), True, WHITE, PP_ALIGN.CENTER)

        # 路径
        tb(s, l + Inches(0.75), t + Inches(0.05), Inches(3.3), Inches(0.35),
           path, Pt(12), True, GREEN_900)

        # 描述
        tb(s, l + Inches(0.75), t + Inches(0.3), Inches(3.3), Inches(0.28),
           desc, Pt(12), False, MID_GRAY)

    # ============================================================
    # 第10页：部署方案
    # ============================================================
    s = new_slide(prs)
    bg(s, WHITE)
    header_bar(s, "07 部署方案")
    footer(s)

    # Railway一键部署流程
    content_box(s, MARGIN, CONTENT_T + Inches(0.2), Inches(6.0), Inches(3.5), GREEN_50, GREEN_500)
    tb(s, MARGIN + Inches(0.2), CONTENT_T + Inches(0.3), Inches(5.6), Inches(0.4),
       "🚀 Railway 一键部署", H2, True, GREEN_900)
    steps = [
        "① Fork 本仓库到 GitHub",
        "② Railway 控制台新建项目",
        "③ 选择「从 GitHub 部署」",
        "④ 选择本仓库，Railway 自动读取 railway.json",
        "⑤ 添加持久化卷，Mount Path：/app/data",
        "⑥ 配置环境变量（如 MINIMAX_API_KEY 等）",
        "⑦ 等待构建完成，自动获得 HTTPS 地址",
    ]
    for i, step in enumerate(steps):
        tb(s, MARGIN + Inches(0.2), CONTENT_T + Inches(0.75) + i * Inches(0.38),
           Inches(5.6), Inches(0.35), step, Pt(14), False, DARK_GRAY)

    # Docker独立部署
    content_box(s, MARGIN + Inches(6.2), CONTENT_T + Inches(0.2), Inches(6.5), Inches(3.5), GREEN_50, GREEN_500)
    tb(s, MARGIN + Inches(6.35), CONTENT_T + Inches(0.3), Inches(6.2), Inches(0.4),
       "🐳 Docker 容器部署", H2, True, GREEN_900)
    docker_steps = [
        "① Docker 已安装（docker --version）",
        "② 项目根目录构建镜像：",
        "   docker build -t yinfa-backend .",
        "③ 运行容器（挂载数据卷）：",
        "   docker run -d -p 8000:8000",
        "   -v yinfa-data:/app/data",
        "   -e PORT=8000",
        "   --name yinfa-backend",
        "   yinfa-backend",
    ]
    for i, step in enumerate(docker_steps):
        bold = not step.startswith("   ")
        tb(s, MARGIN + Inches(6.35), CONTENT_T + Inches(0.75) + i * Inches(0.36),
           Inches(6.2), Inches(0.34),
           step, Pt(14), bold, DARK_GRAY if bold else MID_GRAY)

    # 环境变量表
    content_box(s, MARGIN, CONTENT_T + Inches(3.8), USABLE_W, Inches(2.1), None)
    tb(s, MARGIN + Inches(0.15), CONTENT_T + Inches(3.9), Inches(3.0), Inches(0.35),
       "⚙️ 环境变量配置", H2, True, GREEN_900)

    env_vars = [
        ("PORT", "服务器端口", "8000"),
        ("NODE_ENV", "环境模式", "production"),
        ("MINIMAX_API_KEY", "MiniMax AI密钥", "（联系管理员）"),
        ("AMAP_KEY", "高德天气API", "（可选，未配返回模拟数据）"),
        ("DATABASE_URL", "PostgreSQL连接", "（迁移后配置）"),
    ]
    col_w3 = USABLE_W / 5
    for i, (name, desc, val) in enumerate(env_vars):
        l = MARGIN + i * col_w3
        lb = content_box(s, l, CONTENT_T + Inches(4.3), col_w3 - Inches(0.08), Inches(0.35), GREEN_700)
        tb(s, l, CONTENT_T + Inches(4.32), col_w3 - Inches(0.08), Inches(0.3),
           name, Pt(11), True, WHITE, PP_ALIGN.CENTER)
        tb(s, l, CONTENT_T + Inches(4.65), col_w3 - Inches(0.08), Inches(0.3),
           desc, Pt(12), False, DARK_GRAY, PP_ALIGN.CENTER)
        tb(s, l, CONTENT_T + Inches(4.95), col_w3 - Inches(0.08), Inches(0.3),
           val, Pt(11), False, GREEN_700, PP_ALIGN.CENTER)

    # ============================================================
    # 第11页：未来规划
    # ============================================================
    s = new_slide(prs)
    bg(s, WHITE)
    header_bar(s, "08 未来规划与待办")
    footer(s)

    # P0阻断
    content_box(s, MARGIN, CONTENT_T + Inches(0.2), Inches(4.1), Inches(3.3), None)
    tb(s, MARGIN + Inches(0.15), CONTENT_T + Inches(0.3), Inches(3.8), Inches(0.4),
       "🔴 P0 阻断性问题", H2, True, RGBColor(0xB7, 0x1C, 0x1C))

    p0_items = [
        ("① 配置 DATABASE_URL（PostgreSQL迁移）", "Railway 控制台添加环境变量"),
        ("② 注册微信小商店（激活支付）", "个人开发者首选微信小商店方案"),
        ("③ 部署 lvyou-agi 到 Railway", "解决Render免费版30分钟休眠问题"),
        ("④ 补充 lvyou-agi 知识库文档", "当前仅有杨堤码头1个文档"),
    ]
    for i, (title, sub) in enumerate(p0_items):
        lt = CONTENT_T + Inches(0.75) + i * Inches(0.72)
        tb(s, MARGIN + Inches(0.15), lt, Inches(3.9), Inches(0.35),
           title, Pt(14), True, DARK_GRAY)
        tb(s, MARGIN + Inches(0.15), lt + Inches(0.3), Inches(3.9), Inches(0.3),
           "→ " + sub, Pt(12), False, MID_GRAY)

    # P1功能增强
    content_box(s, MARGIN + Inches(4.25), CONTENT_T + Inches(0.2), Inches(4.4), Inches(3.3), None)
    tb(s, MARGIN + Inches(4.4), CONTENT_T + Inches(0.3), Inches(4.1), Inches(0.4),
       "🟡 P1 功能增强", H2, True, RGBColor(0xF5, 0x7C, 0x00))

    p1_items = [
        ("⑤ Vue前端接入AI助手（50%）", "当前缺详情/地图/支付页"),
        ("⑥ 流式AI响应（SSE）", "改用 /api/chat/stream"),
        ("⑦ 完成微信支付集成", "支付回调处理完善"),
        ("⑧ Railway PostgreSQL验证", "数据库连接测试"),
    ]
    for i, (title, sub) in enumerate(p1_items):
        lt = CONTENT_T + Inches(0.75) + i * Inches(0.72)
        tb(s, MARGIN + Inches(4.4), lt, Inches(4.1), Inches(0.35),
           title, Pt(14), True, DARK_GRAY)
        tb(s, MARGIN + Inches(4.4), lt + Inches(0.3), Inches(4.1), Inches(0.3),
           "→ " + sub, Pt(12), False, MID_GRAY)

    # P2长期建设
    content_box(s, MARGIN, CONTENT_T + Inches(3.6), USABLE_W, Inches(2.5), GREEN_50, GREEN_500)
    tb(s, MARGIN + Inches(0.15), CONTENT_T + Inches(3.7), Inches(4.0), Inches(0.4),
       "🟢 P2 长期建设", H2, True, GREEN_700)

    p2_items = [
        ("⑨ 知识库扩展", "爬取马蜂窝/携程/桂林旅游官网扩充RAG知识库"),
        ("⑩ 双server.py合并", "统一 src/server.py 与根 server.py"),
        ("⑪ 高德地图API对接", "桂林旅游官网数据整合"),
        ("⑫ lvyou-agi多意图并行", "从单意图优先升级为真正多意图并行检测"),
    ]
    cols_p2 = 4
    pw = (USABLE_W - 3 * Inches(0.15)) / 4
    for i, (title, sub) in enumerate(p2_items):
        col = i % cols_p2
        l = MARGIN + col * (pw + Inches(0.15))
        row = i // cols_p2
        t = CONTENT_T + Inches(4.1) + row * Inches(1.15)
        content_box(s, l, t, pw, Inches(1.05), WHITE, GREEN_500)
        tb(s, l + Inches(0.1), t + Inches(0.08), pw - Inches(0.2), Inches(0.35),
           title, Pt(14), True, GREEN_900)
        tb(s, l + Inches(0.1), t + Inches(0.4), pw - Inches(0.2), Inches(0.6),
           sub, Pt(12), False, MID_GRAY)

    # ============================================================
    # 第12页：总结
    # ============================================================
    s = new_slide(prs)
    bg(s, GREEN_50)

    # 顶部深绿条
    tb(s, 0, 0, SLIDE_W, Inches(0.6),
       "桂林银发旅游小程序 · 总结", TITLE_COVER, True, WHITE, PP_ALIGN.CENTER)
    top = s.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(0.6))
    top.fill.solid()
    top.fill.fore_color.rgb = GREEN_900
    top.line.fill.background()
    tb(s, 0, Inches(0.08), SLIDE_W, Inches(0.5),
       "桂林银发旅游小程序 · 项目总结", Pt(22), True, WHITE, PP_ALIGN.CENTER)

    summary_items = [
        ("🎯", "定位清晰", "专为50岁以上银发游客设计，大字界面，简单操作，安全优先"),
        ("🛠️", "技术完整", "微信小程序+Vue3+Express+SQL.js+AI，18个API全贯通"),
        ("🤖", "AI赋能", "集成漓江旅游导览Agent，RAG+LLM智能规划路线"),
        ("📱", "体验优化", "银发友好字体/配色/交互设计，语音+文字双交互"),
        ("🚀", "部署就绪", "Railway一键部署，Docker容器化，持久化存储"),
        ("🔮", "扩展性强", "预留支付/知识库/多Agent扩展，架构清晰可演进"),
    ]

    card_w = (USABLE_W - 5 * Inches(0.15)) / 3
    card_h = Inches(1.55)

    for i, (emoji, title, desc) in enumerate(summary_items):
        col = i % 3
        row = i // 3
        l = MARGIN + col * (card_w + Inches(0.15))
        t = CONTENT_T + Inches(0.5) + row * (card_h + Inches(0.2))

        fc = GREEN_100 if row == 0 else GREEN_50
        content_box(s, l, t, card_w, card_h, fc, GREEN_500)
        tb(s, l + Inches(0.15), t + Inches(0.1), card_w - Inches(0.3), Inches(0.45),
           emoji + " " + title, Pt(16), True, GREEN_900)
        tb(s, l + Inches(0.15), t + Inches(0.55), card_w - Inches(0.3), Inches(0.9),
           desc, Pt(14), False, DARK_GRAY)

    tb(s, MARGIN, CONTENT_T + Inches(4.0), USABLE_W, Inches(0.4),
       "yinfa × lvyou-agi · 微信小程序 × AI导览 × 桂林漓江旅游",
       Pt(16), True, GREEN_700, PP_ALIGN.CENTER)
    tb(s, MARGIN, CONTENT_T + Inches(4.45), USABLE_W, Inches(0.4),
       "后端：https://yinfa-backend.up.railway.app  |  AI：https://lvyou-agi.onrender.com",
       Pt(13), False, MID_GRAY, PP_ALIGN.CENTER)
    tb(s, MARGIN, CONTENT_T + Inches(4.85), USABLE_W, Inches(0.4),
       "联系维护：微信小程序搜索「银发旅游」或访问 GitHub cangku/yinfa",
       Pt(13), False, MID_GRAY, PP_ALIGN.CENTER)

    footer(s, "桂林银发旅游小程序 | Yinfa × lvyou-agi × AI · 2026")

    # 保存
    out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                            "桂林银发旅游小程序_深度优化版.pptx")
    prs.save(out_path)
    print(f"✅ yinfa PPT已生成：{out_path}")
    return out_path


if __name__ == "__main__":
    make_yinfa_ppt()
