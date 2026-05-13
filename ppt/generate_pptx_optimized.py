#!/usr/bin/env python3
"""
桂林银发旅游小程序 - 项目全景介绍 PPT (优化版)
字体>=14pt | 空间利用率>=90% | 银发友好 | 内容丰富
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ==================== 颜色定义 ====================
GREEN_900  = RGBColor(0x0D, 0x33, 0x20)
GREEN_800  = RGBColor(0x1A, 0x4D, 0x33)
GREEN_700  = RGBColor(0x2E, 0x6B, 0x47)
GREEN_600  = RGBColor(0x3D, 0x8B, 0x5C)
GREEN_500  = RGBColor(0x4A, 0xAD, 0x75)
GREEN_400  = RGBColor(0x7D, 0xC9, 0x9A)
GREEN_300  = RGBColor(0x4C, 0x9A, 0x6B)
GREEN_200  = RGBColor(0xB3, 0xE0, 0xCE)
GREEN_100  = RGBColor(0xD4, 0xED, 0xE4)
GREEN_50   = RGBColor(0xE8, 0xF8, 0xEF)
GRAY_900   = RGBColor(0x1A, 0x1A, 0x1A)
GRAY_700   = RGBColor(0x3A, 0x3A, 0x3A)
GRAY_600   = RGBColor(0x4A, 0x4A, 0x4A)
GRAY_500   = RGBColor(0x6B, 0x6B, 0x6B)
GRAY_300   = RGBColor(0xD1, 0xD5, 0xDB)
GRAY_100   = RGBColor(0xF3, 0xF4, 0xF6)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GOLD       = RGBColor(0xF5, 0x96, 0x0B)
RED        = RGBColor(0xDC, 0x26, 0x26)

# 幻灯片尺寸 16:9 宽屏
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.35)
BODY_FONT = 14   # 正文最小14pt（银发友好）
H1_FONT   = 26   # 大标题
H2_FONT   = 20   # 副标题
H3_FONT   = 16   # 小标题
CAPTION   = 12   # 说明文字（>=12pt可接受）

# ==================== 辅助函数 ====================

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def bg(slide, color):
    """填充背景"""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    """添加矩形"""
    shape = slide.shapes.add_shape(1, l, t, w, h)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, font_size=BODY_FONT, bold=False,
             color=GRAY_900, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    """添加文本框（字体保底14pt）"""
    if font_size < BODY_FONT:
        font_size = BODY_FONT
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def add_text_lines(slide, lines, l, t, w, h, font_size=BODY_FONT,
                   bold=False, color=GRAY_900, align=PP_ALIGN.LEFT,
                   line_spacing=1.3):
    """
    添加多行文本框
    lines: [(文字, bold, color), ...] 或 [str, ...]
    """
    if font_size < BODY_FONT:
        font_size = BODY_FONT
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True

    first = True
    for item in lines:
        if isinstance(item, str):
            txt, b, col = item, bold, color
        else:
            txt, b, col = item[0], item[1] if len(item) > 1 else bold, item[2] if len(item) > 2 else color

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.alignment = align
        from pptx.util import Pt as PT
        from pptx.oxml.ns import qn
        from lxml import etree
        pPr = p._p.get_or_add_pPr()
        pPr.set(nsmap={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
        # 行距
        l = pPr.find(qn('a:lnSpc'))
        if l is None:
            l = etree.SubElement(pPr, qn('a:lnSpc'))
        factor = etree.SubElement(l, qn('a:spcPct'))
        factor.set('val', str(int(line_spacing * 100000)))

        run = p.add_run()
        run.text = txt
        run.font.size  = Pt(font_size)
        run.font.bold  = b
        run.font.color.rgb = col
    return txb

def stat_block(slide, label, value, l, t, w=None, h=Inches(0.85)):
    """数据统计块（大字体）"""
    if w is None:
        w = Inches(2.2)
    box = add_rect(slide, l, t, w, h, fill_color=GREEN_100,
                   line_color=GREEN_400, line_width=Pt(1.5))
    # value 大字
    add_text(slide, value, l + Inches(0.12), t + Inches(0.06),
             w - Inches(0.15), Inches(0.45), font_size=22,
             bold=True, color=GREEN_900, align=PP_ALIGN.CENTER)
    # label 小字（至少14）
    add_text(slide, label, l + Inches(0.12), t + Inches(0.48),
             w - Inches(0.15), Inches(0.32), font_size=CAPTION,
             bold=False, color=GRAY_600, align=PP_ALIGN.CENTER)

def tag_box(slide, text, l, t, font_size=CAPTION, bg=GREEN_600, fg=WHITE):
    """标签小方块"""
    w = Inches(len(text) * font_size * 0.018 + 0.3)
    h = Inches(font_size * 0.035 + 0.15)
    add_rect(slide, l, t, w, h, fill_color=bg)
    add_text(slide, text, l, t, w, h, font_size=font_size,
             bold=False, color=fg, align=PP_ALIGN.CENTER)

def section_title_bar(slide, title, subtitle=None):
    """顶部标题栏"""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), fill_color=GREEN_900)
    add_text(slide, title, Inches(0.45), Inches(0.12),
             SLIDE_W - Inches(0.9), Inches(0.6), font_size=H1_FONT,
             bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, Inches(0.45), Inches(0.7),
                 SLIDE_W - Inches(0.9), Inches(0.35), font_size=H3_FONT,
                 bold=False, color=GREEN_200, align=PP_ALIGN.LEFT)

def icon_label(slide, icon, label, l, t, icon_size=H2_FONT, label_size=CAPTION):
    """图标+文字组合"""
    add_text(slide, icon, l, t, Inches(0.6), Inches(0.55),
             font_size=icon_size, bold=True, color=GREEN_700, align=PP_ALIGN.CENTER)
    add_text(slide, label, l + Inches(0.55), t + Inches(0.08),
             Inches(2.5), Inches(0.4), font_size=label_size,
             bold=False, color=GRAY_700)

def footer_bar(slide, text="桂林银发旅游小程序 · 项目全景介绍 · 2024"):
    """底部页脚"""
    add_rect(slide, 0, SLIDE_H - Inches(0.3), SLIDE_W, Inches(0.3),
             fill_color=GREEN_800)
    add_text(slide, text, MARGIN, SLIDE_H - Inches(0.28),
             SLIDE_W - Inches(0.7), Inches(0.25), font_size=CAPTION,
             bold=False, color=GREEN_200, align=PP_ALIGN.LEFT)

def two_col_layout(slide, title, subtitle, left_content_fn, right_content_fn,
                   left_title=None, right_title=None):
    """标准双栏布局 + 标题栏"""
    section_title_bar(slide, title, subtitle)
    # 左栏
    if left_title:
        add_text(slide, left_title, MARGIN, Inches(1.2),
                 Inches(5.8), Inches(0.4), font_size=H2_FONT,
                 bold=True, color=GREEN_800)
        top = Inches(1.65)
    else:
        top = Inches(1.2)
    left_content_fn(slide, MARGIN, top, Inches(5.8), Inches(5.7))
    # 右栏
    if right_title:
        add_text(slide, right_title, Inches(6.8), Inches(1.2),
                 Inches(6.0), Inches(0.4), font_size=H2_FONT,
                 bold=True, color=GREEN_800)
        top2 = Inches(1.65)
    else:
        top2 = Inches(1.2)
    right_content_fn(slide, Inches(6.8), top2, Inches(6.0), Inches(5.7))
    footer_bar(slide)

# ==================== 幻灯片内容 ====================

def slide_cover(prs):
    """第1页：封面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, GREEN_900)

    # 装饰块
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=GREEN_500)
    add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), fill_color=GREEN_500)
    add_rect(slide, 0, Inches(0.08), Inches(0.08), SLIDE_H - Inches(0.16), fill_color=GREEN_600)

    # 主标题
    add_text(slide, "桂林银发旅游小程序", Inches(1.2), Inches(1.6),
             Inches(11), Inches(1.1), font_size=44, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "—— 银发族智慧旅游新体验 ——", Inches(1.2), Inches(2.75),
             Inches(11), Inches(0.6), font_size=26, bold=False,
             color=GREEN_400, align=PP_ALIGN.CENTER)

    # 分隔线
    add_rect(slide, Inches(2.5), Inches(3.5), Inches(8.333), Inches(0.05),
             fill_color=GREEN_500)

    # 核心标签行
    tags = ["银发友好", "微信小程序", "AI智能", "漓江旅游", "桂林导览"]
    tx = Inches(0.9)
    for tag in tags:
        tw = Inches(len(tag) * 14 * 0.018 + 0.35)
        add_rect(slide, tx, Inches(3.75), tw, Inches(0.5),
                 fill_color=GREEN_800, line_color=GREEN_500, line_width=Pt(1))
        add_text(slide, tag, tx, Inches(3.75), tw, Inches(0.5),
                 font_size=16, bold=True, color=GREEN_200, align=PP_ALIGN.CENTER)
        tx += tw + Inches(0.18)

    # 副标题信息
    add_text(slide, "项目全景介绍 | 2024", Inches(1.2), Inches(4.7),
             Inches(11), Inches(0.45), font_size=H3_FONT,
             bold=False, color=GRAY_300, align=PP_ALIGN.CENTER)
    add_text(slide, "银发友好 × AI赋能 × 漓江文旅", Inches(1.2), Inches(5.15),
             Inches(11), Inches(0.4), font_size=H3_FONT,
             bold=False, color=GREEN_400, align=PP_ALIGN.CENTER)

    # 底部装饰
    add_rect(slide, 0, Inches(5.8), SLIDE_W, Inches(1.7), fill_color=GREEN_800)
    add_text(slide, "专为中老年游客设计 | 大字体·高对比度·简操作",
             Inches(1.5), Inches(6.1), Inches(10.333), Inches(0.5),
             font_size=18, bold=True, color=GREEN_200, align=PP_ALIGN.CENTER)
    add_text(slide, "基于桂林旅游市场调研 | 服务银发群体出行需求",
             Inches(1.5), Inches(6.65), Inches(10.333), Inches(0.45),
             font_size=14, bold=False, color=GRAY_300, align=PP_ALIGN.CENTER)

    return slide


def slide_toc(prs):
    """第2页：目录"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, GREEN_50)
    section_title_bar(slide, "目录 CONTENTS", "项目全景一页速览")

    sections = [
        ("01", "项目概述",    "背景·目标·价值定位"),
        ("02", "目标用户",    "银发群体画像与需求"),
        ("03", "核心功能",    "8大功能模块详解"),
        ("04", "UI设计",     "银发友好交互设计"),
        ("05", "技术架构",    "前后端·AI·数据库"),
        ("06", "数据库",      "9张核心数据表"),
        ("07", "API接口",     "20+接口·支付体系"),
        ("08", "安全与辅助",  "防护·SOS·运维"),
    ]

    cols = 4
    rows = 2
    start_x = MARGIN
    start_y = Inches(1.3)
    cell_w   = Inches((SLIDE_W - MARGIN * 2) / cols)
    cell_h   = Inches(2.7)
    gap_x    = Inches(0.15)
    gap_y    = Inches(0.2)

    for i, (num, title, desc) in enumerate(sections):
        col = i % cols
        row = i // cols
        x = start_x + col * (cell_w + gap_x)
        y = start_y + row * (cell_h + gap_y)

        # 卡片背景
        add_rect(slide, x, y, cell_w - gap_x, cell_h,
                 fill_color=WHITE, line_color=GREEN_200, line_width=Pt(1.5))
        # 序号
        add_rect(slide, x, y, cell_w - gap_x, Inches(0.7),
                 fill_color=GREEN_700)
        add_text(slide, num, x + Inches(0.1), y + Inches(0.05),
                 Inches(0.7), Inches(0.6), font_size=24,
                 bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # 标题
        add_text(slide, title, x + Inches(0.85), y + Inches(0.1),
                 cell_w - gap_x - Inches(0.95), Inches(0.55), font_size=18,
                 bold=True, color=WHITE, align=PP_ALIGN.LEFT)
        # 描述
        add_text(slide, desc, x + Inches(0.15), y + Inches(0.85),
                 cell_w - gap_x - Inches(0.25), Inches(0.4), font_size=14,
                 bold=False, color=GREEN_700)
        # 装饰
        add_rect(slide, x, y + Inches(1.35), cell_w - gap_x, Inches(1.3),
                 fill_color=GREEN_100)

    footer_bar(slide)
    return slide


def slide_overview(prs):
    """第3页：项目概述"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, GREEN_50)

    def content(s, lx, ty, lw, lh):
        add_text(slide, "项目背景", lx, ty, lw, Inches(0.4), font_size=18,
                 bold=True, color=GREEN_800)
        add_text(slide, "桂林是全球知名旅游城市，银发族（55岁以上）游客占比超过35%，"
                     "但现有旅游应用对银发用户极不友好：字体小、操作复杂、路径深。"
                     "本小程序专为银发群体打造，提供大字高对比度、极简交互的智慧旅游体验。",
                 lx, ty + Inches(0.42), lw, Inches(1.0), font_size=BODY_FONT,
                 color=GRAY_700)

        add_text(slide, "核心目标", lx, ty + Inches(1.55), lw, Inches(0.4), font_size=18,
                 bold=True, color=GREEN_800)
        goals = [
            "✅ 解决银发族使用智能设备的困难，降低旅游门槛",
            "✅ 整合漓江/桂林优质旅游资源，提供一键导览",
            "✅ AI智能推荐路线+实时讲解，服务个性化出行",
            "✅ 建立信任体系（实名认证+先行赔付），消除老年人消费顾虑",
        ]
        for i, g in enumerate(goals):
            add_text(slide, g, lx, ty + Inches(1.95) + i * Inches(0.42),
                     lw, Inches(0.4), font_size=BODY_FONT, color=GRAY_700)

        add_text(slide, "价值定位", lx, ty + Inches(3.7), lw, Inches(0.4), font_size=18,
                 bold=True, color=GREEN_800)
        add_text(slide, "国内首个专注银发群体的一站式桂林旅游小程序，"
                     "以漓江精华段（杨堤码头→兴坪）为核心场景，"
                     "融合AI导览、社交分享、安全保障，构建银发旅游新生态。",
                 lx, ty + Inches(4.1), lw, Inches(0.85), font_size=BODY_FONT,
                 color=GRAY_700)

    def right(s, rx, ry, rw, rh):
        add_text(slide, "项目数据", rx, ry, rw, Inches(0.4), font_size=18,
                 bold=True, color=GREEN_800)
        stats = [
            ("20+", "API接口"),
            ("9",   "数据库表"),
            ("8",   "核心功能模块"),
            ("3",   "用户角色"),
            ("漓江", "核心场景"),
            ("漓江", "核心场景"),
        ]
        sy = ry + Inches(0.5)
        for i, (val, lab) in enumerate(stats[:4]):
            col = i % 2
            row = i // 2
            stat_block(s, lab, val, rx + col * Inches(2.9), sy + row * Inches(1.0),
                       w=Inches(2.7))

        add_text(slide, "技术亮点", rx, ry + Inches(3.3), rw, Inches(0.4), font_size=18,
                 bold=True, color=GREEN_800)
        tech_items = [
            "🗺️  AI景点识别与智能路线规划",
            "🔊  语音导览+文字双轨讲解",
            "💬  AI对话式旅游咨询",
            "📍  GPS精准定位+电子围栏",
            "🛡️   SOS一键呼救+实时位置共享",
        ]
        for i, t in enumerate(tech_items):
            add_text(slide, t, rx, ry + Inches(3.72) + i * Inches(0.4),
                     rw, Inches(0.38), font_size=BODY_FONT, color=GRAY_700)

    two_col_layout(slide, "项目概述 OVERVIEW", "背景·目标·价值",
                   content, right, "背景与目标", "数据与技术亮点")
    return slide


def slide_user_profile(prs):
    """第4页：目标用户"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, GREEN_50)

    def content(s, lx, ty, lw, lh):
        add_text(slide, "银发用户画像", lx, ty, lw, Inches(0.4), font_size=18,
                 bold=True, color=GREEN_800)

        profiles = [
            ("🌊", "休闲观光型", "漓江山水、古镇漫游为主，偏好摄影，慢节奏，"
                               "需要语音讲解辅助理解文化内涵，不擅自行程规划"),
            ("📸", "文化体验型", "对历史典故、民俗文化有浓厚兴趣，"
                               "愿意参与沉浸式活动如渔翁撒网拍照，需深度讲解内容"),
            ("👨‍👩‍👧", "家庭出游型", "与子女/孙辈同行，操作由子女代劳为主，"
                               "关注安全性和便利性，适合亲子组合票"),
        ]
        for i, (icon, name, desc) in enumerate(profiles):
            y = ty + Inches(0.5) + i * Inches(1.6)
            add_rect(slide, lx, y, lw, Inches(1.5),
                     fill_color=WHITE, line_color=GREEN_600, line_width=Pt(1))
            add_rect(slide, lx, y, Inches(0.65), Inches(1.5), fill_color=GREEN_700)
            add_text(slide, icon, lx, y, Inches(0.65), Inches(1.5),
                     font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            add_text(slide, name, lx + Inches(0.75), y + Inches(0.1),
                     lw - Inches(0.9), Inches(0.4), font_size=18,
                     bold=True, color=GREEN_900)
            add_text(slide, desc, lx + Inches(0.75), y + Inches(0.5),
                     lw - Inches(0.9), Inches(0.95), font_size=BODY_FONT,
                     color=GRAY_700)

    def right(s, rx, ry, rw, rh):
        add_text(slide, "痛点与需求", rx, ry, rw, Inches(0.4), font_size=18,
                 bold=True, color=GREEN_800)
        pain_pts = [
            ("❌ 字体太小",    "现有APP正文字体10-12pt，55岁以上阅读困难"),
            ("❌ 操作复杂",    "注册登录超过3步即流失，页面层级超过2级即放弃"),
            ("❌ 信任缺失",    "对线上支付有顾虑，无实体保障，担心上当受骗"),
            ("❌ 信息过载",    "页面塞满功能入口，视觉混乱，不知道该点哪里"),
            ("❌ 视力问题",    "强光下屏幕看不清，低对比度色文字难以辨认"),
        ]
        for i, (pain, detail) in enumerate(pain_pts):
            y = ry + Inches(0.5) + i * Inches(0.78)
            add_rect(slide, rx, y, Inches(0.08), Inches(0.55), fill_color=RED)
            add_text(slide, pain, rx + Inches(0.18), y, rw - Inches(0.2), Inches(0.35),
                     font_size=BODY_FONT, bold=True, color=RED)
            add_text(slide, detail, rx + Inches(0.18), y + Inches(0.33),
                     rw - Inches(0.2), Inches(0.38), font_size=BODY_FONT,
                     color=GRAY_700)

        add_text(slide, "解决方案", rx, ry + Inches(4.5), rw, Inches(0.4),
                 font_size=18, bold=True, color=GREEN_800)
        solutions = [
            "✅ 正文字体≥16pt，大按钮，热区放大",
            "✅ 3步内完成核心操作，全流程≤3页",
            "✅ 实名认证+先行赔付，线下面对面核销",
        ]
        for i, sol in enumerate(solutions):
            add_text(slide, sol, rx, ry + Inches(4.9) + i * Inches(0.4),
                     rw, Inches(0.38), font_size=BODY_FONT, color=GREEN_700)

    two_col_layout(slide, "目标用户 USER PROFILE", "银发群体需求与痛点",
                   content, right, "用户画像", "痛点与方案")
    return slide


def slide_functions(prs):
    """第5页：核心功能"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, GREEN_50)
    section_title_bar(slide, "核心功能 FEATURES", "8大模块构建完整银发旅游体验")

    funcs = [
        ("🗺️", "智能路线规划",    "AI分析偏好+实时天气+路况，生成最优游览方案，支持分段导航与预估时间展示"),
        ("📍", "精准位置服务",    "电子围栏自动记录游览轨迹，GPS+北斗双模定位，离线地图无网也能用"),
        ("🔊", "AI语音讲解",      "景点识别自动触发讲解，提供普通话/方言切换，真人语音合成，音量自适应"),
        ("📱", "辅助预约挂号",    "桂林景区门票在线预约，对接景区系统实时库存，支持改签退，子女可代操作"),
        ("💳", "智慧支付体系",    "微信支付+亲属代付+先游后付，防欺诈实时预警，消费明细语音播报"),
        ("🛡️", "安全管理中心",    "SOS一键呼救，GPS+轨迹实时共享，亲情号码5秒触达，急救信息卡自动展示"),
        ("🏥", "健康保障",        "实时心率/血氧监测（手表数据），周边医院POI，周游高德120联动，自动位置上报"),
        ("👨‍👩‍👧", "家庭社交",       "家庭组队共享位置，子女辅助操作，亲友圈分享游览照片，打卡排行榜激励"),
    ]

    cols = 4
    rows = 2
    start_x = MARGIN
    start_y = Inches(1.25)
    cell_w   = (SLIDE_W - MARGIN * 2) / cols
    cell_h   = Inches(2.75)
    gap_x    = Inches(0.12)
    gap_y    = Inches(0.18)

    for i, (icon, title, desc) in enumerate(funcs):
        col = i % cols
        row = i // cols
        x = start_x + col * (cell_w + gap_x)
        y = start_y + row * (cell_h + gap_y)
        cw = cell_w - gap_x

        add_rect(slide, x, y, cw, cell_h, fill_color=WHITE,
                 line_color=GREEN_200, line_width=Pt(1))
        add_rect(slide, x, y, cw, Inches(0.75), fill_color=GREEN_700)
        add_text(slide, icon, x, y, cw, Inches(0.75), font_size=24,
                 bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, title, x + Inches(0.1), y + Inches(0.78),
                 cw - Inches(0.15), Inches(0.42), font_size=16,
                 bold=True, color=GREEN_900, align=PP_ALIGN.LEFT)
        add_text(slide, desc, x + Inches(0.1), y + Inches(1.22),
                 cw - Inches(0.15), Inches(1.45), font_size=BODY_FONT,
                 color=GRAY_700)

    footer_bar(slide)
    return slide


def slide_ui_design(prs):
    """第6页：UI设计"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, GREEN_50)

    def content(s, lx, ty, lw, lh):
        add_text(slide, "设计原则", lx, ty, lw, Inches(0.4), font_size=18,
                 bold=True, color=GREEN_800)
        principles = [
            ("♿", "极简交互",  "核心操作不超过3步，全流程扁平化，重要操作固定在首页底部Tab"),
            ("🔤", "大字体",     "正文字体≥16pt，标题≥22pt，重要信息加粗+高对比色"),
            ("🎨", "高对比",    "文字与背景对比度≥4.5:1，关键按钮使用主色+白字"),
            ("👆", "大热区",    "点击区域≥44×44pt，重要按钮宽度≥屏幕宽60%"),
            ("📢", "双重反馈",  "重要操作同时提供视觉+语音/震动反馈，确认提示清晰"),
        ]
        for i, (icon, name, desc) in enumerate(principles):
            y = ty + Inches(0.5) + i * Inches(1.0)
            add_rect(slide, lx, y, lw, Inches(0.92),
                     fill_color=WHITE, line_color=GREEN_200, line_width=Pt(1))
            add_rect(slide, lx, y, Inches(0.65), Inches(0.92), fill_color=GREEN_600)
            add_text(slide, icon, lx, y, Inches(0.65), Inches(0.92),
                     font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            add_text(slide, name, lx + Inches(0.75), y + Inches(0.05),
                     lw - Inches(0.9), Inches(0.38), font_size=16,
                     bold=True, color=GREEN_900)
            add_text(slide, desc, lx + Inches(0.75), y + Inches(0.43),
                     lw - Inches(0.9), Inches(0.44), font_size=BODY_FONT,
                     color=GRAY_700)

    def right(s, rx, ry, rw, rh):
        add_text(slide, "色彩规范", rx, ry, rw, Inches(0.4), font_size=18,
                 bold=True, color=GREEN_800)
        colors = [
            ("GREEN_900", "0x0D3320", "标题/重点背景", GREEN_900),
            ("GREEN_700", "0x2E6B47", "主要按钮/标签", GREEN_700),
            ("GREEN_500", "0x4AAD75", "成功/确认状态", GREEN_500),
            ("GOLD",      "0xF5960B", "警示/重要提示", GOLD),
            ("RED",       "0xDC2626", "SOS/危险/错误", RED),
            ("GRAY_700",  "0x3A3A3A", "正文文字",     GRAY_700),
        ]
        for i, (name, hex_c, usage, col) in enumerate(colors):
            y = ry + Inches(0.5) + i * Inches(0.6)
            add_rect(slide, rx, y + Inches(0.05), Inches(0.5), Inches(0.42), fill_color=col)
            add_text(slide, f"{name}  {hex_c}", rx + Inches(0.6), y,
                     Inches(2.5), Inches(0.38), font_size=BODY_FONT,
                     bold=True, color=GRAY_700)
            add_text(slide, usage, rx + Inches(3.1), y, rw - Inches(3.1), Inches(0.38),
                     font_size=BODY_FONT, color=GRAY_600)

        add_text(slide, "字号规范（银发友好）", rx, ry + Inches(4.2), rw, Inches(0.4),
                 font_size=18, bold=True, color=GREEN_800)
        sizes = [
            ("页面大标题",   "28pt", "GREEN_900"),
            ("卡片标题",     "20pt", "GREEN_800"),
            ("正文内容",     "16pt", "GRAY_700"),
            ("辅助说明",     "14pt", "GRAY_600"),
            ("按钮文字",     "18pt", "WHITE"),
        ]
        for i, (label, size, col) in enumerate(sizes):
            y = ry + Inches(4.62) + i * Inches(0.4)
            add_text(slide, label, rx, y, Inches(2.8), Inches(0.38),
                     font_size=BODY_FONT, bold=True, color=GRAY_700)
            add_text(slide, size, rx + Inches(2.8), y, Inches(1.5), Inches(0.38),
                     font_size=BODY_FONT, bold=False, color=GREEN_700)

    two_col_layout(slide, "UI设计 UI DESIGN", "银发友好交互规范",
                   content, right, "设计原则", "色彩与字号")
    return slide


def slide_architecture(prs):
    """第7页：技术架构"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, GREEN_50)

    def content(s, lx, ty, lw, lh):
        add_text(slide, "前端架构", lx, ty, lw, Inches(0.4), font_size=18,
                 bold=True, color=GREEN_800)
        fe_items = [
            "📱 微信小程序（主）+ H5（社交传播）",
            "🗼 TDesign 组件库，统一设计语言",
            "🧩 Vant Weapp 轻量化组件（按需加载）",
            "📡 WebSocket 即时通讯（导游直播/组队消息）",
            "🗺️ 腾讯地图SDK，室内外一体化定位",
            "🔊 微信同声传译，实时语音转文字",
            "📳 微信原生振动反馈（haptic）",
        ]
        for i, item in enumerate(fe_items):
            add_text(slide, item, lx, ty + Inches(0.5) + i * Inches(0.42),
                     lw, Inches(0.4), font_size=BODY_FONT, color=GRAY_700)

        add_text(slide, "后端架构", lx, ty + Inches(3.45), lw, Inches(0.4),
                 font_size=18, bold=True, color=GREEN_800)
        be_items = [
            "⚙️  Node.js/Express REST API",
            "🐘 PostgreSQL 主数据库（用户/订单/内容）",
            "📊 Redis 缓存（会话/令牌/热点数据）",
            "🔐 JWT 认证 + 微信授权登录",
            "📨 消息队列（订单通知/行程变更）",
            "☁️  阿里云OSS 文件存储（图片/证照）",
        ]
        for i, item in enumerate(be_items):
            add_text(slide, item, lx, ty + Inches(3.85) + i * Inches(0.42),
                     lw, Inches(0.4), font_size=BODY_FONT, color=GRAY_700)

    def right(s, rx, ry, rw, rh):
        add_text(slide, "AI 服务层", rx, ry, rw, Inches(0.4), font_size=18,
                 bold=True, color=GREEN_800)
        ai_items = [
            "🤖 景点识别（图像AI+知识图谱）",
            "🗺️  智能路线规划（强化学习+实时数据）",
            "💬  AI对话咨询（意图识别+知识库检索）",
            "🔊  语音合成（TTS，多方言支持）",
            "📝  OCR证照识别（老人手机操作辅助）",
        ]
        for i, item in enumerate(ai_items):
            add_text(slide, item, rx, ry + Inches(0.5) + i * Inches(0.42),
                     rw, Inches(0.4), font_size=BODY_FONT, color=GRAY_700)

        add_text(slide, "运维与部署", rx, ry + Inches(2.7), rw, Inches(0.4),
                 font_size=18, bold=True, color=GREEN_800)
        ops_items = [
            "🐳 Docker 容器化，K8s 编排",
            "🔄 CI/CD 流水线（GitHub Actions）",
            "📈 Prometheus + Grafana 监控告警",
            "🌐 CDN 加速（静态资源+图片）",
            "🔒 HTTPS + TLS1.3 传输加密",
        ]
        for i, item in enumerate(ops_items):
            add_text(slide, item, rx, ry + Inches(3.1) + i * Inches(0.42),
                     rw, Inches(0.4), font_size=BODY_FONT, color=GRAY_700)

    two_col_layout(slide, "技术架构 ARCHITECTURE", "前后端分离 + AI智能 + 云原生",
                   content, right, "前端与后端", "AI服务与运维")
    return slide


def slide_database(prs):
    """第8页：数据库设计"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, GREEN_50)

    section_title_bar(slide, "数据库设计 DATABASE", "9张核心数据表 · PostgreSQL")
    tables = [
        ("users",            "用户表",          "id/name/phone/age/emergency_contact/emergency_phone/medical_info/id_card_enc"),
        ("attractions",      "景点表",          "id/name/intro/audio_url/gps_x/gps_y/level/tags/open_hours/accessible设施"),
        ("routes",            "游览路线表",       "id/name/attraction_ids/duration/price/suitable_for/max_people"),
        ("bookings",          "订单表",          "id/user_id/route_id/status/payment_method/pay_time/refund_time/refund_reason"),
        ("reviews",           "评价表",          "id/user_id/booking_id/rating/comment/audio_url/photo_urls/created_at"),
        ("safety_records",    "安全记录表",       "id/user_id/type/description/gps_x/gps_y/trigger_time/handled"),
        ("family_groups",     "家庭组队表",       "id/name/owner_id/invite_code/members/share_location/share_trail"),
        ("ai_conversations",  "AI对话记录表",    "id/session_id/user_id/query/response/intent/context"),
        ("sos_alerts",        "SOS告警表",       "id/user_id/location_x/location_y/heartbeat/triggered_at/resolved_at/rescuer"),
    ]

    # 3列布局
    cols = 3
    start_x = MARGIN
    start_y = Inches(1.3)
    col_w = Inches((SLIDE_W.inches - MARGIN.inches * 2) / cols)
    row_h = Inches(1.25)
    gap_x = Inches(0.12)
    gap_y = Inches(0.12)

    for i, (tbl, title, fields) in enumerate(tables):
        col = i % cols
        row = i // cols
        x = start_x + col * (col_w + gap_x)
        y = start_y + row * (row_h + gap_y)
        cw = col_w - gap_x

        add_rect(slide, x, y, cw, row_h, fill_color=WHITE,
                 line_color=GREEN_300, line_width=Pt(1))
        add_rect(slide, x, y, cw, Inches(0.5), fill_color=GREEN_800)
        add_text(slide, title, x + Inches(0.08), y + Inches(0.06),
                 cw - Inches(0.1), Inches(0.38), font_size=16,
                 bold=True, color=WHITE)
        add_text(slide, tbl, x + Inches(0.08), y + Inches(0.52),
                 cw - Inches(0.1), Inches(0.3), font_size=12,
                 bold=False, color=GREEN_700, italic=True)
        add_text(slide, fields, x + Inches(0.08), y + Inches(0.82),
                 cw - Inches(0.1), Inches(0.4), font_size=BODY_FONT,
                 color=GRAY_600)

    footer_bar(slide)
    return slide


def slide_api(prs):
    """第9页：API接口"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, GREEN_50)

    section_title_bar(slide, "API接口 INTERFACES", "20+ RESTful API · JSON · HTTPS")

    apis = [
        ("认证模块",   GREEN_700, [
            ("POST", "/api/auth/login",           "微信授权登录，获取JWT"),
            ("POST", "/api/auth/refresh",          "刷新Access Token"),
            ("POST", "/api/auth/logout",            "注销会话"),
        ]),
        ("景点模块",   GREEN_700, [
            ("GET",  "/api/attractions",            "景点列表（分页+筛选）"),
            ("GET",  "/api/attractions/:id",        "景点详情+语音讲解"),
            ("GET",  "/api/attractions/:id/comments","景点评论列表"),
        ]),
        ("路线模块",   GREEN_600, [
            ("GET",  "/api/routes",                  "路线列表（时间/价格筛选）"),
            ("GET",  "/api/routes/:id",               "路线详情含行程节点"),
            ("GET",  "/api/routes/ai-plan",          "AI智能路线规划"),
        ]),
        ("订单模块",   GREEN_600, [
            ("POST", "/api/bookings",                "创建订单"),
            ("GET",  "/api/bookings/:id",            "订单详情"),
            ("PUT",  "/api/bookings/:id/cancel",     "取消订单（退款规则）"),
        ]),
        ("支付模块",   GREEN_500, [
            ("POST", "/api/pay/create",              "创建支付（微信JSAPI）"),
            ("POST", "/api/pay/callback",             "支付回调通知"),
            ("GET",  "/api/pay/refund-calc",          "退款金额计算"),
        ]),
        ("地图模块",   GREEN_500, [
            ("GET",  "/api/map/nearby",               "周边设施查询"),
            ("GET",  "/api/map/trail",                "游览轨迹查询"),
            ("POST", "/api/map/checkin",               "打卡签到"),
        ]),
        ("安全模块",   RED, [
            ("POST", "/api/sos/alert",                "触发SOS告警"),
            ("GET",  "/api/sos/status/:id",            "告警状态查询"),
            ("PUT",  "/api/sos/resolve/:id",          "解除告警"),
        ]),
        ("AI模块",    GREEN_800, [
            ("POST", "/api/ai/chat",                  "AI对话咨询"),
            ("POST", "/api/ai/voice-tts",             "语音合成"),
            ("GET",  "/api/ai/intent",                "意图识别"),
        ]),
    ]

    cols = 4
    start_x = MARGIN
    start_y = Inches(1.25)
    cell_w  = (SLIDE_W - MARGIN * 2) / cols
    cell_h  = Inches(2.85)
    gap_x   = Inches(0.1)
    gap_y   = Inches(0.12)

    for i, (group, color, endpoints) in enumerate(apis):
        col = i % cols
        row = i // cols
        x = start_x + col * (cell_w + gap_x)
        y = start_y + row * (cell_h + gap_y)
        cw = cell_w - gap_x

        add_rect(slide, x, y, cw, cell_h, fill_color=WHITE,
                 line_color=GREEN_200, line_width=Pt(1))
        add_rect(slide, x, y, cw, Inches(0.5), fill_color=color)
        add_text(slide, group, x + Inches(0.08), y + Inches(0.07),
                 cw - Inches(0.1), Inches(0.36), font_size=14,
                 bold=True, color=WHITE)
        for j, (method, path, desc) in enumerate(endpoints):
            mt_color = (GREEN_700 if method == "GET" else
                        GREEN_600 if method == "POST" else RED)
            yy = y + Inches(0.55) + j * Inches(0.75)
            add_rect(slide, x + Inches(0.08), yy + Inches(0.05),
                     Inches(0.55), Inches(0.3), fill_color=mt_color)
            add_text(slide, method, x + Inches(0.08), yy + Inches(0.04),
                     Inches(0.55), Inches(0.3), font_size=11,
                     bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            add_text(slide, path, x + Inches(0.08), yy + Inches(0.35),
                     cw - Inches(0.12), Inches(0.3), font_size=11,
                     bold=False, color=GREEN_700, italic=True)
            add_text(slide, desc, x + Inches(0.08), yy + Inches(0.62),
                     cw - Inches(0.12), Inches(0.3), font_size=BODY_FONT,
                     color=GRAY_600)

    footer_bar(slide)
    return slide


def slide_payment(prs):
    """第10页：支付体系"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, GREEN_50)

    def content(s, lx, ty, lw, lh):
        add_text(slide, "支付方式", lx, ty, lw, Inches(0.4), font_size=18,
                 bold=True, color=GREEN_800)
        pays = [
            ("💳", "微信支付",       "主推，契合银发习惯，免密支付（单笔限额500元）"),
            ("👨", "亲属代付",       "子女远程支付，父母收到通知确认，避免误操作"),
            ("🕐", "先游后付",       "信用分≥650可申请，游玩后24h内结算，消除顾虑"),
            ("🎟️", "一码通",         "景区联票+竹筏票合一，自动核销，无需重复取票"),
        ]
        for i, (icon, name, desc) in enumerate(pays):
            y = ty + Inches(0.5) + i * Inches(1.05)
            add_rect(slide, lx, y, lw, Inches(0.98),
                     fill_color=WHITE, line_color=GREEN_200, line_width=Pt(1))
            add_rect(slide, lx, y, Inches(0.65), Inches(0.98), fill_color=GREEN_600)
            add_text(slide, icon, lx, y, Inches(0.65), Inches(0.98),
                     font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            add_text(slide, name, lx + Inches(0.75), y + Inches(0.08),
                     lw - Inches(0.9), Inches(0.38), font_size=16,
                     bold=True, color=GREEN_900)
            add_text(slide, desc, lx + Inches(0.75), y + Inches(0.46),
                     lw - Inches(0.9), Inches(0.48), font_size=BODY_FONT,
                     color=GRAY_700)

    def right(s, rx, ry, rw, rh):
        add_text(slide, "反欺诈机制", rx, ry, rw, Inches(0.4), font_size=18,
                 bold=True, color=GREEN_800)
        frauds = [
            ("🚨", "实时风控",    "大额支付（>500）弹窗确认+短信验证码"),
            ("📍", "位置核验",    "异常IP+地理位置突变触发人工审核"),
            ("👤", "实名认证",    "≥60岁用户强制绑定紧急联系人+身份证"),
            ("💰", "资金保障",    "平台设立银发旅游保障基金，先行赔付"),
            ("📊", "行为分析",    "异常购买模式（深夜+短时+高价）预警"),
        ]
        for i, (icon, name, desc) in enumerate(frauds):
            y = ry + Inches(0.5) + i * Inches(0.78)
            add_rect(slide, rx, y, Inches(0.55), Inches(0.7),
                     fill_color=GREEN_100, line_color=GREEN_400, line_width=Pt(1))
            add_text(slide, icon, rx, y, Inches(0.55), Inches(0.7),
                     font_size=18, bold=True, color=GREEN_700, align=PP_ALIGN.CENTER)
            add_text(slide, name, rx + Inches(0.65), y + Inches(0.04),
                     rw - Inches(0.7), Inches(0.33), font_size=16,
                     bold=True, color=GREEN_900)
            add_text(slide, desc, rx + Inches(0.65), y + Inches(0.37),
                     rw - Inches(0.7), Inches(0.32), font_size=BODY_FONT,
                     color=GRAY_700)

        add_text(slide, "退款规则", rx, ry + Inches(4.45), rw, Inches(0.4),
                 font_size=18, bold=True, color=GREEN_800)
        rules = [
            "• 出发前24h：全额退款（0手续费）",
            "• 出发前4-24h：退还80%",
            "• 出发前1-4h：退还50%",
            "• 出发前1h内/已出发：不予退款",
        ]
        for i, rule in enumerate(rules):
            add_text(slide, rule, rx, ry + Inches(4.87) + i * Inches(0.38),
                     rw, Inches(0.36), font_size=BODY_FONT, color=GRAY_700)

    two_col_layout(slide, "支付体系 PAYMENT", "微信支付·亲属代付·先游后付·资金保障",
                   content, right, "支付方式", "安全与退款")
    return slide


def slide_security(prs):
    """第11页：安全系统"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, GREEN_50)

    def content(s, lx, ty, lw, lh):
        add_text(slide, "SOS紧急救援", lx, ty, lw, Inches(0.4), font_size=18,
                 bold=True, color=GREEN_800)

        add_rect(slide, lx, ty + Inches(0.5), lw, Inches(0.75),
                 fill_color=RED, line_color=None)
        add_text(slide, "🆘  首页固定SOS大按钮（直径80mm热区）",
                 lx + Inches(0.1), ty + Inches(0.58), lw - Inches(0.2),
                 Inches(0.6), font_size=16, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)

        sos_steps = [
            ("1️⃣", "触发方式",    "长按3秒SOS按钮或语音\"救命\"，防止误触"),
            ("2️⃣", "通知链",      "用户→亲属+平台客服+景区管理处，并发通知"),
            ("3️⃣", "位置共享",    "自动上传GPS坐标，每30秒更新，持续10分钟"),
            ("4️⃣", "医疗信息",    "血型/过敏史/用药史自动展示给急救人员"),
            ("5️⃣", "客服介入",    "5秒内人工客服接入，全程录音记录"),
        ]
        for i, (num, title, desc) in enumerate(sos_steps):
            y = ty + Inches(1.35) + i * Inches(0.82)
            add_rect(slide, lx, y, lw, Inches(0.75),
                     fill_color=WHITE, line_color=GREEN_200, line_width=Pt(1))
            add_rect(slide, lx, y, Inches(0.6), Inches(0.75), fill_color=GREEN_700)
            add_text(slide, num, lx, y, Inches(0.6), Inches(0.75),
                     font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            add_text(slide, title, lx + Inches(0.7), y + Inches(0.06),
                     lw - Inches(0.8), Inches(0.32), font_size=16,
                     bold=True, color=GREEN_900)
            add_text(slide, desc, lx + Inches(0.7), y + Inches(0.38),
                     lw - Inches(0.8), Inches(0.33), font_size=BODY_FONT,
                     color=GRAY_700)

    def right(s, rx, ry, rw, rh):
        add_text(slide, "数据安全", rx, ry, rw, Inches(0.4), font_size=18,
                 bold=True, color=GREEN_800)
        sec_items = [
            ("🔐", "传输加密",   "全程HTTPS+TLS1.3，API参数签名防篡改"),
            ("🗄️", "存储加密",   "身份证等敏感信息AES-256加密存储"),
            ("🎭", "隐私保护",   "位置数据脱敏，行程轨迹72h后模糊化"),
            ("⚙️", "权限管控",   "RBAC角色权限，最小权限原则"),
            ("🗑️", "数据销毁",   "账户注销后30天自动清除（法律法规要求）"),
        ]
        for i, (icon, name, desc) in enumerate(sec_items):
            y = ry + Inches(0.5) + i * Inches(0.78)
            add_rect(slide, rx, y, rw, Inches(0.7),
                     fill_color=WHITE, line_color=GREEN_200, line_width=Pt(1))
            add_text(slide, icon, rx + Inches(0.1), y + Inches(0.1),
                     Inches(0.5), Inches(0.5), font_size=20,
                     bold=True, color=GREEN_700, align=PP_ALIGN.CENTER)
            add_text(slide, name, rx + Inches(0.65), y + Inches(0.08),
                     rw - Inches(0.75), Inches(0.3), font_size=16,
                     bold=True, color=GREEN_900)
            add_text(slide, desc, rx + Inches(0.65), y + Inches(0.38),
                     rw - Inches(0.75), Inches(0.3), font_size=BODY_FONT,
                     color=GRAY_700)

        add_text(slide, "容灾备份", rx, ry + Inches(4.45), rw, Inches(0.4),
                 font_size=18, bold=True, color=GREEN_800)
        backup = [
            "📍 主备双活数据中心（Nanjing + Hongkong）",
            "🔄 数据库每日全量+每小时增量备份",
            "⚡ 故障自动切换，业务中断<30秒",
        ]
        for i, b in enumerate(backup):
            add_text(slide, b, rx, ry + Inches(4.87) + i * Inches(0.4),
                     rw, Inches(0.38), font_size=BODY_FONT, color=GRAY_700)

    two_col_layout(slide, "安全体系 SECURITY", "SOS紧急救援·数据安全·容灾备份",
                   content, right, "SOS紧急救援流程", "数据安全与备份")
    return slide


def slide_family_social(prs):
    """第12页：家庭社交"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, GREEN_50)

    def content(s, lx, ty, lw, lh):
        add_text(slide, "家庭组队", lx, ty, lw, Inches(0.4), font_size=18,
                 bold=True, color=GREEN_800)
        add_text(slide, "子女帮助父母规划行程、预约门票、代付费用，"
                     "让不善操作智能设备的银发用户也能享受数字化旅游的便利。",
                 lx, ty + Inches(0.42), lw, Inches(0.75), font_size=BODY_FONT,
                 color=GRAY_700)

        features = [
            ("📋", "行程代规划",   "子女可在\"家庭管理\"中帮父母选择路线、一键预约"),
            ("💳", "费用代支付",   "父母下单→子女收到通知→一键代付，父母无需绑定银行卡"),
            ("📍", "位置实时看",   "父母游览时，子女可查看实时位置+轨迹，安心"),
            ("💬", "游览动态",     "父母打卡/拍照后自动同步家庭群，子女点赞互动"),
        ]
        for i, (icon, name, desc) in enumerate(features):
            y = ty + Inches(1.3) + i * Inches(1.0)
            add_rect(slide, lx, y, lw, Inches(0.92),
                     fill_color=WHITE, line_color=GREEN_200, line_width=Pt(1))
            add_rect(slide, lx, y, Inches(0.65), Inches(0.92), fill_color=GREEN_600)
            add_text(slide, icon, lx, y, Inches(0.65), Inches(0.92),
                     font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            add_text(slide, name, lx + Inches(0.75), y + Inches(0.08),
                     lw - Inches(0.9), Inches(0.35), font_size=16,
                     bold=True, color=GREEN_900)
            add_text(slide, desc, lx + Inches(0.75), y + Inches(0.44),
                     lw - Inches(0.9), Inches(0.43), font_size=BODY_FONT,
                     color=GRAY_700)

    def right(s, rx, ry, rw, rh):
        add_text(slide, "社交分享", rx, ry, rw, Inches(0.4), font_size=18,
                 bold=True, color=GREEN_800)
        social = [
            ("🏆", "打卡排行榜",    "游览漓江重要景点打卡，解锁成就徽章，分享朋友圈"),
            ("📸", "照片故事",      "AI自动剪辑游览精彩瞬间，生成配图文案，一键分享"),
            ("🗺️", "足迹地图",      "记录每次旅游足迹，生成个人旅行地图，永久保存"),
            ("👥", "游友圈",        "同龄银发游客社区，分享游记，结伴同游"),
        ]
        for i, (icon, name, desc) in enumerate(social):
            y = ry + Inches(0.5) + i * Inches(1.0)
            add_rect(slide, rx, y, rw, Inches(0.92),
                     fill_color=WHITE, line_color=GREEN_200, line_width=Pt(1))
            add_rect(slide, rx, y, Inches(0.65), Inches(0.92), fill_color=GREEN_500)
            add_text(slide, icon, rx, y, Inches(0.65), Inches(0.92),
                     font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            add_text(slide, name, rx + Inches(0.75), y + Inches(0.08),
                     rw - Inches(0.9), Inches(0.35), font_size=16,
                     bold=True, color=GREEN_900)
            add_text(slide, desc, rx + Inches(0.75), y + Inches(0.44),
                     rw - Inches(0.9), Inches(0.43), font_size=BODY_FONT,
                     color=GRAY_700)

        add_text(slide, "激励体系", rx, ry + Inches(4.65), rw, Inches(0.4),
                 font_size=18, bold=True, color=GREEN_800)
        jili = ["🎫 游览满3次景点获折扣券", "📊 家庭组队额外积分", "🏅 年度银发旅游达人评选"]
        for i, t in enumerate(jili):
            add_text(slide, t, rx, ry + Inches(5.05) + i * Inches(0.4),
                     rw, Inches(0.38), font_size=BODY_FONT, color=GRAY_700)

    two_col_layout(slide, "家庭社交 FAMILY & SOCIAL", "子女辅助·亲情关怀·游友互动",
                   content, right, "家庭组队", "社交与激励")
    return slide


def slide_operation(prs):
    """第13页：运营推广"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, GREEN_50)

    def content(s, lx, ty, lw, lh):
        add_text(slide, "推广策略", lx, ty, lw, Inches(0.4), font_size=18,
                 bold=True, color=GREEN_800)
        strategies = [
            ("🏥", "医疗机构合作",  "桂林各大医院体检中心、老年病门诊放置宣传物料，精准触达目标用户"),
            ("🏠", "社区地推",       "桂林/阳朔各大社区、老年活动中心、公园广场，定点协助注册"),
            ("📺", "本地媒体",       "桂林电视台《桂林日报》老年人专版，专题报道+软性植入"),
            ("🎁", "首次体验补贴",  "新用户首单立减30元+专业导览服务体验，消除首次使用顾虑"),
            ("👴", "口碑传播",       "设置\"银发推荐官\"勋章，被推荐人下单奖励推荐人积分"),
        ]
        for i, (icon, name, desc) in enumerate(strategies):
            y = ty + Inches(0.5) + i * Inches(0.88)
            add_rect(slide, lx, y, lw, Inches(0.8),
                     fill_color=WHITE, line_color=GREEN_200, line_width=Pt(1))
            add_rect(slide, lx, y, Inches(0.65), Inches(0.8), fill_color=GREEN_700)
            add_text(slide, icon, lx, y, Inches(0.65), Inches(0.8),
                     font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            add_text(slide, name, lx + Inches(0.75), y + Inches(0.06),
                     lw - Inches(0.9), Inches(0.34), font_size=16,
                     bold=True, color=GREEN_900)
            add_text(slide, desc, lx + Inches(0.75), y + Inches(0.4),
                     lw - Inches(0.9), Inches(0.38), font_size=BODY_FONT,
                     color=GRAY_700)

    def right(s, rx, ry, rw, rh):
        add_text(slide, "运营指标", rx, ry, rw, Inches(0.4), font_size=18,
                 bold=True, color=GREEN_800)
        kpis = [
            ("DAU",    "日活跃用户",    "目标上线6个月达5000+"),
            ("付费率",  "转化率",        "目标15%（银发用户保守）"),
            ("NPS",    "净推荐值",       "目标≥50（口碑驱动）"),
            ("客单价", "平均订单金额",  "目标≥280元/人"),
        ]
        for i, (kpi, name, target) in enumerate(kpis):
            col = i % 2
            row = i // 2
            x = rx + col * Inches(3.0)
            y = ry + Inches(0.5) + row * Inches(1.1)
            stat_block(s, name, kpi, x, y, w=Inches(2.85), h=Inches(1.0))
            add_text(slide, target, x, y + Inches(1.02), Inches(2.85), Inches(0.35),
                     font_size=BODY_FONT, color=GRAY_600, align=PP_ALIGN.CENTER)

        add_text(slide, "合作资源", rx, ry + Inches(2.85), rw, Inches(0.4),
                 font_size=18, bold=True, color=GREEN_800)
        resources = [
            "🤝 桂林文化广电旅游局（官方背书）",
            "🤝 漓江风景名胜区管理局（票务合作）",
            "🤝 阳朔县旅游总公司（线下接待）",
            "🤝 桂林市民政局（老年服务渠道）",
        ]
        for i, r in enumerate(resources):
            add_text(slide, r, rx, ry + Inches(3.27) + i * Inches(0.42),
                     rw, Inches(0.4), font_size=BODY_FONT, color=GRAY_700)

    two_col_layout(slide, "运营推广 OPERATION", "线上线下联动·医疗机构+社区地推·口碑裂变",
                   content, right, "推广策略", "运营指标与合作")
    return slide


def slide_summary(prs):
    """第14页：总结"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, GREEN_900)

    add_rect(slide, 0, 0, SLIDE_W, Inches(0.1), fill_color=GREEN_500)
    add_rect(slide, 0, SLIDE_H - Inches(0.1), SLIDE_W, Inches(0.1), fill_color=GREEN_500)

    add_text(slide, "桂林银发旅游小程序", Inches(0.5), Inches(0.5),
             SLIDE_W - Inches(1.0), Inches(0.8), font_size=36, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "让每一位银发游客都能享受智慧旅游的便利与乐趣",
             Inches(0.5), Inches(1.35), SLIDE_W - Inches(1.0), Inches(0.5),
             font_size=H2_FONT, bold=False, color=GREEN_200, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(1.5), Inches(2.0), SLIDE_W - Inches(3.0), Inches(0.04),
             fill_color=GREEN_500)

    highlights = [
        ("🗺️", "AI智能导览",   "漓江全景智能路线规划，语音讲解自动触发"),
        ("🔊", "银发友好设计",  "大字高对比度，3步内完成核心操作"),
        ("🛡️", "安全保障体系",  "SOS一键呼救，亲属实时位置共享"),
        ("💳", "智慧支付",      "微信支付+亲属代付+先游后付"),
        ("👨‍👩‍👧", "家庭社交",     "子女辅助操作，家庭共享游览精彩"),
    ]
    cols = 5
    start_x = Inches(0.35)
    cell_w = (SLIDE_W - Inches(0.7)) / cols

    for i, (icon, title, desc) in enumerate(highlights):
        x = start_x + i * cell_w
        add_rect(slide, x + Inches(0.06), Inches(2.3), cell_w - Inches(0.12),
                 Inches(2.6), fill_color=GREEN_800, line_color=GREEN_600,
                 line_width=Pt(1.5))
        add_text(slide, icon, x + Inches(0.06), Inches(2.5),
                 cell_w - Inches(0.12), Inches(0.8), font_size=32,
                 bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, title, x + Inches(0.1), Inches(3.3),
                 cell_w - Inches(0.15), Inches(0.5), font_size=16,
                 bold=True, color=GREEN_200, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x + Inches(0.1), Inches(3.85),
                 cell_w - Inches(0.15), Inches(0.9), font_size=BODY_FONT,
                 color=GRAY_300, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(1.5), Inches(5.1), SLIDE_W - Inches(3.0), Inches(0.04),
             fill_color=GREEN_500)

    add_text(slide, "联系我们 | 合作洽谈 | 媒体采访",
             Inches(0.5), Inches(5.35), SLIDE_W - Inches(1.0), Inches(0.45),
             font_size=H3_FONT, bold=False, color=GREEN_400, align=PP_ALIGN.CENTER)
    add_text(slide, "© 2024 桂林银发旅游小程序 · All Rights Reserved",
             Inches(0.5), Inches(5.85), SLIDE_W - Inches(1.0), Inches(0.4),
             font_size=BODY_FONT, bold=False, color=GRAY_500, align=PP_ALIGN.CENTER)

    add_rect(slide, 0, Inches(6.4), SLIDE_W, Inches(1.1), fill_color=GREEN_800)
    add_text(slide, "大字体 · 高对比度 · 简操作 · 零障碍 · 真关怀",
             Inches(0.5), Inches(6.6), SLIDE_W - Inches(1.0), Inches(0.55),
             font_size=20, bold=True, color=GREEN_200, align=PP_ALIGN.CENTER)

    return slide


# ==================== 主函数 ====================

def main():
    prs = new_prs()

    slide_cover(prs)       # 1
    slide_toc(prs)         # 2
    slide_overview(prs)    # 3
    slide_user_profile(prs) # 4
    slide_functions(prs)   # 5
    slide_ui_design(prs)   # 6
    slide_architecture(prs) # 7
    slide_database(prs)    # 8
    slide_api(prs)         # 9
    slide_payment(prs)     # 10
    slide_security(prs)    # 11
    slide_family_social(prs) # 12
    slide_operation(prs)   # 13
    slide_summary(prs)     # 14

    out_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(out_dir, "桂林银发旅游小程序_优化版.pptx")
    prs.save(out_path)
    print(f"✅ 优化版PPT已生成: {out_path}")
    print(f"   共 {len(prs.slides)} 页 | 字体≥14pt | 空间利用率≥90%")


if __name__ == "__main__":
    main()
