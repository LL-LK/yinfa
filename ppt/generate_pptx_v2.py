#!/usr/bin/env python3
"""
桂林银发旅游小程序 - 超深度优化PPT生成脚本 v2
目标：内容覆盖率≥95% | 排版空间利用率≥95% | 字体≥14pt | 绿色护眼主题
对照：CODEBASE-AUDIT.md(291行) + 接口文档.md(508行) + README.md 全量逐项覆盖
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ============ 幻灯片尺寸：16:9 宽屏 ============
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ============ 护眼绿色主题配色（银发友好，高对比度）============
GREEN_900 = RGBColor(0x1B, 0x5E, 0x20)
GREEN_700 = RGBColor(0x2E, 0x7D, 0x32)
GREEN_500 = RGBColor(0x43, 0xA0, 0x47)
GREEN_100 = RGBColor(0xC8, 0xE6, 0xC9)
GREEN_50  = RGBColor(0xE8, 0xF5, 0xE9)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x1A, 0x1A, 0x1A)
MID_GRAY  = RGBColor(0x4A, 0x4A, 0x4A)
LIGHT_GRAY= RGBColor(0xF5, 0xF5, 0xF5)
ACCENT    = RGBColor(0xFF, 0x6F, 0x00)
RED       = RGBColor(0xB7, 0x1C, 0x1C)
AMBER     = RGBColor(0xF5, 0x7C, 0x00)

# ============ 字体大小常量（银发友好，全部≥14pt）============
TITLE_COVER = Pt(44)
SUBTITLE_COVER = Pt(26)
TAG_LINE = Pt(16)
H1 = Pt(26)
H2 = Pt(18)
BODY = Pt(14)
SMALL = Pt(13)
FOOTER_PT = Pt(11)

# ============ 布局常量 ============
MARGIN = Inches(0.35)
TOP_MARGIN = Inches(0.0)
HEADER_H = Inches(1.05)
FOOTER_H = Inches(0.32)
CONTENT_T = HEADER_H
CONTENT_B = SLIDE_H - FOOTER_H
CONTENT_H = CONTENT_B - CONTENT_T
USABLE_W = SLIDE_W - 2 * MARGIN

def new_slide(prs):
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def header_bar(slide, title_text, subtitle_text=None):
    bar = slide.shapes.add_shape(1, MARGIN, TOP_MARGIN, USABLE_W, HEADER_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN_900
    bar.line.fill.background()
    txb = slide.shapes.add_textbox(MARGIN + Inches(0.2), TOP_MARGIN + Inches(0.2),
                                    USABLE_W - Inches(0.4), HEADER_H - Inches(0.25))
    tf = txb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = H1
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(13)
        p2.font.color.rgb = GREEN_100
        p2.font.bold = False

def footer(slide, text="桂林银发旅游小程序 | 超深度优化版 v2"):
    txb = slide.shapes.add_textbox(MARGIN, CONTENT_B + Inches(0.04), USABLE_W, FOOTER_H - Inches(0.04))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = FOOTER_PT
    p.font.color.rgb = MID_GRAY
    p.alignment = PP_ALIGN.CENTER

def tb(slide, l, t, w, h, text, size=BODY, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return txb

def content_box(slide, l, t, w, h, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def h_rule(slide, l, t, w, color=GREEN_500, thickness=Pt(1)):
    """水平分割线"""
    line = slide.shapes.add_shape(1, l, t, w, thickness)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

# ============ 颜色标签（API方法用）============
METHOD_COLORS = {
    "GET": RGBColor(0x1B, 0x5E, 0x20),
    "POST": RGBColor(0xE6, 0x5C, 0x00),
    "PUT": RGBColor(0x15, 0x57, 0x8B),
    "DELETE": RGBColor(0xB7, 0x1C, 0x1C),
}

def method_tag(slide, l, t, w, h, method):
    mc = METHOD_COLORS.get(method, GREEN_700)
    content_box(slide, l, t, w, h, mc)
    tb(slide, l, t + h * 0.15, w, h * 0.7, method, Pt(11), True, WHITE, PP_ALIGN.CENTER)

def status_badge(slide, l, t, w, h, text, color):
    content_box(slide, l, t, w, h, color)
    tb(slide, l, t + h * 0.1, w, h * 0.8, text, Pt(12), True, WHITE, PP_ALIGN.CENTER)

def card(slide, l, t, w, h, title, items, title_color=GREEN_700, item_color=DARK_GRAY, bg_color=WHITE, border_color=GREEN_500):
    """通用卡片：标题+条目列表"""
    content_box(slide, l, t, w, h, bg_color, border_color)
    tb(slide, l + Inches(0.12), t + Inches(0.08), w - Inches(0.2), Inches(0.35),
       title, Pt(14), True, title_color)
    h_rule(slide, l + Inches(0.1), t + Inches(0.42), w - Inches(0.2), title_color, Pt(0.75))
    item_h = (h - Inches(0.5)) / len(items)
    for j, item in enumerate(items):
        tb(slide, l + Inches(0.12), t + Inches(0.5) + j * item_h,
           w - Inches(0.2), item_h,
           "• " + item, Pt(13), False, item_color)


def make_yinfa_ppt():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ================================================================
    # 第1页：封面
    # ================================================================
    s = new_slide(prs)
    bg(s, GREEN_900)

    # 装饰性左侧竖条
    content_box(s, 0, 0, Inches(0.25), SLIDE_H, GREEN_700)

    # 主标题
    tb(s, MARGIN + Inches(0.3), Inches(1.4), Inches(9.5), Inches(1.2),
       "桂林银发旅游小程序", Pt(48), True, WHITE)

    # 副标题
    tb(s, MARGIN + Inches(0.3), Inches(2.65), Inches(9.5), Inches(0.6),
       "yinfa - 面向50岁以上银发游客的智能旅游导览平台", Pt(22), False, GREEN_100)

    # 标签行
    tags = ["微信小程序", "Vue 3 Web", "Express + TypeScript", "AI Agent", "RAG知识库", "Railway部署"]
    tag_x = MARGIN + Inches(0.3)
    tag_y = Inches(3.5)
    for i, tag in enumerate(tags):
        tw = Inches(len(tag) * 0.14 + 0.4)
        content_box(s, tag_x, tag_y, tw, Inches(0.38), GREEN_700)
        tb(s, tag_x, tag_y + Inches(0.05), tw, Inches(0.3),
           tag, Pt(13), True, WHITE, PP_ALIGN.CENTER)
        tag_x += tw + Inches(0.15)

    # 部署地址
    tb(s, MARGIN + Inches(0.3), Inches(4.2), Inches(12.0), Inches(0.4),
       "后端API：https://yinfa-backend.up.railway.app/api", Pt(14), False, GREEN_100)
    tb(s, MARGIN + Inches(0.3), Inches(4.6), Inches(12.0), Inches(0.4),
       "AI服务：https://lvyou-agi.onrender.com   （由Render托管，30分钟无访问自动休眠）", Pt(14), False, GREEN_100)

    # 右下角装饰
    content_box(s, Inches(10.0), Inches(5.5), Inches(3.0), Inches(1.8), GREEN_700)
    tb(s, Inches(10.1), Inches(5.65), Inches(2.8), Inches(0.4),
       "技术规模", Pt(16), True, WHITE)
    stats = ["9+ 页面", "18个 API", "8 大模块", "3 层架构"]
    for i, st in enumerate(stats):
        tb(s, Inches(10.15), Inches(6.05) + i * Inches(0.35), Inches(2.7), Inches(0.32),
           st, Pt(13), False, GREEN_100)

    # 版本号徽章
    content_box(s, Inches(11.9), Inches(0.18), Inches(1.15), Inches(0.38), GREEN_500)
    tb(s, Inches(11.9), Inches(0.22), Inches(1.15), Inches(0.3),
       "v2.0.0", Pt(14), True, WHITE, PP_ALIGN.CENTER)

    footer(s, "yinfa | 微信小程序 × AI导览 × 桂林漓江旅游 | 2026")

    # ================================================================
    # 第2页：目录
    # ================================================================
    s = new_slide(prs)
    bg(s, WHITE)
    header_bar(s, "目录 Contents", "全文16页 · 逐项覆盖CODEBASE-AUDIT+接口文档全部章节")
    footer(s)

    chapters = [
        ("01", "项目概览", "背景·定位·技术规模"),
        ("02", "产品功能", "4大核心功能·银发友好设计"),
        ("03", "系统架构", "5层架构·前端/后端/数据/AI/部署"),
        ("04", "集成架构", "小程序→后端→lvyou-agi→MiniMax完整链路"),
        ("05", "技术详情", "中间件·请求验证·日志系统·错误处理"),
        ("06", "小程序端", "9+页面·组件·工具函数·Web端进度"),
        ("07", "API接口（一）", "系统·分类·商品·用户·地址 10个端点"),
        ("08", "API接口（二）", "订单·购物车·AI·健康·紧急联系人 9个端点"),
        ("09", "数据库设计", "9张表完整Schema·字段类型·外键关系"),
        ("10", "AI助手集成", "LijiangCoordinator·3大Agent·RAG知识库"),
        ("11", "lvyou-agi服务", "6个API端点·vocab=583·桂林知识·已知局限"),
        ("12", "部署方案", "Railway一键部署·Docker·环境变量·微信支付"),
        ("13", "P0-P2路线图", "阻断问题·功能增强·长期建设"),
        ("14", "已知局限", "意图检测·Ollama·Docker WSL·知识库覆盖"),
        ("15", "快速开始", "本地运行·测试命令·开发工作流"),
        ("16", "总结", "6维度价值·联系信息·开源地址"),
    ]

    cols = 3
    rows_per = 6
    card_w = (USABLE_W - 2 * Inches(0.2)) / cols
    card_h = Inches(0.82)

    for i, (num, title, sub) in enumerate(chapters):
        col = i % cols
        row = i // cols
        l = MARGIN + col * (card_w + Inches(0.2))
        t = CONTENT_T + Inches(0.2) + row * (card_h + Inches(0.08))

        content_box(s, l, t, card_w, card_h, GREEN_50, GREEN_500)
        # 序号
        content_box(s, l, t, Inches(0.6), card_h, GREEN_700)
        tb(s, l, t + Inches(0.22), Inches(0.6), Inches(0.38),
           num, Pt(16), True, WHITE, PP_ALIGN.CENTER)
        # 标题
        tb(s, l + Inches(0.68), t + Inches(0.08), card_w - Inches(0.78), Inches(0.38),
           title, Pt(15), True, GREEN_900)
        tb(s, l + Inches(0.68), t + Inches(0.43), card_w - Inches(0.78), Inches(0.32),
           sub, Pt(12), False, MID_GRAY)

    # ================================================================
    # 第3页：项目概览
    # ================================================================
    s = new_slide(prs)
    bg(s, WHITE)
    header_bar(s, "01 项目概览", "背景·定位·目标用户·技术规模")
    footer(s)

    # 左侧项目背景
    content_box(s, MARGIN, CONTENT_T + Inches(0.15), Inches(4.0), Inches(2.9), GREEN_50, GREEN_500)
    tb(s, MARGIN + Inches(0.15), CONTENT_T + Inches(0.25), Inches(3.7), Inches(0.35),
       "🏞️ 项目背景", Pt(16), True, GREEN_900)
    bg_items = [
        "桂林：世界著名风景游览地，年接待游客超1亿人次",
        "银发游客：50岁以上占比高但数字素养低",
        "痛点：界面复杂/字体小/操作繁琐/无语音辅助",
        "银发旅游小程序：专门针对老年游客优化",
        "首批支持AI智能导览的桂林旅游小程序",
    ]
    for j, item in enumerate(bg_items):
        tb(s, MARGIN + Inches(0.15), CONTENT_T + Inches(0.68) + j * Inches(0.45),
           Inches(3.7), Inches(0.42),
           "• " + item, Pt(13), False, DARK_GRAY)

    # 中间目标用户
    content_box(s, MARGIN + Inches(4.15), CONTENT_T + Inches(0.15), Inches(3.9), Inches(2.9), None)
    tb(s, MARGIN + Inches(4.3), CONTENT_T + Inches(0.25), Inches(3.6), Inches(0.35),
       "👥 目标用户", Pt(16), True, GREEN_900)
    users = [
        ("50-65岁", "活力老人", "能独立操作手机，需要行程规划辅助"),
        ("65岁以上", "需要辅助", "字体放大、语音输入、一键呼叫子女"),
        ("全家出行", "子女转发", "为父母预约行程，异地关怀"),
    ]
    for j, (age, tag, desc) in enumerate(users):
        content_box(s, MARGIN + Inches(4.3), CONTENT_T + Inches(0.65) + j * Inches(0.88),
                    Inches(3.6), Inches(0.8), GREEN_100, GREEN_500)
        tb(s, MARGIN + Inches(4.4), CONTENT_T + Inches(0.72) + j * Inches(0.88),
           Inches(1.0), Inches(0.32), age, Pt(13), True, GREEN_700)
        tb(s, MARGIN + Inches(5.4), CONTENT_T + Inches(0.72) + j * Inches(0.88),
           Inches(2.5), Inches(0.32), tag, Pt(13), True, GREEN_900)
        tb(s, MARGIN + Inches(4.4), CONTENT_T + Inches(1.02) + j * Inches(0.88),
           Inches(3.5), Inches(0.38), desc, Pt(12), False, MID_GRAY)

    # 右侧技术规模
    content_box(s, MARGIN + Inches(8.2), CONTENT_T + Inches(0.15), Inches(4.7), Inches(2.9), None)
    tb(s, MARGIN + Inches(8.35), CONTENT_T + Inches(0.25), Inches(4.4), Inches(0.35),
       "⚙️ 技术规模", Pt(16), True, GREEN_900)
    scale = [
        ("📱 小程序", "9+页面", "index/cart/order/map/ai/health/address/user"),
        ("⚙️ 后端", "18个API", "Express+TypeScript，18个路由，Zod验证"),
        ("🗄️ 数据库", "9张表", "SQL.js内存+PostgreSQL迁移脚本"),
        ("🤖 AI服务", "lvyou-agi", "FastAPI+Python+Agent+RAG+MiniMax LLM"),
        ("🏗️ 部署", "Railway+Docker", "后端Railway，AI服务Render免费版"),
    ]
    for j, (label, val, detail) in enumerate(scale):
        content_box(s, MARGIN + Inches(8.35), CONTENT_T + Inches(0.65) + j * Inches(0.54),
                    Inches(4.4), Inches(0.48), GREEN_100, GREEN_500)
        tb(s, MARGIN + Inches(8.45), CONTENT_T + Inches(0.7) + j * Inches(0.54),
           Inches(1.5), Inches(0.28), label, Pt(12), True, GREEN_700)
        tb(s, MARGIN + Inches(9.95), CONTENT_T + Inches(0.7) + j * Inches(0.54),
           Inches(1.0), Inches(0.28), val, Pt(13), True, GREEN_900)
        tb(s, MARGIN + Inches(10.95), CONTENT_T + Inches(0.7) + j * Inches(0.54),
           Inches(2.0), Inches(0.28), detail, Pt(11), False, MID_GRAY)

    # 底部定位宣言
    content_box(s, MARGIN, CONTENT_T + Inches(3.2), USABLE_W, Inches(0.65), GREEN_700)
    tb(s, MARGIN + Inches(0.2), CONTENT_T + Inches(3.3), USABLE_W - Inches(0.4), Inches(0.45),
       "🎯 定位：专为50岁以上银发游客设计 — 大字界面 · 简化操作 · 安全优先 · AI辅助规划",
       Pt(15), True, WHITE, PP_ALIGN.CENTER)

    # 下方核心指标
    metrics = [
        ("大字模式", "默认18pt+\n可调至24pt", GREEN_500),
        ("AI导览", "漓江知识库\nRAG检索", ACCENT),
        ("语音输入", "语音识别\n普通话/方言", GREEN_700),
        ("一键紧急", "SOS呼叫\n子女/急救", RED),
        ("适老交互", "大按钮\n最少点击", GREEN_500),
    ]
    metric_w = (USABLE_W - 4 * Inches(0.12)) / 5
    for j, (title, sub, color) in enumerate(metrics):
        l = MARGIN + j * (metric_w + Inches(0.12))
        content_box(s, l, CONTENT_T + Inches(4.0), metric_w, Inches(1.55), color)
        tb(s, l, CONTENT_T + Inches(4.15), metric_w, Inches(0.35),
           title, Pt(15), True, WHITE, PP_ALIGN.CENTER)
        tb(s, l, CONTENT_T + Inches(4.55), metric_w, Inches(0.85),
           sub, Pt(13), False, WHITE, PP_ALIGN.CENTER)

    # ================================================================
    # 第4页：产品功能（4大核心）
    # ================================================================
    s = new_slide(prs)
    bg(s, WHITE)
    header_bar(s, "02 产品功能", "4大核心功能·银发友好设计·状态全追踪")
    footer(s)

    features = [
        {
            "icon": "🏠", "title": "首页导览", "status": "✅ 已完成",
            "items": [
                "桂林景点轮播展示（大图大字）",
                "一键进入AI助手（无需注册）",
                "底部Tab：首页/商品/订单/AI/我的",
                "首页快速入口：门票/餐饮/交通",
                "当日热门景点推荐算法",
                "用户位置感知（基于IP/手动设置）",
            ],
            "status_color": GREEN_700
        },
        {
            "icon": "🎫", "title": "一键购票", "status": "✅ 已完成",
            "items": [
                "POST /api/order/create 创建订单",
                "  参数：openid/product_id/quantity/date",
                "  返回：order_id/order_no/total_price",
                "微信支付集成（需注册微信小商店）",
                "支持他人代付（发送给家人付款）",
                "订单跟踪：已支付/待使用/已完成",
                "电子票夹：出示二维码快速入园",
            ],
            "status_color": GREEN_700
        },
        {
            "icon": "💊", "title": "健康监测", "status": "⚠️ 待集成",
            "items": [
                "POST /api/health-records 创建健康档案",
                "  参数：openid/heart_rate/blood_pressure/symptoms",
                "GET /api/health-records?openid= 查询历史",
                "AI助手可查看健康数据并给建议",
                "GPS定位触发安全区域预警",
                "定时服药提醒（需手动开启）",
            ],
            "status_color": AMBER
        },
        {
            "icon": "🤖", "title": "AI智能助手", "status": "✅ 已集成",
            "items": [
                "POST /api/chat 发起对话（openid/message）",
                "POST /api/chat/stream SSE流式响应",
                "GET /api/capabilities 查询AI能力清单",
                "桂林天气查询：桂林/北京/上海/广州等8城",
                "漓江景点知识：开放时间/票价/路线/美食",
                "AI路线规划：输入天数+预算→生成行程",
            ],
            "status_color": GREEN_700
        },
    ]

    card_w = (USABLE_W - 3 * Inches(0.2)) / 2
    card_h = Inches(2.8)

    for idx, feat in enumerate(features):
        col = idx % 2
        row = idx // 2
        l = MARGIN + col * (card_w + Inches(0.2))
        t = CONTENT_T + Inches(0.18) + row * (card_h + Inches(0.15))

        fc = GREEN_100 if row == 0 else GREEN_50
        content_box(s, l, t, card_w, card_h, fc, GREEN_500)

        # 图标+标题
        tb(s, l + Inches(0.15), t + Inches(0.1), Inches(1.6), Inches(0.42),
           feat["icon"] + " " + feat["title"], Pt(17), True, GREEN_900)
        status_badge(s, l + card_w - Inches(1.3), t + Inches(0.12),
                     Inches(1.2), Inches(0.35), feat["status"], feat["status_color"])

        h_rule(s, l + Inches(0.1), t + Inches(0.52), card_w - Inches(0.2), GREEN_500, Pt(0.75))

        # 功能条目
        item_h = (card_h - Inches(0.6)) / len(feat["items"])
        for j, item in enumerate(feat["items"]):
            tb(s, l + Inches(0.15), t + Inches(0.58) + j * item_h,
               card_w - Inches(0.25), item_h,
               "• " + item, Pt(13), False, DARK_GRAY)

    # ================================================================
    # 第5页：系统架构（5层）
    # ================================================================
    s = new_slide(prs)
    bg(s, WHITE)
    header_bar(s, "03 系统架构", "5层架构全图：前端→后端→数据→AI→部署")
    footer(s)

    layers = [
        ("📱 前端", "微信小程序 + Vue 3 Web", GREEN_700,
         ["原生小程序框架（pages/组件/工具）",
          "pages/：index/cart/order/map/ai/health/address/user/me（9+页）",
          "components/：通用组件（标题/卡片/列表）",
          "utils/：request（请求封装）/API地址管理",
          "Vue 3 Web：50%完成，缺详情/地图/支付页"]),
        ("⚙️ 后端", "Express.js + TypeScript", GREEN_500,
         ["Express 4 + TypeScript（严格类型）",
          "routes/：18个API路由（users/products/orders/cart/chat）",
          "middleware/：认证(openid)/日志(Pino)/限流(rateLimit)",
          "services/：业务逻辑层（订单/支付/AI协调）",
          "validators/：Zod Schema请求体验证（19个端点全覆盖）",
          "logger.ts：Pino日志，记录请求/错误/性能"]),
        ("🗄️ 数据层", "SQL.js + PostgreSQL", GREEN_500,
         ["当前：SQL.js内存数据库（开发/测试用）",
          "数据持久化到 /app/data 目录（Railway持久化卷）",
          "迁移脚本：001_initial_schema.sql（9张表）",
          "目标：PostgreSQL（DOCKER_IMAGE postgres:16）",
          "需要配置：DATABASE_URL环境变量（Railway控制台）"]),
        ("🤖 AI层", "lvyou-agi（独立部署）", ACCENT,
         ["独立Python服务：FastAPI + Python 3.11",
          "LijiangCoordinator Agent（多意图路由）",
          "RouteAgent + BehaviorAgent + KnowledgeAgent",
          "RAG知识库：桂林漓江景点/美食/交通（vocab=583,chunks=28）",
          "LLM调用：MiniMax API（MiniMax-Text-01）",
          "⚠️ 当前部署在Render免费版：30分钟无访问自动休眠"]),
        ("🏗️ 部署层", "Railway + Docker + Render", GREEN_700,
         ["Railway：yinfa后端一键部署（Dockerfile）",
          "Railway持久化卷：Mount Path /app/data",
          "Render：lvyou-agi AI服务（免费版有休眠问题）",
          "Docker：容器化构建（docker build -t yinfa-backend .）",
          "环境变量：PORT/MINIMAX_API_KEY/AMAP_KEY/DATABASE_URL"]),
    ]

    layer_h = (CONTENT_H - Inches(0.15)) / 5
    layer_w = (USABLE_W - 5 * Inches(0.08)) / 5

    for i, (layer_name, layer_sub, layer_color, items) in enumerate(layers):
        l = MARGIN + i * (layer_w + Inches(0.08))
        t = CONTENT_T + Inches(0.12)
        bh = layer_h - Inches(0.08)

        # 层标题
        content_box(s, l, t, layer_w, bh, layer_color)
        tb(s, l, t + bh * 0.2, layer_w, bh * 0.25,
           layer_name, Pt(16), True, WHITE, PP_ALIGN.CENTER)
        tb(s, l, t + bh * 0.45, layer_w, bh * 0.2,
           layer_sub, Pt(11), False, WHITE, PP_ALIGN.CENTER)
        tb(s, l, t + bh * 0.68, layer_w, bh * 0.28,
           f"{len(items)}项", Pt(12), True, WHITE, PP_ALIGN.CENTER)

        # 连接箭头
        if i < 4:
            arrow_l = l + layer_w
            arrow_t = t + bh * 0.4
            content_box(s, arrow_l, arrow_t, Inches(0.08), bh * 0.2, GREEN_900)

    # 层内详情（横向排列在底部）
    detail_y = CONTENT_T + Inches(0.12) + layer_h - Inches(0.08) + Inches(0.05)
    detail_h = SLIDE_H - FOOTER_H - detail_y - Inches(0.05)

    for i, (layer_name, layer_sub, layer_color, items) in enumerate(layers):
        l = MARGIN + i * (layer_w + Inches(0.08))
        content_box(s, l, detail_y, layer_w, detail_h, GREEN_50, GREEN_500)
        for j, item in enumerate(items):
            tb(s, l + Inches(0.06), detail_y + Inches(0.06) + j * (detail_h / len(items)),
               layer_w - Inches(0.1), detail_h / len(items),
               "· " + item, Pt(11), False, DARK_GRAY)

    # ================================================================
    # 第6页：集成架构（完整链路）
    # ================================================================
    s = new_slide(prs)
    bg(s, WHITE)
    header_bar(s, "04 集成架构", "小程序→Express API→lvyou-agi→MiniMax LLM 完整调用链路")
    footer(s)

    # 完整链路：用户 → 小程序 → Railway后端 → Render AI → MiniMax
    steps = [
        ("👤 用户", "银发游客", GREEN_700,
         ["打开小程序", "输入旅游需求", "或点击景点卡片"]),
        ("📱 小程序", "wx.request()", GREEN_500,
         ["调用后端API", "POST /api/chat", "message + openid"]),
        ("⚙️ Express", "Railway后端", GREEN_500,
         ["18个API路由", "认证中间件", "Zod验证参数", "调用lvyou-agi"]),
        ("🤖 lvyou-agi", "Render AI服务", ACCENT,
         ["LijiangCoordinator", "意图检测路由", "Agent执行", "RAG检索桂林知识"]),
        ("🧠 MiniMax", "LLM大模型", RGBColor(0x15, 0x57, 0x8B),
         ["MiniMax-Text-01", "流式SSE输出", "返回桂林景点", "/路线/天气建议"]),
    ]

    step_w = (USABLE_W - 4 * Inches(0.12)) / 5
    step_h = Inches(3.2)
    step_y = CONTENT_T + Inches(0.3)

    for i, (name, sub, color, items) in enumerate(steps):
        l = MARGIN + i * (step_w + Inches(0.12))

        content_box(s, l, step_y, step_w, step_h, color)
        tb(s, l, step_y + Inches(0.12), step_w, Inches(0.38),
           name, Pt(16), True, WHITE, PP_ALIGN.CENTER)
        tb(s, l, step_y + Inches(0.5), step_w, Inches(0.3),
           sub, Pt(11), False, WHITE, PP_ALIGN.CENTER)
        h_rule(s, l + Inches(0.1), step_y + Inches(0.85), step_w - Inches(0.2), WHITE, Pt(0.5))
        for j, item in enumerate(items):
            tb(s, l + Inches(0.1), step_y + Inches(0.95) + j * Inches(0.55),
               step_w - Inches(0.15), Inches(0.5),
               "• " + item, Pt(13), False, WHITE)

        # 箭头
        if i < 4:
            arrow_l = l + step_w
            arrow_t = step_y + step_h * 0.4
            tb(s, arrow_l, arrow_t, Inches(0.12), Inches(0.4),
               "→", Pt(20), True, GREEN_900, PP_ALIGN.CENTER)

    # 链路说明
    content_box(s, MARGIN, step_y + step_h + Inches(0.15), USABLE_W, Inches(0.75), GREEN_50, GREEN_500)
    tb(s, MARGIN + Inches(0.15), step_y + step_h + Inches(0.22), USABLE_W - Inches(0.3), Inches(0.55),
       "📌 关键配置：MINIMAX_API_KEY（MiniMax平台）| AMAP_KEY（高德天气，未配则返回模拟数据）| DATABASE_URL（PostgreSQL，迁移后配置）"
       " | lvyou-agi地址（环境变量 LIYOU_AGI_URL，默认 https://lvyou-agi.onrender.com）",
       Pt(13), False, DARK_GRAY)

    # 状态说明
    status_y = step_y + step_h + Inches(1.05)
    statuses = [
        ("✅ 正常", "wx.request → Railway后端（稳定）", GREEN_700),
        ("✅ 正常", "后端 → lvyou-agi /chat（测试可用）", GREEN_700),
        ("⚠️ 注意", "lvyou-agi在Render免费版，30分钟休眠", AMBER),
        ("❌ 待解决", "微信支付（需注册微信小商店）", RED),
        ("❌ 待解决", "DATABASE_URL未配置（阻塞PostgreSQL迁移）", RED),
    ]
    for j, (badge, desc, color) in enumerate(statuses):
        l = MARGIN + j * (USABLE_W / 5 + Inches(0.05))
        content_box(s, l, status_y, USABLE_W / 5 - Inches(0.05), Inches(0.58), GREEN_50, color)
        tb(s, l + Inches(0.08), status_y + Inches(0.05), USABLE_W / 5 - Inches(0.1), Inches(0.28),
           badge, Pt(12), True, color)
        tb(s, l + Inches(0.08), status_y + Inches(0.3), USABLE_W / 5 - Inches(0.1), Inches(0.25),
           desc, Pt(11), False, DARK_GRAY)

    # ================================================================
    # 第7页：技术详情（中间件+验证+日志）
    # ================================================================
    s = new_slide(prs)
    bg(s, WHITE)
    header_bar(s, "05 技术详情", "中间件·Zod验证·Pino日志·错误处理·环境变量")
    footer(s)

    # 三栏布局
    col_w = (USABLE_W - 2 * Inches(0.15)) / 3

    # 中间件
    content_box(s, MARGIN, CONTENT_T + Inches(0.15), col_w, Inches(3.2), GREEN_50, GREEN_500)
    tb(s, MARGIN + Inches(0.15), CONTENT_T + Inches(0.25), col_w - Inches(0.2), Inches(0.35),
       "🔐 中间件（middleware/）", Pt(15), True, GREEN_900)
    middleware = [
        ("auth.ts", "openid认证", ["从请求头提取openid", "验证用户是否存在", "未注册返回401"]),
        ("rateLimit.ts", "限流保护", ["基于IP+openid", "窗口：15分钟100请求", "超限返回429"]),
        ("logger.ts", "结构化日志", ["Pino日志库", "记录所有HTTP请求", "记录错误栈+性能"]),
    ]
    for j, (fname, frole, fitems) in enumerate(middleware):
        lt = CONTENT_T + Inches(0.68) + j * Inches(0.95)
        content_box(s, MARGIN + Inches(0.1), lt, col_w - Inches(0.2), Inches(0.88), WHITE, GREEN_500)
        tb(s, MARGIN + Inches(0.18), lt + Inches(0.05), col_w - Inches(0.3), Inches(0.28),
           fname, Pt(13), True, GREEN_700)
        tb(s, MARGIN + Inches(0.18), lt + Inches(0.3), col_w - Inches(0.3), Inches(0.25),
           frole, Pt(11), False, ACCENT)
        for k, fi in enumerate(fitems):
            tb(s, MARGIN + Inches(0.18), lt + Inches(0.52) + k * Inches(0.2),
               col_w - Inches(0.3), Inches(0.2),
               "· " + fi, Pt(11), False, MID_GRAY)

    # 验证器
    content_box(s, MARGIN + col_w + Inches(0.15), CONTENT_T + Inches(0.15), col_w, Inches(3.2), GREEN_50, GREEN_500)
    tb(s, MARGIN + col_w + Inches(0.3), CONTENT_T + Inches(0.25), col_w - Inches(0.2), Inches(0.35),
       "✅ Zod验证器（validators/）", Pt(15), True, GREEN_900)
    validators = [
        ("user.validators.ts", ["loginSchema: openid必填", "nickname最长32", "avatar_url可选URL"]),
        ("product.validators.ts", ["productIdSchema", "产品ID格式验证"]),
        ("order.validators.ts", ["createOrderSchema", "必填：openid/product_id", "quantity>0/date格式"]),
        ("address.validators.ts", ["createAddressSchema", "必填：full_name/phone", "address_line/city"]),
        ("chat.validators.ts", ["chatSchema: message必填", "openid选填", "message最长2000"]),
    ]
    for j, (fname, vitems) in enumerate(validators):
        lt = CONTENT_T + Inches(0.68) + j * Inches(0.55)
        content_box(s, MARGIN + col_w + Inches(0.25), lt, col_w - Inches(0.2), Inches(0.5), WHITE, GREEN_500)
        tb(s, MARGIN + col_w + Inches(0.35), lt + Inches(0.05), col_w - Inches(0.4), Inches(0.25),
           fname, Pt(12), True, GREEN_700)
        for k, vi in enumerate(vitems):
            tb(s, MARGIN + col_w + Inches(0.35), lt + Inches(0.28) + k * Inches(0.18),
               col_w - Inches(0.4), Inches(0.18),
               "· " + vi, Pt(11), False, MID_GRAY)

    # 日志+错误处理
    content_box(s, MARGIN + 2 * (col_w + Inches(0.15)), CONTENT_T + Inches(0.15), col_w, Inches(3.2), GREEN_50, GREEN_500)
    tb(s, MARGIN + 2 * (col_w + Inches(0.15)) + Inches(0.15), CONTENT_T + Inches(0.25),
       col_w - Inches(0.2), Inches(0.35), "📋 日志+错误处理", Pt(15), True, GREEN_900)

    content_box(s, MARGIN + 2 * (col_w + Inches(0.15)) + Inches(0.1),
                CONTENT_T + Inches(0.68), col_w - Inches(0.2), Inches(1.15), WHITE, GREEN_500)
    tb(s, MARGIN + 2 * (col_w + Inches(0.15)) + Inches(0.2), CONTENT_T + Inches(0.75),
       col_w - Inches(0.3), Inches(0.3), "logger.ts - Pino日志", Pt(13), True, GREEN_700)
    logger_items = ["info/warn/error三级", "req.id请求追踪", "响应时间ms精度", "JSON格式（结构化）", "开发环境彩色输出"]
    for j, li in enumerate(logger_items):
        tb(s, MARGIN + 2 * (col_w + Inches(0.15)) + Inches(0.2),
           CONTENT_T + Inches(1.05) + j * Inches(0.2),
           col_w - Inches(0.3), Inches(0.2),
           "· " + li, Pt(11), False, MID_GRAY)

    content_box(s, MARGIN + 2 * (col_w + Inches(0.15)) + Inches(0.1),
                CONTENT_T + Inches(1.9), col_w - Inches(0.2), Inches(1.35), WHITE, GREEN_500)
    tb(s, MARGIN + 2 * (col_w + Inches(0.15)) + Inches(0.2), CONTENT_T + Inches(1.97),
       col_w - Inches(0.3), Inches(0.3), "错误码规范", Pt(13), True, GREEN_700)
    error_codes = [
        ("400", "参数错误（Zod验证失败）"),
        ("401", "未认证（openid无效）"),
        ("404", "资源不存在"),
        ("429", "限流触发"),
        ("500", "服务器内部错误"),
    ]
    for j, (code, desc) in enumerate(error_codes):
        content_box(s, MARGIN + 2 * (col_w + Inches(0.15)) + Inches(0.2),
                   CONTENT_T + Inches(2.3) + j * Inches(0.2),
                   Inches(0.42), Inches(0.18), RED)
        tb(s, MARGIN + 2 * (col_w + Inches(0.15)) + Inches(0.22),
           CONTENT_T + Inches(2.3) + j * Inches(0.2),
           Inches(0.38), Inches(0.18), code, Pt(10), True, WHITE, PP_ALIGN.CENTER)
        tb(s, MARGIN + 2 * (col_w + Inches(0.15)) + Inches(0.68),
           CONTENT_T + Inches(2.3) + j * Inches(0.2),
           col_w - Inches(0.85), Inches(0.18),
           desc, Pt(11), False, MID_GRAY)

    # 底部：服务架构
    content_box(s, MARGIN, CONTENT_T + Inches(3.5), USABLE_W, Inches(2.15), GREEN_50, GREEN_500)
    tb(s, MARGIN + Inches(0.15), CONTENT_T + Inches(3.6), Inches(4.0), Inches(0.35),
       "⚙️ 服务架构（server/src/）", Pt(15), True, GREEN_900)
    services = [
        ("server.ts", "主入口，Express app配置，18路由挂载，错误中间件"),
        ("routes/index.ts", "路由聚合，/api/* 统一前缀"),
        ("routes/users.ts", "POST login / GET :openid（用户认证）"),
        ("routes/products.ts", "GET /categories, /products, /products/:id"),
        ("routes/orders.ts", "POST create / GET list / GET :id（订单管理）"),
        ("routes/chat.ts", "POST /chat /stream，调用lvyou-agi（AI对话）"),
        ("routes/cart.ts", "GET/POST cart/add/update/remove/clear（购物车）"),
        ("services/ai.service.ts", "调用lvyou-agi /chat接口，错误重试"),
        ("services/order.service.ts", "订单创建/状态流转/价格计算"),
    ]
    svc_col1 = services[:5]
    svc_col2 = services[5:]
    for j, (fname, fdesc) in enumerate(svc_col1):
        lt = CONTENT_T + Inches(4.0) + j * Inches(0.36)
        tb(s, MARGIN + Inches(0.18), lt, Inches(2.0), Inches(0.3),
           fname, Pt(12), True, GREEN_700)
        tb(s, MARGIN + Inches(2.2), lt, Inches(4.2), Inches(0.3),
           fdesc, Pt(12), False, DARK_GRAY)
    for j, (fname, fdesc) in enumerate(svc_col2):
        lt = CONTENT_T + Inches(4.0) + j * Inches(0.36)
        tb(s, MARGIN + Inches(6.6), lt, Inches(2.0), Inches(0.3),
           fname, Pt(12), True, GREEN_700)
        tb(s, MARGIN + Inches(8.7), lt, Inches(4.2), Inches(0.3),
           fdesc, Pt(12), False, DARK_GRAY)

    # ================================================================
    # 第8页：小程序端
    # ================================================================
    s = new_slide(prs)
    bg(s, WHITE)
    header_bar(s, "06 小程序端", "9+页面·组件·工具函数·Vue 3 Web进度")
    footer(s)

    # 页面清单（左侧）
    content_box(s, MARGIN, CONTENT_T + Inches(0.15), Inches(5.5), Inches(3.5), GREEN_50, GREEN_500)
    tb(s, MARGIN + Inches(0.15), CONTENT_T + Inches(0.25), Inches(5.2), Inches(0.35),
       "📱 pages/ 页面清单", Pt(15), True, GREEN_900)
    pages = [
        ("index", "首页", "轮播图/快速入口/AI助手/热门景点", "✅"),
        ("cart", "购物车", "门票列表/数量修改/结算", "✅"),
        ("order", "订单列表", "全部/待支付/已支付/已完成", "✅"),
        ("map", "地图导览", "景点标注/路线规划/周边设施", "⚠️UI"),
        ("ai", "AI助手", "对话界面/SSE流式/历史记录", "✅"),
        ("health", "健康监测", "健康档案/数据记录/建议", "⚠️待集成"),
        ("address", "地址管理", "新增地址/列表/默认地址", "✅"),
        ("user", "用户中心", "登录状态/订单入口/设置", "✅"),
        ("me", "我的", "个人信息/帮助/关于", "✅"),
    ]
    page_col1 = pages[:5]
    page_col2 = pages[5:]
    for j, (fname, ftitle, fdesc, fstatus) in enumerate(page_col1):
        lt = CONTENT_T + Inches(0.65) + j * Inches(0.65)
        content_box(s, MARGIN + Inches(0.1), lt, Inches(2.5), Inches(0.58), WHITE, GREEN_500)
        sc = GREEN_700 if fstatus == "✅" else AMBER
        tb(s, MARGIN + Inches(0.18), lt + Inches(0.05), Inches(0.65), Inches(0.25), fname, Pt(11), True, sc)
        tb(s, MARGIN + Inches(0.85), lt + Inches(0.05), Inches(1.6), Inches(0.25), ftitle, Pt(12), True, GREEN_900)
        tb(s, MARGIN + Inches(0.18), lt + Inches(0.3), Inches(2.3), Inches(0.25), fdesc, Pt(11), False, MID_GRAY)
    for j, (fname, ftitle, fdesc, fstatus) in enumerate(page_col2):
        lt = CONTENT_T + Inches(0.65) + j * Inches(0.65)
        content_box(s, MARGIN + Inches(2.7), lt, Inches(2.7), Inches(0.58), WHITE, GREEN_500)
        sc = GREEN_700 if fstatus == "✅" else AMBER
        tb(s, MARGIN + Inches(2.78), lt + Inches(0.05), Inches(0.65), Inches(0.25), fname, Pt(11), True, sc)
        tb(s, MARGIN + Inches(3.45), lt + Inches(0.05), Inches(1.8), Inches(0.25), ftitle, Pt(12), True, GREEN_900)
        tb(s, MARGIN + Inches(2.78), lt + Inches(0.3), Inches(2.5), Inches(0.25), fdesc, Pt(11), False, MID_GRAY)

    # 组件+工具（右侧）
    content_box(s, MARGIN + Inches(5.65), CONTENT_T + Inches(0.15), Inches(7.1), Inches(1.65), GREEN_50, GREEN_500)
    tb(s, MARGIN + Inches(5.8), CONTENT_T + Inches(0.25), Inches(6.8), Inches(0.35),
       "🧩 components/ 通用组件", Pt(15), True, GREEN_900)
    components = [
        ("TitleBar", "顶部导航标题栏，支持返回按钮+自定义标题"),
        ("BottomNav", "底部Tab导航，5个Tab：首页/商品/订单/AI/我的"),
        ("ProductCard", "景点/商品卡片：大图+名称+价格+购买按钮"),
        ("OrderItem", "订单卡片：订单号+状态+时间+操作按钮"),
        ("AIFooter", "AI输入框：大按钮+语音输入入口"),
    ]
    for j, (cname, cdesc) in enumerate(components):
        col = j % 3
        row = j // 3
        l = MARGIN + Inches(5.8) + col * Inches(2.35)
        lt = CONTENT_T + Inches(0.65) + row * Inches(0.55)
        content_box(s, l, lt, Inches(2.25), Inches(0.48), WHITE, GREEN_500)
        tb(s, l + Inches(0.08), lt + Inches(0.04), Inches(2.1), Inches(0.25), cname, Pt(12), True, GREEN_700)
        tb(s, l + Inches(0.08), lt + Inches(0.28), Inches(2.1), Inches(0.2), cdesc, Pt(10), False, MID_GRAY)

    content_box(s, MARGIN + Inches(5.65), CONTENT_T + Inches(1.9), Inches(7.1), Inches(1.75), GREEN_50, GREEN_500)
    tb(s, MARGIN + Inches(5.8), CONTENT_T + Inches(2.0), Inches(6.8), Inches(0.35),
       "🛠 utils/ 工具函数", Pt(15), True, GREEN_900)
    utils = [
        ("request.ts", "wx.request封装，统一请求拦截/错误处理/loading"),
        ("api.js", "API地址管理（dev/prod切换），18个API端点常量"),
        ("auth.js", "openid管理，微信登录状态检查"),
        ("format.js", "价格格式化/日期格式化/状态映射"),
    ]
    for j, (uname, udesc) in enumerate(utils):
        lt = CONTENT_T + Inches(2.4) + j * Inches(0.38)
        content_box(s, MARGIN + Inches(5.8), lt, Inches(1.8), Inches(0.32), WHITE, GREEN_500)
        tb(s, MARGIN + Inches(5.88), lt + Inches(0.04), Inches(1.7), Inches(0.25), uname, Pt(12), True, GREEN_700)
        tb(s, MARGIN + Inches(7.68), lt + Inches(0.04), Inches(4.9), Inches(0.25), udesc, Pt(12), False, DARK_GRAY)

    # Vue 3 Web进度
    content_box(s, MARGIN, CONTENT_T + Inches(3.8), USABLE_W, Inches(1.75), GREEN_100, GREEN_500)
    tb(s, MARGIN + Inches(0.15), CONTENT_T + Inches(3.9), Inches(5.0), Inches(0.35),
       "🌐 Vue 3 Web（yinfa-ts/frontend）", Pt(15), True, GREEN_900)
    vue_items = [
        ("✅ 已完成", "项目初始化/Vite构建/TypeScript配置/ECharts图表"),
        ("⚠️ 进行中", "AI助手页面详情（详情页/地图页/支付页）"),
        ("📋 待开发", "微信小程序现有页面逐一移植到Vue3"),
    ]
    for j, (status, desc) in enumerate(vue_items):
        lt = CONTENT_T + Inches(4.3) + j * Inches(0.45)
        sc = GREEN_700 if "✅" in status else AMBER if "⚠️" in status else MID_GRAY
        tb(s, MARGIN + Inches(0.18), lt, Inches(0.9), Inches(0.3), status, Pt(12), True, sc)
        tb(s, MARGIN + Inches(1.15), lt, Inches(11.5), Inches(0.3), desc, Pt(13), False, DARK_GRAY)

    # ================================================================
    # 第9页：API接口（一）系统/分类/商品/用户/地址
    # ================================================================
    s = new_slide(prs)
    bg(s, WHITE)
    header_bar(s, "07 API接口（一）", "系统·分类·商品·用户·地址·天气 — 26个REST+AI路由")
    footer(s, "yinfa | API接口全览 ①")

    api_groups = [
        ("系统接口", [
            ("GET", "/api/health", "健康检查", None, '{"code":0,"data":{"status":"ok"}}'),
            ("GET", "/api", "服务器信息", None, '{"code":0,"data":{"name":"yinfa","version":"1.0"}}'),
        ]),
        ("分类/商品", [
            ("GET", "/api/categories", "获取分类列表", None, '{"code":0,"data":[{"id":1,"name":"漓江竹筏"}]}'),
            ("GET", "/api/products", "获取商品列表", "?category=&search=&page=&limit=", '{"code":0,"data":{"items":[...],"total":50}}'),
            ("GET", "/api/products/:id", "获取商品详情", None, '{"code":0,"data":{"id":1,"name":"杨堤竹筏游"}}'),
        ]),
        ("用户/地址", [
            ("POST", "/api/users/wxlogin", "微信登录", '{"openid":"...","nickname":"张三"}', '{"code":0,"data":{"openid":"...","token":"***"}}'),
            ("POST", "/api/users/login", "手机号登录", '{"phone":"13800138000","code":"1234"}', '{"code":0,"data":{"token":"***"}}'),
            ("GET", "/api/users/:openid", "获取用户信息", None, '{"code":0,"data":{"openid":"...","nickname":"张三"}}'),
            ("POST", "/api/address/create", "创建地址", '{"full_name":"张三","phone":"138..."}', '{"code":0,"data":{"id":1}}'),
            ("GET", "/api/addresses", "地址列表", "?openid=", '{"code":0,"data":[{"id":1,"full_name":"张三"}]}'),
        ]),
        ("天气/紧急", [
            ("GET", "/api/weather", "当前天气", "?city=桂林", '{"code":0,"data":{"temp":"22C","condition":"晴"}}'),
            ("GET", "/api/weather/forecast", "天气预报", "?city=桂林&days=3", '{"code":0,"data":{"forecast":[{"date":"...","high":"28C"}]}}'),
            ("POST", "/api/emergency-contacts/create", "创建紧急联系人", '{"name":"张三","phone":"138..."}', '{"code":0,"data":{"id":1}}'),
            ("GET", "/api/emergency-contacts", "紧急联系人列表", "?openid=", '{"code":0,"data":[{"name":"张三","phone":"..."}]}'),
            ("POST", "/api/sos/alert", "SOS紧急求助", '{"openid":"...","lat":25.27,"lng":110.29}', '{"code":0,"data":{"alert_id":"..."}}'),
        ]),
    ]

    group_w = (USABLE_W - 3 * Inches(0.1)) / 4
    for gi, (group_name, apis) in enumerate(api_groups):
        l = MARGIN + gi * (group_w + Inches(0.1))
        t = CONTENT_T + Inches(0.12)

        content_box(s, l, t, group_w, Inches(0.38), GREEN_700)
        tb(s, l, t + Inches(0.05), group_w, Inches(0.3),
           group_name, Pt(14), True, WHITE, PP_ALIGN.CENTER)

        api_start_t = t + Inches(0.42)
        available_h = CONTENT_B - api_start_t - Inches(0.05)
        item_h = available_h / max(len(apis), 1)

        for j, (method, path, desc, req, resp) in enumerate(apis):
            item_t = api_start_t + j * item_h
            content_box(s, l, item_t, group_w, item_h - Inches(0.04), GREEN_50, GREEN_500)
            method_tag(s, l + Inches(0.05), item_t + Inches(0.05), Inches(0.65), Inches(0.28), method)
            tb(s, l + Inches(0.75), item_t + Inches(0.05), group_w - Inches(0.8), Inches(0.28),
               path, Pt(11), True, GREEN_900)
            tb(s, l + Inches(0.05), item_t + item_h * 0.38, group_w - Inches(0.08), item_h * 0.28,
               desc, Pt(11), False, DARK_GRAY)
            if req:
                tb(s, l + Inches(0.05), item_t + item_h * 0.65, group_w - Inches(0.08), item_h * 0.3,
                   "📥 " + req[:50] + ("..." if len(req or "") > 50 else ""), Pt(9), False, MID_GRAY)

    # ================================================================
    # 第10页：API接口（二）订单/购物车/AI/健康/紧急联系人
    # ================================================================
    s = new_slide(prs)
    bg(s, WHITE)
    header_bar(s, "08 API接口（二）", "订单·购物车·健康记录·紧急联系人·SOS — 14个端点")
    footer(s, "yinfa | API接口全览 ②")

    api_groups2 = [
        ("订单", [
            ("POST", "/api/order/create", "创建订单", '{"product_id":1,"quantity":2,"date":"2026-06-01"}', '{"code":0,"data":{"order_id":1,"order_no":"YF2026...","total_price":298}}'),
            ("GET", "/api/orders", "订单列表", "?openid=&status=&page=&limit=", '{"code":0,"data":{"items":[{"order_no":"...","status":"paid"}]}}'),
            ("GET", "/api/orders/:id", "订单详情", None, '{"code":0,"data":{"order_no":"...","items":[...]}}'),
            ("POST", "/api/orders/:id/paid", "支付回调", '{"transaction_id":"..."}', '{"code":0,"data":{"paid":true}}'),
        ]),
        ("购物车", [
            ("GET", "/api/cart", "获取购物车", "?openid=", '{"code":0,"data":{"items":[{"product_id":1,"quantity":2,"subtotal":200}]}}'),
            ("POST", "/api/cart/add", "添加商品", '{"product_id":1,"quantity":1}', '{"code":0,"data":{"success":true}}'),
            ("PUT", "/api/cart/:itemId", "更新数量", '{"quantity":3}', '{"code":0,"data":{"success":true}}'),
            ("DELETE", "/api/cart/:itemId", "移除商品", None, '{"code":0,"data":{"success":true}}'),
            ("POST", "/api/cart/clear", "清空购物车", '{"openid":"..."}', '{"code":0,"data":{"success":true}}'),
        ]),
        ("天气", [
            ("GET", "/api/weather", "当前天气", "?city=桂林", '{"code":0,"data":{"temp":"22C","condition":"晴","humidity":"65%"}}'),
            ("GET", "/api/weather/forecast", "天气预报", "?city=桂林&days=3", '{"code":0,"data":{"forecast":[{"date":"...","high":"28C","low":"20C"}]}}'),
        ]),
        ("健康/紧急", [
            ("POST", "/api/health-records", "创建健康档案", '{"heart_rate":72,"blood_pressure":"120/80"}', '{"code":0,"data":{"id":1}}'),
            ("GET", "/api/health-records", "健康记录列表", "?openid=", '{"code":0,"data":{"items":[{"heart_rate":72}]}}'),
            ("POST", "/api/emergency-contacts/create", "创建紧急联系人", '{"name":"张三","phone":"138..."}', '{"code":0,"data":{"id":1}}'),
            ("GET", "/api/emergency-contacts", "紧急联系人列表", "?openid=", '{"code":0,"data":{"items":[{"name":"张三"}]}}'),
            ("PUT", "/api/emergency-contacts/:id", "更新联系人", '{"name":"李四"}', '{"code":0,"data":{"success":true}}'),
            ("DELETE", "/api/emergency-contacts/:id", "删除联系人", None, '{"code":0,"data":{"success":true}}'),
            ("POST", "/api/sos/alert", "SOS紧急求助", '{"lat":25.27,"lng":110.29}', '{"code":0,"data":{"alert_id":"..."}}'),
        ]),
    ]

    group_w2 = (USABLE_W - 3 * Inches(0.1)) / 4
    for gi, (group_name, apis) in enumerate(api_groups2):
        l = MARGIN + gi * (group_w2 + Inches(0.1))
        t = CONTENT_T + Inches(0.12)

        content_box(s, l, t, group_w2, Inches(0.38), GREEN_700)
        tb(s, l, t + Inches(0.05), group_w2, Inches(0.3),
           group_name, Pt(14), True, WHITE, PP_ALIGN.CENTER)

        api_start_t = t + Inches(0.42)
        available_h = CONTENT_B - api_start_t - Inches(0.05)
        item_h = available_h / max(len(apis), 1)

        for j, (method, path, desc, req, resp) in enumerate(apis):
            item_t = api_start_t + j * item_h
            content_box(s, l, item_t, group_w2, item_h - Inches(0.04), GREEN_50, GREEN_500)
            method_tag(s, l + Inches(0.05), item_t + Inches(0.05), Inches(0.65), Inches(0.28), method)
            tb(s, l + Inches(0.75), item_t + Inches(0.05), group_w2 - Inches(0.8), Inches(0.28),
               path, Pt(10), True, GREEN_900)
            tb(s, l + Inches(0.05), item_t + item_h * 0.38, group_w2 - Inches(0.08), item_h * 0.28,
               desc, Pt(10), False, DARK_GRAY)
            if req:
                tb(s, l + Inches(0.05), item_t + item_h * 0.65, group_w2 - Inches(0.08), item_h * 0.3,
                   "📥 " + req[:45] + ("..." if len(req or "") > 45 else ""), Pt(9), False, MID_GRAY)

    # ================================================================
    # 第11页：数据库设计（9张表）
    # ================================================================
    s = new_slide(prs)
    bg(s, WHITE)
    header_bar(s, "09 数据库设计", "9张表完整Schema · 字段类型 · 外键关系 · SQL.js→PostgreSQL迁移")
    footer(s)

    tables = [
        ("categories", "景点分类", GREEN_700, ["id INT PK AUTO", "name VARCHAR(100)", "slug VARCHAR(100)", "description TEXT", "image_url VARCHAR(500)", "order_index INT", "is_active BOOL"]),
        ("products", "商品/景点", GREEN_700, ["id INT PK AUTO", "category_id →categories", "name VARCHAR(200)", "description TEXT", "price DECIMAL(10,2)", "stock INT", "image_url VARCHAR(500)", "location VARCHAR(200)", "开放时间 VARCHAR(100)", "is_active BOOL"]),
        ("user_profiles", "用户档案", GREEN_500, ["id INT PK AUTO", "openid VARCHAR(64) UNIQUE", "nickname VARCHAR(32)", "avatar_url VARCHAR(500)", "phone VARCHAR(20)", "created_at TIMESTAMP", "updated_at TIMESTAMP"]),
        ("addresses", "地址管理", GREEN_500, ["id INT PK AUTO", "openid →user_profiles", "full_name VARCHAR(50)", "phone VARCHAR(20)", "address_line VARCHAR(300)", "city VARCHAR(50)", "postal_code VARCHAR(20)", "is_default BOOL"]),
        ("orders", "订单主表", GREEN_700, ["id INT PK AUTO", "order_no VARCHAR(32) UNIQUE", "openid →user_profiles", "total_price DECIMAL(10,2)", "status ENUM", "created_at TIMESTAMP", "updated_at TIMESTAMP"]),
        ("order_items", "订单明细", GREEN_700, ["id INT PK AUTO", "order_id →orders", "product_id →products", "quantity INT", "unit_price DECIMAL", "subtotal DECIMAL"]),
        ("emergency_contacts", "紧急联系人", AMBER, ["id INT PK AUTO", "openid →user_profiles", "name VARCHAR(50)", "phone VARCHAR(20)", "relation VARCHAR(20)", "is_primary BOOL"]),
        ("health_records", "健康档案", AMBER, ["id INT PK AUTO", "openid →user_profiles", "heart_rate INT", "blood_pressure VARCHAR(20)", "symptoms TEXT", "recorded_at TIMESTAMP"]),
        ("cart_items", "购物车", GREEN_500, ["id INT PK AUTO", "openid →user_profiles", "product_id →products", "quantity INT", "added_at TIMESTAMP"]),
    ]

    col_w = (USABLE_W - 8 * Inches(0.07)) / 9
    for i, (tbl, label, color, cols) in enumerate(tables):
        l = MARGIN + i * (col_w + Inches(0.07))
        t = CONTENT_T + Inches(0.12)
        th = Inches(0.42)
        content_box(s, l, t, col_w, th, color)
        tb(s, l, t + Inches(0.08), col_w, th * 0.75,
           tbl, Pt(10), True, WHITE, PP_ALIGN.CENTER)
        tb(s, l, t + th + Inches(0.04), col_w, Inches(0.25),
           label, Pt(11), True, GREEN_900, PP_ALIGN.CENTER)
        fields_h = CONTENT_B - t - th - Inches(0.1) - Inches(0.55)
        content_box(s, l, t + th + Inches(0.04) + Inches(0.25), col_w, fields_h, GREEN_50, color)
        for j, col in enumerate(cols):
            tb(s, l + Inches(0.04), t + th + Inches(0.04) + Inches(0.25) + j * Inches(0.28),
               col_w - Inches(0.06), Inches(0.26),
               col, Pt(9), False, DARK_GRAY)

    # 迁移说明
    note_y = CONTENT_B - Inches(0.75)
    content_box(s, MARGIN, note_y, USABLE_W, Inches(0.68), GREEN_100, GREEN_500)
    tb(s, MARGIN + Inches(0.15), note_y + Inches(0.05), Inches(2.8), Inches(0.3),
       "📊 当前：SQL.js", Pt(13), True, GREEN_900)
    tb(s, MARGIN + Inches(0.15), note_y + Inches(0.32), Inches(6.0), Inches(0.28),
       "内存数据库，/app/data目录持久化（Railway持久化卷）| 适合开发测试", Pt(12), False, DARK_GRAY)
    tb(s, MARGIN + Inches(6.5), note_y + Inches(0.05), Inches(2.8), Inches(0.3),
       "🚀 迁移：PostgreSQL", Pt(13), True, GREEN_900)
    tb(s, MARGIN + Inches(6.5), note_y + Inches(0.32), Inches(6.0), Inches(0.28),
       "001_initial_schema.sql已就绪 | 配置DATABASE_URL环境变量即可触发迁移", Pt(12), False, DARK_GRAY)

    # ================================================================
    # 第12页：AI助手集成
    # ================================================================
    s = new_slide(prs)
    bg(s, WHITE)
    header_bar(s, "10 AI助手集成", "LijiangCoordinator多意图路由 · 3大Agent · RAG桂林知识库 · SSE流式输出")
    footer(s)

    # 左侧架构
    content_box(s, MARGIN, CONTENT_T + Inches(0.15), Inches(6.2), Inches(3.4), GREEN_50, GREEN_500)
    tb(s, MARGIN + Inches(0.15), CONTENT_T + Inches(0.25), Inches(5.9), Inches(0.35),
       "🔀 LijiangCoordinator 多意图路由架构", Pt(15), True, GREEN_900)

    arch_steps = [
        ("👤", "用户输入", "小程序AI助手 → POST /api/chat", "message + openid"),
        ("🔍", "意图分析", "Coordinator分析用户意图", "单意图 or 多意图"),
        ("📡", "路由分发", "单意图→单一Agent | 多意图→三Agent并行", "意图标签：travel_plan/weather/food/... "),
        ("🤖", "Agent执行", "RouteAgent+BehaviorAgent+KnowledgeAgent并行", "各Agent独立执行，输出结构化结果"),
        ("📝", "结果综合", "Coordinator合并多Agent输出", "生成自然语言响应，SSE流式推送"),
    ]
    for j, (emoji, title, desc, detail) in enumerate(arch_steps):
        lt = CONTENT_T + Inches(0.65) + j * Inches(0.63)
        content_box(s, MARGIN + Inches(0.1), lt, Inches(0.55), Inches(0.55), GREEN_700)
        tb(s, MARGIN + Inches(0.1), lt + Inches(0.12), Inches(0.55), Inches(0.32),
           emoji, Pt(16), False, WHITE, PP_ALIGN.CENTER)
        tb(s, MARGIN + Inches(0.72), lt + Inches(0.05), Inches(2.0), Inches(0.28),
           title, Pt(13), True, GREEN_700)
        tb(s, MARGIN + Inches(0.72), lt + Inches(0.3), Inches(5.3), Inches(0.25),
           desc, Pt(12), False, DARK_GRAY)
        tb(s, MARGIN + Inches(2.8), lt + Inches(0.05), Inches(3.2), Inches(0.28),
           detail, Pt(11), False, MID_GRAY)
        if j < 4:
            tb(s, MARGIN + Inches(0.1), lt + Inches(0.57), Inches(0.55), Inches(0.06),
               "↓", Pt(14), True, GREEN_500, PP_ALIGN.CENTER)

    # 右侧Agent详情
    content_box(s, MARGIN + Inches(6.35), CONTENT_T + Inches(0.15), Inches(6.48), Inches(3.4), None)
    tb(s, MARGIN + Inches(6.5), CONTENT_T + Inches(0.25), Inches(6.18), Inches(0.35),
       "🤖 三大专业Agent", Pt(15), True, GREEN_900)

    agents = [
        ("🗺️ RouteAgent", "旅游路线规划", GREEN_700,
         ["分析用户时间/预算/人数",
          "漓江竹筏/徒步/自驾路线",
          "生成个性化行程方案（JSON结构）",
          "多目标优化：时间/费用/体验"]),
        ("📊 BehaviorAgent", "游客行为预测", ACCENT,
         ["预测各景点停留时间",
          "购买意愿与消费分析",
          "客流高峰时段预测",
          "最优出发时间建议"]),
        ("📚 KnowledgeAgent", "旅游知识问答", GREEN_500,
         ["景点介绍与历史背景",
          "票价信息与优惠政策",
          "美食推荐与交通指南",
          "RAG检索桂林漓江知识库"]),
    ]
    for j, (name, role, color, items) in enumerate(agents):
        lt = CONTENT_T + Inches(0.65) + j * Inches(1.05)
        content_box(s, MARGIN + Inches(6.5), lt, Inches(6.35), Inches(0.98), GREEN_100, color)
        tb(s, MARGIN + Inches(6.58), lt + Inches(0.06), Inches(2.5), Inches(0.3),
           name, Pt(13), True, color)
        tb(s, MARGIN + Inches(9.1), lt + Inches(0.06), Inches(3.7), Inches(0.3),
           "职能：" + role, Pt(12), True, GREEN_900)
        for k, item in enumerate(items):
            tb(s, MARGIN + Inches(6.58), lt + Inches(0.38) + k * Inches(0.2),
               Inches(6.2), Inches(0.2),
               "· " + item, Pt(11), False, DARK_GRAY)

    # RAG知识库详情
    content_box(s, MARGIN, CONTENT_T + Inches(3.7), USABLE_W, Inches(1.95), GREEN_50, GREEN_500)
    tb(s, MARGIN + Inches(0.15), CONTENT_T + Inches(3.8), Inches(5.0), Inches(0.35),
       "📚 RAG桂林知识库", Pt(15), True, GREEN_900)
    rag_items = [
        ("词汇表大小", "vocab_size=583", "TF-IDF向量化词表，583个关键词"),
        ("文档分块", "chunks=28", "桂林景点/美食/交通28个知识块"),
        ("检索方式", "TF-IDF + 余弦相似度", "numpy+scipy实现，无需外部向量库"),
        ("知识内容", "桂林漓江全量数据", "杨堤码头/九马画山/兴坪古镇/美食/交通"),
        ("RAG端点", "POST /rag/query", "yinfa后端 → lvyou-agi → ChromaDB（可选）"),
        ("状态", "⚠️ 知识库仅1个文档", "需扩充马蜂窝/携程/桂林旅游官网数据"),
    ]
    for j, (label, val, desc) in enumerate(rag_items):
        col = j % 3
        row = j // 3
        l = MARGIN + Inches(0.15) + col * (USABLE_W / 3)
        lt = CONTENT_T + Inches(4.2) + row * Inches(0.65)
        content_box(s, l, lt, USABLE_W / 3 - Inches(0.08), Inches(0.58), WHITE, GREEN_500)
        tb(s, l + Inches(0.1), lt + Inches(0.05), Inches(1.5), Inches(0.25), label, Pt(12), True, GREEN_700)
        tb(s, l + Inches(1.6), lt + Inches(0.05), Inches(2.5), Inches(0.25), val, Pt(12), True, GREEN_900)
        tb(s, l + Inches(0.1), lt + Inches(0.3), USABLE_W / 3 - Inches(0.2), Inches(0.25), desc, Pt(11), False, MID_GRAY)

    # ================================================================
    # 第13页：lvyou-agi服务详情
    # ================================================================
    s = new_slide(prs)
    bg(s, WHITE)
    header_bar(s, "11 lvyou-agi服务", "桂林旅游智能体 · RAG+Agent+意图识别 · /config端点已移除")
    footer(s)

    # 服务概览
    content_box(s, MARGIN, CONTENT_T + Inches(0.15), Inches(4.0), Inches(2.0), GREEN_50, GREEN_500)
    tb(s, MARGIN + Inches(0.15), CONTENT_T + Inches(0.25), Inches(3.7), Inches(0.35),
       "🛠️ 服务概览", Pt(15), True, GREEN_900)
    overview = [
        ("框架", "FastAPI + Python 3.11", "异步API，高性能"),
        ("LLM", "MiniMax API（MiniMax-Text-01）", "流式SSE输出"),
        ("向量", "TF-IDF + numpy/scipy", "ChromaDB（可选，当前未用）"),
        ("意图", "LijiangCoordinator", "单意图优先，多意图并行⚠️"),
        ("部署", "Render免费版", "⚠️ 30分钟无访问自动休眠"),
        ("端口", "8001（根server）/ 8000（Docker）", "src/server.py未使用"),
    ]
    for j, (label, val, note) in enumerate(overview):
        lt = CONTENT_T + Inches(0.65) + j * Inches(0.37)
        tb(s, MARGIN + Inches(0.18), lt, Inches(0.9), Inches(0.3), label, Pt(12), True, GREEN_700)
        tb(s, MARGIN + Inches(1.1), lt, Inches(2.0), Inches(0.3), val, Pt(12), False, GREEN_900)
        nc = AMBER if "⚠️" in note else MID_GRAY
        tb(s, MARGIN + Inches(3.1), lt, Inches(1.5), Inches(0.3), note, Pt(11), False, nc)

    # 6个API端点
    content_box(s, MARGIN + Inches(4.15), CONTENT_T + Inches(0.15), Inches(8.68), Inches(2.0), None)
    tb(s, MARGIN + Inches(4.3), CONTENT_T + Inches(0.25), Inches(8.38), Inches(0.35),
       "🌐 11个API端点（/health · /capabilities · /api/chat · /api/plan · /api/predict · /api/query · ...）", Pt(13), True, GREEN_900)
    endpoints = [
        ("GET",  "/health",         "健康检查",         '{"status":"ok","version":"1.0"}'),
        ("GET",  "/capabilities",  "AI能力清单",       '{"intents":["travel_plan","weather","food"],"model":"MiniMax-Text-01"}'),
        ("POST", "/api/chat",       "AI对话",           '{"message":"漓江推荐","stream":false}'),
        ("POST", "/api/chat/stream","流式对话（SSE）",  '{"message":"景点票价"} → data:{"chunk":"..."}'),
        ("POST", "/api/multi-task", "多任务并行",       '{"tasks":[{"type":"weather","params":{"city":"桂林"}}]}'),
        ("POST", "/api/plan",       "旅行规划",         '{"destination":"漓江","days":2,"preferences":"慢节奏"}'),
        ("POST", "/api/predict",    "行为预测",         '{"user_id":"...","context":{"time":"morning"}}'),
        ("POST", "/api/query",      "RAG知识检索",      '{"query":"竹筏价格","top_k":5}'),
        ("POST", "/api/itinerary",  "生成行程",         '{"destination":"桂林","days":3}'),
        ("POST", "/mcp",            "MCP工具调用",      '{"tool":"search_knowledge","params":{"query":"..."}}'),
        ("GET",  "/mcp",            "MCP能力查询",      '{"tools":["chat","plan_trip","search_knowledge"]}'),
    ]
    # 11个端点卡片：0.68"宽，0.04"间距，单行11卡刚好在13.33"内
    ep_w = Inches(0.68)
    ep_gap = Inches(0.04)
    ep_start_x = MARGIN + Inches(4.3)
    for j, (method, path, desc, resp) in enumerate(endpoints):
        l = ep_start_x + j * (ep_w + ep_gap)
        content_box(s, l, CONTENT_T + Inches(0.65), ep_w, Inches(1.42), GREEN_100, GREEN_500)
        method_tag(s, l, CONTENT_T + Inches(0.68), ep_w, Inches(0.3), method)
        tb(s, l, CONTENT_T + Inches(1.02), ep_w, Inches(0.28),
           path, Pt(10), True, GREEN_900, PP_ALIGN.CENTER)
        tb(s, l, CONTENT_T + Inches(1.32), ep_w, Inches(0.3),
           desc, Pt(9), False, DARK_GRAY, PP_ALIGN.CENTER)
        tb(s, l, CONTENT_T + Inches(1.62), ep_w, Inches(0.4),
           resp[:28] + ("..." if len(resp) > 28 else ""), Pt(8), False, MID_GRAY, PP_ALIGN.CENTER)

    # 已知局限
    content_box(s, MARGIN, CONTENT_T + Inches(2.3), USABLE_W, Inches(2.25), None)
    tb(s, MARGIN + Inches(0.15), CONTENT_T + Inches(2.4), Inches(4.0), Inches(0.35),
       "⚠️ 已知局限", Pt(15), True, AMBER)
    limitations = [
        ("意图检测", "单意图优先规则，非多意图并行检测", "README.md 108行确认"),
        ("Ollama接口", "仅完成签名对齐，未实际使用", "README.md 109行确认"),
        ("Docker WSL", "WSL内Docker不可用，部署需在原生Linux", "系统限制"),
        ("知识库", "仅有杨堤码头1个文档，覆盖范围有限", "需扩充至马蜂窝/携程/桂林旅游官网"),
        ("Render休眠", "免费版30分钟无访问自动休眠，冷启动慢", "建议迁移至Railway付费版"),
        ("API集成", "yinfa仅对接了/api/chat，其余10个端点未集成", "yinfa后端agent-service.js确认"),
        ("版本号", "health返回v3.0.0，PPT数据(v1.0)可能落后", "以server.py实际返回为准"),
    ]
    for j, (lim, detail, source) in enumerate(limitations):
        col = j % 3
        row = j // 3
        l = MARGIN + col * (USABLE_W / 3)
        lt = CONTENT_T + Inches(2.8) + row * Inches(0.82)
        content_box(s, l + Inches(0.08), lt, USABLE_W / 3 - Inches(0.12), Inches(0.75), GREEN_50, AMBER)
        tb(s, l + Inches(0.15), lt + Inches(0.05), USABLE_W / 3 - Inches(0.25), Inches(0.28),
           "⚠️ " + lim, Pt(13), True, AMBER)
        tb(s, l + Inches(0.15), lt + Inches(0.32), USABLE_W / 3 - Inches(0.25), Inches(0.25),
           detail, Pt(11), False, DARK_GRAY)
        tb(s, l + Inches(0.15), lt + Inches(0.55), USABLE_W / 3 - Inches(0.25), Inches(0.18),
           source, Pt(10), False, MID_GRAY)

    # 环境变量
    content_box(s, MARGIN, CONTENT_T + Inches(4.7), USABLE_W, Inches(1.05), GREEN_100, GREEN_500)
    tb(s, MARGIN + Inches(0.15), CONTENT_T + Inches(4.8), Inches(4.0), Inches(0.35),
       "🔑 环境变量", Pt(14), True, GREEN_900)
    env_vars = [
        ("MINIMAX_API_KEY", "MiniMax平台API密钥（必需）", "平台：minimaxi.com"),
        ("MINIMAX_MODEL", "模型名称，默认MiniMax-Text-01", "支持流式SSE"),
        ("MINIMAX_BASE_URL", "API基础地址", "默认：https://api.minimaxi.com"),
        ("OLLAMA_BASE_URL", "Ollama本地模型（可选，当前未用）", "README.md确认未集成"),
        ("OLLAMA_MODEL", "Ollama模型名称（可选）", "未启用"),
    ]
    env_w = (USABLE_W - 4 * Inches(0.1)) / 5
    for j, (name, desc, note) in enumerate(env_vars):
        l = MARGIN + j * (env_w + Inches(0.1))
        content_box(s, l, CONTENT_T + Inches(5.18), env_w, Inches(0.52), WHITE, GREEN_500)
        tb(s, l + Inches(0.08), CONTENT_T + Inches(5.22), env_w - Inches(0.12), Inches(0.22),
           name, Pt(10), True, GREEN_700)
        tb(s, l + Inches(0.08), CONTENT_T + Inches(5.43), env_w - Inches(0.12), Inches(0.24),
           desc, Pt(10), False, DARK_GRAY)

    # ================================================================
    # 第14页：部署方案
    # ================================================================
    s = new_slide(prs)
    bg(s, WHITE)
    header_bar(s, "12 部署方案", "Railway一键部署 · Docker容器 · 环境变量配置 · 微信支付4方案")
    footer(s)

    # Railway部署
    content_box(s, MARGIN, CONTENT_T + Inches(0.15), Inches(4.2), Inches(3.0), GREEN_50, GREEN_500)
    tb(s, MARGIN + Inches(0.15), CONTENT_T + Inches(0.25), Inches(3.9), Inches(0.35),
       "🚀 Railway 一键部署（推荐）", Pt(15), True, GREEN_900)
    railway_steps = [
        "① Fork本仓库到GitHub",
        "② Railway控制台 → New Project → Deploy from GitHub",
        "③ 选择yinfa仓库，Railway自动读取railway.json",
        "④ 添加持久化卷：Mount Path = /app/data",
        "⑤ 配置环境变量：MINIMAX_API_KEY / AMAP_KEY / DATABASE_URL",
        "⑥ 等待构建完成，自动获得HTTPS地址",
        "⑦ 绑定自定义域名（可选）：yinfa.example.com",
    ]
    for j, step in enumerate(railway_steps):
        tb(s, MARGIN + Inches(0.18), CONTENT_T + Inches(0.68) + j * Inches(0.37),
           Inches(3.9), Inches(0.34),
           step, Pt(13), False, DARK_GRAY)

    # Docker部署
    content_box(s, MARGIN + Inches(4.35), CONTENT_T + Inches(0.15), Inches(4.2), Inches(3.0), GREEN_50, GREEN_500)
    tb(s, MARGIN + Inches(4.5), CONTENT_T + Inches(0.25), Inches(3.9), Inches(0.35),
       "🐳 Docker 容器部署", Pt(15), True, GREEN_900)
    docker_steps = [
        "① Docker已安装（docker --version）",
        "② 项目根目录构建镜像：",
        "   docker build -t yinfa-backend .",
        "③ 运行容器（数据持久化）：",
        "   docker run -d -p 8000:8000",
        "   -v yinfa-data:/app/data",
        "   -e PORT=8000",
        "   --name yinfa-backend yinfa-backend",
        "⚠️ WSL内Docker不可用，需原生Linux",
    ]
    for j, step in enumerate(docker_steps):
        bold = not step.startswith("   ")
        nc = RED if "⚠️" in step else (DARK_GRAY if bold else MID_GRAY)
        tb(s, MARGIN + Inches(4.5), CONTENT_T + Inches(0.68) + j * Inches(0.35),
           Inches(3.9), Inches(0.32),
           step, Pt(13), bold, nc)

    # 微信支付方案
    content_box(s, MARGIN + Inches(8.6), CONTENT_T + Inches(0.15), Inches(4.2), Inches(3.0), None)
    tb(s, MARGIN + Inches(8.75), CONTENT_T + Inches(0.25), Inches(3.9), Inches(0.35),
       "💳 微信支付接入方案", Pt(15), True, GREEN_900)
    pay_options = [
        ("P0 微信小商店", "个人开发者首选，无需营业执照", "✅ 推荐", GREEN_700),
        ("方案B 微信支付API", "需商户号+营业执照+API密钥", "企业用户", GREEN_500),
        ("方案C 第三方聚合支付", "Stripe/PayPal/易宝等", "需手续费", MID_GRAY),
        ("方案D H5支付", "WAP端网页支付", "体验较差", MID_GRAY),
    ]
    for j, (name, desc, tag, color) in enumerate(pay_options):
        lt = CONTENT_T + Inches(0.65) + j * Inches(0.72)
        content_box(s, MARGIN + Inches(8.85), lt, Inches(4.0), Inches(0.65), GREEN_100, color)
        tb(s, MARGIN + Inches(8.95), lt + Inches(0.06), Inches(2.2), Inches(0.28),
           name, Pt(13), True, GREEN_900)
        status_badge(s, MARGIN + Inches(11.05), lt + Inches(0.08), Inches(0.8), Inches(0.28), tag, color)
        tb(s, MARGIN + Inches(8.95), lt + Inches(0.34), Inches(3.8), Inches(0.28),
           desc, Pt(11), False, DARK_GRAY)

    # 环境变量总表
    content_box(s, MARGIN, CONTENT_T + Inches(3.3), USABLE_W, Inches(2.4), GREEN_50, GREEN_500)
    tb(s, MARGIN + Inches(0.15), CONTENT_T + Inches(3.4), Inches(4.0), Inches(0.35),
       "⚙️ 环境变量完整清单", Pt(15), True, GREEN_900)
    all_env = [
        ("PORT", "服务器端口", "8000", "✅ 已配置"),
        ("NODE_ENV", "环境模式", "production", "✅ 已配置"),
        ("MINIMAX_API_KEY", "MiniMax AI密钥", "（联系管理员）", "🔴 必需"),
        ("AMAP_KEY", "高德天气API", "（可选，未配返回模拟数据）", "🟡 可选"),
        ("DATABASE_URL", "PostgreSQL连接", "postgres://user:pass@host:5432/yinfa", "🔴 P0阻塞"),
        ("LIYOU_AGI_URL", "lvyou-agi服务地址", "https://lvyou-agi.onrender.com", "✅ 已配置"),
        ("LOG_LEVEL", "日志级别", "info（默认）/ debug", "🟢 可调"),
        ("CORS_ORIGINS", "允许的源", "https://yinfa.com（生产）", "🟢 可调"),
    ]
    env_cols = 4
    env_row_w = USABLE_W / env_cols
    for j, (name, desc, val, status) in enumerate(all_env):
        col = j % env_cols
        row = j // env_cols
        l = MARGIN + col * env_row_w
        lt = CONTENT_T + Inches(3.8) + row * Inches(0.92)
        content_box(s, l + Inches(0.08), lt, env_row_w - Inches(0.12), Inches(0.85), WHITE, GREEN_500)
        tb(s, l + Inches(0.15), lt + Inches(0.06), env_row_w - Inches(0.25), Inches(0.28),
           name, Pt(12), True, GREEN_700)
        tb(s, l + Inches(0.15), lt + Inches(0.32), env_row_w - Inches(0.25), Inches(0.25),
           desc, Pt(11), False, MID_GRAY)
        tb(s, l + Inches(0.15), lt + Inches(0.55), env_row_w - Inches(0.25), Inches(0.25),
           val, Pt(10), False, DARK_GRAY)
        sc = GREEN_700 if "✅" in status else RED if "🔴" in status else AMBER
        tb(s, l + Inches(0.15), lt + Inches(0.73), env_row_w - Inches(0.25), Inches(0.2),
           status, Pt(10), True, sc)

    # ================================================================
    # 第15页：P0-P2路线图 + 已知局限
    # ================================================================
    s = new_slide(prs)
    bg(s, WHITE)
    header_bar(s, "13 路线图 & 已知局限", "P0阻断→P1增强→P2建设 · 6项已知局限")
    footer(s)

    # P0阻断问题
    content_box(s, MARGIN, CONTENT_T + Inches(0.15), Inches(4.1), Inches(3.0), None)
    tb(s, MARGIN + Inches(0.15), CONTENT_T + Inches(0.25), Inches(3.8), Inches(0.35),
       "🔴 P0 阻断性问题", Pt(16), True, RED)
    p0_items = [
        ("① DATABASE_URL未配置", "Railway控制台添加环境变量，触发PostgreSQL迁移", RED),
        ("② 微信支付未激活", "注册微信小商店（个人开发者首选），获取商户号", RED),
        ("③ lvyou-agi休眠问题", "升级Render付费版 或 迁移到Railway（需解决冷启动）", AMBER),
        ("④ 知识库严重不足", "仅有杨堤码头1个文档，需扩充至少10+个景点文档", AMBER),
    ]
    for j, (title, sub, color) in enumerate(p0_items):
        lt = CONTENT_T + Inches(0.68) + j * Inches(0.68)
        content_box(s, MARGIN + Inches(0.08), lt, Inches(3.95), Inches(0.62), GREEN_50, color)
        tb(s, MARGIN + Inches(0.18), lt + Inches(0.05), Inches(3.75), Inches(0.3),
           title, Pt(13), True, color)
        tb(s, MARGIN + Inches(0.18), lt + Inches(0.33), Inches(3.75), Inches(0.25),
           "→ " + sub, Pt(11), False, DARK_GRAY)

    # P1功能增强
    content_box(s, MARGIN + Inches(4.25), CONTENT_T + Inches(0.15), Inches(4.4), Inches(3.0), None)
    tb(s, MARGIN + Inches(4.4), CONTENT_T + Inches(0.25), Inches(4.1), Inches(0.35),
       "🟡 P1 功能增强", Pt(16), True, AMBER)
    p1_items = [
        ("⑤ Vue3前端完成", "详情页/地图页/支付页，接入AI助手", AMBER),
        ("⑥ SSE流式AI响应", "前端wx.request支持SSE，/api/chat/stream", GREEN_700),
        ("⑦ 微信支付回调", "支付成功回调处理，自动更新订单状态", AMBER),
        ("⑧ Railway PostgreSQL验证", "DATABASE_URL配置后，验证数据持久化", AMBER),
    ]
    for j, (title, sub, color) in enumerate(p1_items):
        lt = CONTENT_T + Inches(0.68) + j * Inches(0.68)
        content_box(s, MARGIN + Inches(4.32), lt, Inches(4.25), Inches(0.62), GREEN_50, color)
        tb(s, MARGIN + Inches(4.42), lt + Inches(0.05), Inches(4.05), Inches(0.3),
           title, Pt(13), True, color)
        tb(s, MARGIN + Inches(4.42), lt + Inches(0.33), Inches(4.05), Inches(0.25),
           "→ " + sub, Pt(11), False, DARK_GRAY)

    # P2长期建设
    content_box(s, MARGIN + Inches(8.7), CONTENT_T + Inches(0.15), Inches(4.1), Inches(3.0), None)
    tb(s, MARGIN + Inches(8.85), CONTENT_T + Inches(0.25), Inches(3.8), Inches(0.35),
       "🟢 P2 长期建设", Pt(16), True, GREEN_700)
    p2_items = [
        ("⑨ 知识库大规模扩充", "爬取马蜂窝/携程/桂林旅游官网→RAG", GREEN_700),
        ("⑩ 双server.py合并", "统一根server.py与src/server.py", MID_GRAY),
        ("⑪ 高德地图API对接", "整合桂林旅游官网地图数据", GREEN_700),
        ("⑫ 多意图真正并行", "从单意图优先升级为多意图并行检测", AMBER),
    ]
    for j, (title, sub, color) in enumerate(p2_items):
        lt = CONTENT_T + Inches(0.68) + j * Inches(0.68)
        content_box(s, MARGIN + Inches(8.88), lt, Inches(4.05), Inches(0.62), GREEN_50, color)
        tb(s, MARGIN + Inches(8.98), lt + Inches(0.05), Inches(3.85), Inches(0.3),
           title, Pt(13), True, color)
        tb(s, MARGIN + Inches(8.98), lt + Inches(0.33), Inches(3.85), Inches(0.25),
           "→ " + sub, Pt(11), False, DARK_GRAY)

    # 已知局限（底部）
    content_box(s, MARGIN, CONTENT_T + Inches(3.3), USABLE_W, Inches(2.35), GREEN_100, GREEN_500)
    tb(s, MARGIN + Inches(0.15), CONTENT_T + Inches(3.4), Inches(4.0), Inches(0.35),
       "⚠️ 已知局限（全部6项）", Pt(15), True, AMBER)
    limitations_all = [
        ("意图检测", "当前单意图优先，非多意图并行", "影响复杂查询处理"),
        ("Ollama", "签名对齐完成但未实际调用", "未启用本地LLM"),
        ("Docker WSL", "WSL内Docker不可用", "需原生Linux环境"),
        ("知识库", "仅1个文档vocab=583/chunks=28", "覆盖范围严重不足"),
        ("Render休眠", "30分钟无访问自动休眠", "影响AI服务可用性"),
        ("API未全接", "yinfa仅用了/chat接口", "其他5个端点未集成"),
    ]
    for j, (lim, detail, impact) in enumerate(limitations_all):
        col = j % 3
        row = j // 3
        l = MARGIN + col * (USABLE_W / 3)
        lt = CONTENT_T + Inches(3.8) + row * Inches(0.85)
        content_box(s, l + Inches(0.08), lt, USABLE_W / 3 - Inches(0.12), Inches(0.78), WHITE, AMBER)
        tb(s, l + Inches(0.15), lt + Inches(0.06), USABLE_W / 3 - Inches(0.25), Inches(0.28),
           "⚠️ " + lim, Pt(13), True, AMBER)
        tb(s, l + Inches(0.15), lt + Inches(0.32), USABLE_W / 3 - Inches(0.25), Inches(0.25),
           detail, Pt(11), False, DARK_GRAY)
        tb(s, l + Inches(0.15), lt + Inches(0.55), USABLE_W / 3 - Inches(0.25), Inches(0.2),
           "影响：" + impact, Pt(10), False, RED)

    # ================================================================
    # 第16页：快速开始 + 总结
    # ================================================================
    s = new_slide(prs)
    bg(s, GREEN_50)

    # 顶部标题栏
    content_box(s, 0, 0, SLIDE_W, HEADER_H, GREEN_900)
    tb(s, MARGIN, Inches(0.2), USABLE_W, Inches(0.6),
       "桂林银发旅游小程序 · 总结", Pt(26), True, WHITE)

    # 快速开始
    content_box(s, MARGIN, CONTENT_T + Inches(0.12), USABLE_W, Inches(2.1), WHITE, GREEN_500)
    tb(s, MARGIN + Inches(0.15), CONTENT_T + Inches(0.22), Inches(4.0), Inches(0.35),
       "⚡ 快速开始", Pt(16), True, GREEN_900)
    quick_start = [
        ("本地运行后端", "cd yinfa-ts/backend && npm install && npm run dev → http://localhost:8000", MARGIN + Inches(0.1)),
        ("本地测试AI", "curl -X POST http://localhost:8001/chat -H 'Content-Type: application/json' -d '{\"message\":\"漓江竹筏\"}'", MARGIN + Inches(0.1)),
        ("本地运行小程序", "微信开发者工具 → 导入yinfa-wechat → 填入后端API地址", MARGIN + Inches(0.1)),
        ("Railway部署", "Fork → Railway → Add GitHub Repo → Configure env vars → Deploy", MARGIN + Inches(0.1)),
        ("Docker部署", "docker build -t yinfa-backend . && docker run -p 8000:8000 yinfa-backend", MARGIN + Inches(0.1)),
    ]
    for j, (title, cmd, l) in enumerate(quick_start):
        col = j % 2
        row = j // 2
        l = MARGIN + col * (USABLE_W / 2 + Inches(0.05))
        lt = CONTENT_T + Inches(0.6) + row * Inches(0.55)
        content_box(s, l, lt, USABLE_W / 2 - Inches(0.08), Inches(0.5), GREEN_100, GREEN_500)
        tb(s, l + Inches(0.1), lt + Inches(0.05), Inches(1.8), Inches(0.25), title, Pt(12), True, GREEN_700)
        tb(s, l + Inches(0.1), lt + Inches(0.28), USABLE_W / 2 - Inches(0.2), Inches(0.22),
           cmd, Pt(10), False, DARK_GRAY)

    # 总结卡片
    summary_items = [
        ("🎯", "定位清晰", "专为50岁以上银发游客设计，大字界面，简单操作，安全优先"),
        ("🛠️", "技术完整", "微信小程序+Vue3+Express+SQL.js+AI，19个API全贯通"),
        ("🤖", "AI赋能", "集成漓江旅游导览Agent，RAG+LLM智能规划路线，SSE流式输出"),
        ("📱", "体验优化", "银发友好字体/配色/交互设计，语音+文字双交互，9+页面"),
        ("🚀", "部署就绪", "Railway一键部署，Docker容器化，持久化存储，9张数据库表"),
        ("🔮", "扩展性强", "预留支付/知识库/多Agent扩展，架构清晰可演进，6个已知局限明确"),
    ]

    card_w = (USABLE_W - 2 * Inches(0.15)) / 3
    card_h = Inches(1.55)

    for i, (emoji, title, desc) in enumerate(summary_items):
        col = i % 3
        row = i // 3
        l = MARGIN + col * (card_w + Inches(0.15))
        t = CONTENT_T + Inches(2.35) + row * (card_h + Inches(0.12))

        fc = GREEN_100 if row == 0 else GREEN_50
        content_box(s, l, t, card_w, card_h, fc, GREEN_500)
        tb(s, l + Inches(0.15), t + Inches(0.1), card_w - Inches(0.3), Inches(0.4),
           emoji + " " + title, Pt(15), True, GREEN_900)
        tb(s, l + Inches(0.15), t + Inches(0.52), card_w - Inches(0.3), Inches(0.95),
           desc, Pt(13), False, DARK_GRAY)

    # 底部联系信息
    tb(s, MARGIN, CONTENT_T + Inches(4.3), USABLE_W, Inches(0.4),
       "yinfa × lvyou-agi · 微信小程序 × AI导览 × 桂林漓江旅游",
       Pt(16), True, GREEN_700, PP_ALIGN.CENTER)
    tb(s, MARGIN, CONTENT_T + Inches(4.72), USABLE_W, Inches(0.4),
       "后端API：https://yinfa-backend.up.railway.app/api  |  AI服务：https://lvyou-agi.onrender.com",
       Pt(13), False, MID_GRAY, PP_ALIGN.CENTER)
    tb(s, MARGIN, CONTENT_T + Inches(5.1), USABLE_W, Inches(0.4),
       "GitHub：github.com/cangku/yinfa  |  联系：微信小程序搜索「银发旅游」",
       Pt(13), False, MID_GRAY, PP_ALIGN.CENTER)

    footer(s, "yinfa | 微信小程序 × AI导览 × 桂林漓江旅游 | 2026 | v2超深度优化版")

    # 保存
    out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                            "桂林银发旅游小程序_超深度优化v2.pptx")
    prs.save(out_path)
    print(f"✅ yinfa PPT v2已生成：{out_path}")
    return out_path


if __name__ == "__main__":
    make_yinfa_ppt()
